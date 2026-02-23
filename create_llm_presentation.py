#!/usr/bin/env python3
"""
Create a PowerPoint presentation on LLMs as a more efficient way of gathering information.
Content extracted solely from Nicolas Bustamante's article "The LLM Context Tax: Best Tips for Tax Avoidance"
(https://www.nicolasbustamante.com/p/the-llm-context-tax-best-tips-for)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color scheme - modern tech aesthetic
TITLE_COLOR = RGBColor(0x1a, 0x1a, 0x2e)  # Dark navy
ACCENT_COLOR = RGBColor(0x16, 0x6f, 0x97)  # Teal blue
BODY_COLOR = RGBColor(0x3d, 0x3d, 0x5c)    # Slate
LIGHT_BG = RGBColor(0xf8, 0xf9, 0xfa)     # Off-white


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullet_points, notes=""):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(12)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(8)
    
    return slide


def add_quote_slide(prs, quote, attribution=""):
    """Add a slide with a key quote or statistic."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    if attribution:
        attr_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
        tf = attr_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"— {attribution}"
        p.font.size = Pt(14)
        p.font.color.rgb = BODY_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "The LLM Context Tax: Best Tips for Tax Avoidance",
        "Optimizing Context, Cost & Performance\nNicolas Bustamante — nicolasbustamante.com"
    )
    
    # Slide 2: The Context Tax Problem
    add_content_slide(
        prs,
        "What Is the Context Tax?",
        [
            "The cumulative costs of filling an LLM's context window with unnecessary tokens",
            "Every token sent to an LLM incurs three penalties: higher costs, slower responses, degraded performance",
            "Context rot: Agents become confused by accumulated noise",
            "Most models show sharp performance degradation past 32K tokens",
        ]
    )
    
    # Slide 3: The Triple Penalty
    add_content_slide(
        prs,
        "The LLM Context Tax: Triple Penalty",
        [
            "Higher costs — Cached tokens cost 10x less than uncached (Claude Opus 4.6)",
            "Increased latency — More tokens = slower response times",
            "Degraded performance — Models show significant confusion past 32K tokens",
            "The difference between $0.50 and $5.00 per query often comes down to context management",
        ]
    )
    
    # Slide 4: Most Critical Optimizations (from article)
    add_content_slide(
        prs,
        "Most Critical Optimizations",
        [
            "Stable prefixes for KV cache hits — Keep system prompts identical; move dynamic content to the end (up to 10x cost reduction)",
            "Design precise tools — Smart tool design can reduce token consumption by 10x",
            "Parallel tool calls — Fewer round trips = less context accumulation",
            "Store outputs externally — Avoid context bloat by keeping tool outputs in filesystem",
            "Application-level response caching — The cheapest token is one you never send",
        ]
    )
    
    # Slide 5: Key Quote
    add_quote_slide(
        prs,
        "The most important optimization for production agents: maintaining stable prefixes for KV cache hits—offering up to 10x cost reduction",
        "Nicolas Bustamante, The LLM Context Tax"
    )
    
    # Slide 7: Output Token Budgeting
    add_content_slide(
        prs,
        "Output Token Budgeting: The Hidden Cost",
        [
            "Output tokens cost 5x more than uncached input tokens",
            "The most expensive tokens are those you generate",
            "Budget outputs carefully — request concise responses when possible",
            "Delegate token-heavy operations to cheaper subagents",
            "Reusable templates — stop regenerating the same code repeatedly",
        ]
    )
    
    # Slide 8: Additional Tactics from the Article
    add_content_slide(
        prs,
        "Additional Optimization Tactics",
        [
            "Clean data before it enters context",
            "Use strategic information placement to avoid the lost-in-the-middle problem",
            "Avoid the 200K token pricing cliff",
            "Append-only context — mutating context destroys cache efficiency",
            "Never include timestamps precise to the second in prompts; move all dynamic content to the end",
        ]
    )
    
    # Slide 9: Key Takeaways
    add_content_slide(
        prs,
        "Key Takeaways",
        [
            "With Claude Opus 4.6, cached input tokens cost 90% less than uncached (10x cost difference)",
            "The difference between $0.50 and $5.00 per query often comes down to context management",
            "Stable prefixes for KV cache hits is the single most important optimization for production agents",
            "Output tokens cost 5x more than uncached inputs — budget them carefully",
        ]
    )
    
    # Slide 10: Source
    add_content_slide(
        prs,
        "Source",
        [
            "Nicolas Bustamante — The LLM Context Tax: Best Tips for Tax Avoidance",
            "https://www.nicolasbustamante.com/p/the-llm-context-tax-best-tips-for",
        ]
    )
    
    output_path = "/workspace/LLM_Information_Efficiency_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    main()
