"""
Generate a PowerPoint presentation on Agentic AI using content from
Databricks Data + AI Summit events (2024, 2025, 2026), company filings,
and competitive analysis.
Dark theme with accent colors matching the Databricks visual identity.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1B, 0x1F, 0x2B)
ACCENT_ORANGE = RGBColor(0xFF, 0x6F, 0x00)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x3D)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)
ACCENT_PURPLE = RGBColor(0xA7, 0x55, 0xF7)
ACCENT_TEAL = RGBColor(0x14, 0xB8, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
CARD_BG = RGBColor(0x25, 0x2A, 0x3A)
MUTED = RGBColor(0x70, 0x78, 0x90)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_info.get("text", "")
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.alignment = line_info.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_info.get("space_after", 6))
    return txBox


def add_stat_card(slide, left, top, width, height, stat_value, stat_label, accent_color):
    add_card(slide, left, top, width, height)
    add_accent_bar(slide, left, top, width, Pt(4), accent_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4),
                 Inches(0.6), stat_value, font_size=28, color=accent_color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), width - Inches(0.4),
                 height - Inches(0.9), stat_label, font_size=12, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)


def slide_header(slide, title, bar_color=ACCENT_ORANGE, bar_width=Inches(2)):
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 title, font_size=36, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.15), bar_width, Pt(4), bar_color)


# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── SLIDE 1: TITLE ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_accent_bar(slide, Inches(1), Inches(3.0), Inches(1.5), Pt(5), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
             "DATABRICKS DATA + AI SUMMIT SERIES", font_size=18, color=ACCENT_ORANGE,
             bold=True)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Agentic AI", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             "Enterprise Adoption, Architecture & Production Insights from\n"
             "Databricks Data + AI Summit 2024, 2025 & 2026",
             font_size=20, color=LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
             "Sourced from summit keynotes, sessions, and the State of AI Agents report",
             font_size=14, color=MUTED)

stats = [
    ("20,000+", "Databricks Customers\nWorldwide", ACCENT_ORANGE),
    ("$134B", "Valuation\n(Latest Raise)", ACCENT_BLUE),
    ("60%+", "Fortune 500\non Platform", ACCENT_GREEN),
    ("800+", "Sessions at\nSummit 2026", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(stats):
    add_stat_card(slide, Inches(1) + i * Inches(2.85), Inches(5.2), Inches(2.6),
                  Inches(1.6), val, label, color)


# ── SLIDE 2: AGENDA ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Agenda")

agenda_items = [
    ("01", "What Is Agentic AI?",
     "Definition, evolution from chatbots to autonomous multi-agent systems",
     ACCENT_ORANGE),
    ("02", "Databricks' Agentic AI Platform",
     "Mosaic AI, Agent Bricks, MLflow 3.0, Lakebase, and MCP integration",
     ACCENT_BLUE),
    ("03", "Summit Sessions Across Three Years",
     "Key agentic AI presentations from 2024, 2025, and 2026 summits",
     ACCENT_GREEN),
    ("04", "Enterprise Case Studies & State of AI Agents",
     "adidas, Walmart, AstraZeneca, 7-Eleven, Zillow, RBC + governance metrics",
     ACCENT_PURPLE),
    ("05", "AI Agent Usage in Company Filings",
     "Revenue, adoption, and agent metrics from SEC filings and earnings calls",
     ACCENT_TEAL),
    ("06", "Private AI Companies & Agents",
     "OpenAI, Anthropic, Cursor, Perplexity, xAI, Mistral, Cohere, Cognition (Devin)",
     ACCENT_RED),
    ("07", "Competitive Landscape & Takeaways",
     "Databricks positioning, token economics, architecture, and strategic implications",
     RGBColor(0xF5, 0xA6, 0x23)),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(1.55) + i * Inches(0.8)
    add_accent_bar(slide, Inches(0.8), y, Pt(4), Inches(0.65), color)
    add_text_box(slide, Inches(1.1), y, Inches(0.6), Inches(0.65), num,
                 font_size=22, color=color, bold=True)
    add_text_box(slide, Inches(1.8), y - Inches(0.02), Inches(5), Inches(0.35), title,
                 font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.3), Inches(8), Inches(0.35), desc,
                 font_size=13, color=LIGHT_GRAY)


# ── SLIDE 3: WHAT IS AGENTIC AI? ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "What Is Agentic AI?", bar_color=ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "From chatbots to autonomous systems — the evolution presented across "
             "Databricks summits",
             font_size=14, color=LIGHT_GRAY)

levels = [
    ("Level 0: Chatbot", "Single-turn Q&A, no tools or memory",
     "Basic FAQ bots", ACCENT_RED),
    ("Level 1: Copilot", "Human-directed, tool-assisted responses",
     "GitHub Copilot, AI/BI Genie", ACCENT_ORANGE),
    ("Level 2: Task Agent", "Autonomous single-task execution with tool use",
     "RAG agents, coding agents", ACCENT_BLUE),
    ("Level 3: Multi-Agent", "Multiple agents collaborating via orchestration",
     "Supervisor + specialist patterns", ACCENT_GREEN),
    ("Level 4: Autonomous", "Fully autonomous, long-running systems",
     "Digital workers, agent fleets", ACCENT_PURPLE),
]
for i, (stage, desc, example, color) in enumerate(levels):
    y = Inches(2.0) + i * Inches(1.0)
    card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(0.85))
    add_accent_bar(slide, Inches(0.8), y, Pt(5), Inches(0.85), color)
    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(3), Inches(0.35),
                 stage, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.08), Inches(4.5), Inches(0.7),
                 desc, font_size=13, color=WHITE)
    add_text_box(slide, Inches(9.2), y + Inches(0.08), Inches(3), Inches(0.7),
                 example, font_size=12, color=LIGHT_GRAY)

add_multi_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4), [
    {"text": "Source: Composite maturity model — Databricks Summit, Google Cloud Next, "
             "McKinsey, Sequoia Capital", "size": 10, "color": MUTED},
])


# ── SLIDE 4: DATABRICKS CEO VISION ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Databricks CEO Vision: The Agentic Era",
             bar_color=ACCENT_ORANGE)

card = add_card(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.6))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Pt(5), Inches(1.6), ACCENT_ORANGE)
add_multi_text(slide, Inches(1.3), Inches(1.6), Inches(10.5), Inches(1.4), [
    {"text": '"It\'s very elusive to get AI that really works and understands that '
             'proprietary data', "size": 18, "color": WHITE, "space_after": 2},
    {"text": 'that\'s inside enterprise... this could maybe get all the way to a '
             'trillion."', "size": 18, "color": WHITE, "space_after": 8},
    {"text": "— Ali Ghodsi, CEO & Co-Founder, Databricks  |  Fortune Brainstorm AI, "
             "December 2025", "size": 13, "color": ACCENT_ORANGE, "bold": True},
])

pillars = [
    ("Enterprise Data + AI Agents",
     "General AI knowledge is commoditized — agents with proprietary enterprise data "
     "create unique, defensible value. Databricks positions itself as the data "
     "intelligence platform for this era.",
     ACCENT_ORANGE),
    ("AI Agents Building Infrastructure",
     "80% of databases on the Databricks platform are now created by AI agents, not "
     "humans. This drove the launch of Lakebase — a Postgres-compatible database "
     "built for agent-scale workloads.",
     ACCENT_BLUE),
    ("Open Platform Strategy",
     "Cloud-agnostic, open-source approach (MLflow, Unity Catalog, Apache Spark) "
     "attracts enterprises. Customers choose flexibility across AWS, GCP, Azure, "
     "and on-premise deployments.",
     ACCENT_GREEN),
]
for i, (title, desc, color) in enumerate(pillars):
    x = Inches(0.8) + i * Inches(3.95)
    y = Inches(3.5)
    card = add_card(slide, x, y, Inches(3.7), Inches(2.6))
    add_accent_bar(slide, x, y, Inches(3.7), Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.4), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.4), Inches(2.0),
                 desc, font_size=12, color=LIGHT_GRAY)

growth_stats = [
    ("$134B", "Latest Valuation", ACCENT_ORANGE),
    ("80%", "DBs Created\nby AI Agents", ACCENT_BLUE),
    ("60%+", "Fortune 500\non Platform", ACCENT_GREEN),
    ("327%", "Multi-Agent\nGrowth (4 mo)", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(growth_stats):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(6.3), Inches(2.7),
                  Inches(1.0), val, label, color)


# ── SLIDE 5: MOSAIC AI AGENT PLATFORM ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Databricks Mosaic AI: The Agent Platform",
             bar_color=ACCENT_BLUE)

card = add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.7))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(5.5), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Compound AI Systems", font_size=18, color=ACCENT_ORANGE, bold=True)
add_multi_text(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(2.0), [
    {"text": "•  Tuned foundation models + fine-tuning", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  Retrieval-Augmented Generation (RAG)", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  Tool use via Model Context Protocol (MCP)", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  Reasoning agents with multi-step planning", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  Multi-agent orchestration (supervisor pattern)", "size": 13,
     "color": WHITE, "space_after": 4},
    {"text": "Production systems use multiple components rather", "size": 12,
     "color": LIGHT_GRAY, "space_after": 2},
    {"text": "than monolithic models.", "size": 12, "color": LIGHT_GRAY},
])

card = add_card(slide, Inches(6.6), Inches(1.5), Inches(5.5), Inches(2.7))
add_accent_bar(slide, Inches(6.6), Inches(1.5), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(1.65), Inches(5), Inches(0.4),
             "MCP & Tool Calling", font_size=18, color=ACCENT_BLUE, bold=True)
add_multi_text(slide, Inches(6.9), Inches(2.1), Inches(5), Inches(2.0), [
    {"text": "•  Managed MCP servers with Unity Catalog", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  UC Functions, Genie spaces, Vector Search", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  Parallel tool calling on Model Serving", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  On-behalf-of-user authentication", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "•  Custom MCP servers via Databricks Apps", "size": 13, "color": WHITE,
     "space_after": 4},
    {"text": "Standard protocol for equipping agents with", "size": 12,
     "color": LIGHT_GRAY, "space_after": 2},
    {"text": "enterprise tools and data access.", "size": 12, "color": LIGHT_GRAY},
])

card = add_card(slide, Inches(0.8), Inches(4.5), Inches(3.6), Inches(2.5))
add_accent_bar(slide, Inches(0.8), Inches(4.5), Inches(3.6), Pt(4), ACCENT_GREEN)
add_text_box(slide, Inches(1.1), Inches(4.65), Inches(3.2), Inches(0.4),
             "Agent Bricks (No-Code)", font_size=16, color=ACCENT_GREEN, bold=True)
add_multi_text(slide, Inches(1.1), Inches(5.05), Inches(3.2), Inches(1.8), [
    {"text": "•  Auto-generates domain agents", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  High-level task descriptions", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Auto-builds evaluations", "size": 12, "color": WHITE, "space_after": 3},
    {"text": "•  Optimizes for quality & cost", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Supports multi-agent systems", "size": 12, "color": WHITE},
])

card = add_card(slide, Inches(4.65), Inches(4.5), Inches(3.6), Inches(2.5))
add_accent_bar(slide, Inches(4.65), Inches(4.5), Inches(3.6), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(4.95), Inches(4.65), Inches(3.2), Inches(0.4),
             "MLflow 3.0 (Observability)", font_size=16, color=ACCENT_PURPLE,
             bold=True)
add_multi_text(slide, Inches(4.95), Inches(5.05), Inches(3.2), Inches(1.8), [
    {"text": "•  Agent tracing & monitoring", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Prompt versioning & registry", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Cross-platform (AWS/GCP/prem)", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Unity Catalog integration", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Open-source, cloud-agnostic", "size": 12, "color": WHITE},
])

card = add_card(slide, Inches(8.5), Inches(4.5), Inches(3.6), Inches(2.5))
add_accent_bar(slide, Inches(8.5), Inches(4.5), Inches(3.6), Pt(4), ACCENT_TEAL)
add_text_box(slide, Inches(8.8), Inches(4.65), Inches(3.2), Inches(0.4),
             "Lakebase (Agent DB)", font_size=16, color=ACCENT_TEAL, bold=True)
add_multi_text(slide, Inches(8.8), Inches(5.05), Inches(3.2), Inches(1.8), [
    {"text": "•  Postgres-compatible, managed", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Sub-10ms latency", "size": 12, "color": WHITE, "space_after": 3},
    {"text": "•  10,000+ queries per second", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Zero-ETL with copy-on-write", "size": 12, "color": WHITE,
     "space_after": 3},
    {"text": "•  Built for agent-scale workloads", "size": 12, "color": WHITE},
])


# ── SLIDE 6: SUMMIT 2024 HIGHLIGHTS ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Data + AI Summit 2024: Agentic AI Highlights",
             bar_color=ACCENT_ORANGE, bar_width=Inches(3))

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "16,000+ in-person  |  60,000+ virtual  |  140 countries",
             font_size=13, color=LIGHT_GRAY)

sessions_2024 = [
    ("Mosaic AI: Compound AI Systems",
     "Databricks introduced Mosaic AI for building production-quality compound AI "
     "systems. The paradigm shifted from monolithic models to multi-component systems "
     "combining tuned models, retrieval, tool use, and reasoning agents.",
     ["New SDK for building, deploying & evaluating AI agents",
      "Foundation Model APIs with pay-per-token pricing",
      "Mosaic AI Tool Catalog for enterprise function registries",
      "Model Serving support for agents and RAG applications"],
     ACCENT_ORANGE),
    ("GenAI Announcements Keynote",
     "Focus on production-grade compound AI and fine-tuning. Open-source adoption "
     "accelerated rapidly — Llama 3 captured 39% of open-source usage within 4 weeks "
     "of launch.",
     ["Fine-tuning support for foundation models",
      "70% of GenAI companies use vector databases for RAG",
      "Hugging Face Transformers jumped from #4 to #2 adoption",
      "Open-source models preferred by 76% of LLM users"],
     ACCENT_BLUE),
]
for i, (title, desc, bullets, color) in enumerate(sessions_2024):
    y_start = Inches(1.8) + i * Inches(2.7)
    card = add_card(slide, Inches(0.8), y_start, Inches(11.5), Inches(2.5))
    add_accent_bar(slide, Inches(0.8), y_start, Pt(5), Inches(2.5), color)
    add_text_box(slide, Inches(1.2), y_start + Inches(0.15), Inches(10), Inches(0.5),
                 title, font_size=22, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y_start + Inches(0.6), Inches(5), Inches(1.8),
                 desc, font_size=13, color=LIGHT_GRAY)
    bullet_lines = [{"text": f"•  {b}", "size": 13, "color": WHITE, "space_after": 4}
                    for b in bullets]
    add_multi_text(slide, Inches(6.5), y_start + Inches(0.55), Inches(5.5), Inches(2),
                   bullet_lines)


# ── SLIDE 7: SUMMIT 2025 HIGHLIGHTS ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Data + AI Summit 2025: Agentic AI Highlights",
             bar_color=ACCENT_BLUE, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "June 9–12, 2025  |  San Francisco  |  22,000+ Attendees  |  700+ Sessions",
             font_size=13, color=LIGHT_GRAY)

sessions_2025 = [
    ("Autonomous AI Agents in Infrastructure",
     "Apurva Kumar — Walmart Global Tech",
     "Architecture & design principles for autonomous agents: containerization, "
     "orchestration, robustness, and real-time feedback loops.",
     ACCENT_ORANGE),
    ("Agent Bricks: Multi-Agent Systems",
     "Databricks",
     "Natural language interactions with structured databases and unstructured "
     "documents. KBQA evaluation, data governance, and vector search.",
     ACCENT_BLUE),
    ("Multi-Agents in Production",
     "Tony Holdstock-Brown — Inngest",
     "Techniques for building effective multi-agent networks deterministically "
     "while reducing compounding error rates across AI calls.",
     ACCENT_GREEN),
    ("Quality Monitoring for Agents",
     "Databricks",
     "Mosaic Agent Monitoring: user feedback integration, monitoring dashboards, "
     "quality evaluation using AI judges, and operational metrics.",
     ACCENT_PURPLE),
    ("Building Responsible AI Agents",
     "Databricks",
     "Fairness, transparency, regulatory compliance. Bias monitoring, "
     "explainability (SHAP/LIME), Unity Catalog governance, and LLM guardrails.",
     ACCENT_TEAL),
    ("Agentic AI in Finance",
     "Qubika (Sponsored)",
     "AI finance agents using LangChain, RAG, and Databricks unified platform "
     "including MLflow and Mosaic AI integration.",
     ACCENT_RED),
]

card_w = Inches(3.6)
card_h = Inches(2.2)
for i, (title, speaker, desc, color) in enumerate(sessions_2025):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (card_w + Inches(0.2))
    y = Inches(1.8) + row * (card_h + Inches(0.2))
    add_card(slide, x, y, card_w, card_h)
    add_accent_bar(slide, x, y, card_w, Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3),
                 Inches(0.5), title, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), card_w - Inches(0.3),
                 Inches(0.3), speaker, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.95), card_w - Inches(0.3),
                 Inches(1.1), desc, font_size=11, color=LIGHT_GRAY)


# ── SLIDE 8: SUMMIT 2026 HIGHLIGHTS ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Data + AI Summit 2026: Agentic AI Highlights",
             bar_color=ACCENT_GREEN, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "June 15–18, 2026  |  San Francisco  |  800+ Sessions Expected",
             font_size=13, color=LIGHT_GRAY)

sessions_2026 = [
    ("Beyond the Trace: adidas' Agent Digital Twin",
     "adidas",
     "Maps agent→tool→prompt→retrieval→model call→post-processing chains. "
     "Rollup of risk, quality, and unit economics across the fleet. "
     "Uses MLflow Tracing + Unity Catalog for audit evidence.",
     ["200+ serving endpoints", "300+ data & AI products",
      "6,000+ registered models", "600,000+ pipeline runs"],
     ACCENT_ORANGE),
    ("Deploying and Monitoring Agents",
     "Databricks",
     "Hands-on course covering batch (AI Functions) and real-time (Model Serving) "
     "agent deployment. MLflow tracing, monitoring dashboards, and observability.",
     ["Batch mode via AI Functions", "Real-time REST endpoints",
      "Built-in trace collection", "Monitoring dashboards"],
     ACCENT_BLUE),
    ("Building RAG Agents with Agent Bricks",
     "Databricks",
     "Hands-on training using Vector Search, AI Playground, and Agent Bricks "
     "for building production RAG agents.",
     ["Vector Search integration", "AI Playground testing",
      "Auto-optimized agents", "No-code building"],
     ACCENT_GREEN),
    ("Building Tool-Calling Agents with MCP",
     "Databricks",
     "Creating agents with tool-calling capabilities using the Databricks Agent "
     "Framework and Model Context Protocol (MCP).",
     ["Model Context Protocol", "Tool-calling capabilities",
      "Agent Framework SDK", "Enterprise integration"],
     ACCENT_PURPLE),
    ("Scaling Agentic AI with Lakebase",
     "Zillow",
     "Agentic AI production strategy using Databricks Lakebase and enterprise data "
     "mesh patterns to support autonomous AI workloads at scale.",
     ["Production agentic AI", "Enterprise data mesh",
      "Lakebase integration", "Scalable architecture"],
     ACCENT_TEAL),
    ("AI Agents for the Frontline",
     "7-Eleven",
     "GenAI maintenance assistant for frontline store operations — demonstrating "
     "agentic AI applied outside traditional knowledge-worker domains.",
     ["Frontline operations", "Maintenance assistant",
      "Real-world deployment", "Non-desk workers"],
     ACCENT_RED),
]

card_w = Inches(3.7)
card_h = Inches(2.5)
for i, (title, company, desc, bullets, color) in enumerate(sessions_2026):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (card_w + Inches(0.15))
    y = Inches(1.8) + row * (card_h + Inches(0.15))
    add_card(slide, x, y, card_w, card_h)
    add_accent_bar(slide, x, y, card_w, Pt(4), color)
    add_text_box(slide, x + Inches(0.12), y + Inches(0.12), card_w - Inches(0.24),
                 Inches(0.45), title, font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.12), y + Inches(0.48), Inches(0.8), Inches(0.25),
                 company, font_size=10, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.12), y + Inches(0.75), Inches(2.1), Inches(1.3),
                 desc, font_size=10, color=LIGHT_GRAY)
    bullet_lines = [{"text": f"•  {b}", "size": 10, "color": WHITE, "space_after": 2}
                    for b in bullets]
    add_multi_text(slide, x + Inches(2.3), y + Inches(0.75), Inches(1.3), Inches(1.3),
                   bullet_lines)


# ── SLIDE 9: STATE OF AI AGENTS REPORT ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "State of AI Agents Report (2026)", bar_color=ACCENT_TEAL)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Based on 20,000+ global customers  |  Published January 2026",
             font_size=13, color=LIGHT_GRAY)

report_stats = [
    ("327%", "Multi-Agent Workflow\nGrowth in 4 Months", ACCENT_ORANGE),
    ("80%", "New Databases Built\nby AI Agents", ACCENT_BLUE),
    ("19%", "Organizations with Agents\nDeployed at Scale", ACCENT_GREEN),
    ("37%", "Use Supervisor Agent\nArchitecture", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(report_stats):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(1.8), Inches(2.7),
                  Inches(1.6), val, label, color)

# Governance Impact
card = add_card(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(3.2))
add_accent_bar(slide, Inches(0.8), Inches(3.8), Inches(5.5), Pt(4), ACCENT_TEAL)
add_text_box(slide, Inches(1.1), Inches(3.95), Inches(5), Inches(0.4),
             "Governance Impact on Production Success", font_size=18,
             color=ACCENT_TEAL, bold=True)
add_multi_text(slide, Inches(1.1), Inches(4.4), Inches(5), Inches(2.4), [
    {"text": "12x", "size": 36, "color": ACCENT_ORANGE, "bold": True, "space_after": 2},
    {"text": "More AI projects in production for companies using governance tools",
     "size": 14, "color": WHITE, "space_after": 16},
    {"text": "6x", "size": 36, "color": ACCENT_BLUE, "bold": True, "space_after": 2},
    {"text": "More production deployments for companies using evaluation tools",
     "size": 14, "color": WHITE, "space_after": 16},
    {"text": "7x", "size": 36, "color": ACCENT_GREEN, "bold": True, "space_after": 2},
    {"text": "Growth in AI Gateway governance product usage in 9 months",
     "size": 14, "color": WHITE},
])

# Use Cases & Model Diversity
card = add_card(slide, Inches(6.6), Inches(3.8), Inches(5.5), Inches(3.2))
add_accent_bar(slide, Inches(6.6), Inches(3.8), Inches(5.5), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(6.9), Inches(3.95), Inches(5), Inches(0.4),
             "Use Cases & Model Diversity", font_size=18, color=ACCENT_PURPLE,
             bold=True)
add_multi_text(slide, Inches(6.9), Inches(4.4), Inches(5), Inches(2.4), [
    {"text": "40%  Customer Experience", "size": 16, "color": ACCENT_ORANGE,
     "bold": True, "space_after": 4},
    {"text": "       Support, onboarding, and engagement agents", "size": 12,
     "color": LIGHT_GRAY, "space_after": 10},
    {"text": "35%  Predictive Maintenance", "size": 16, "color": ACCENT_BLUE,
     "bold": True, "space_after": 4},
    {"text": "       Manufacturing & automotive monitoring", "size": 12,
     "color": LIGHT_GRAY, "space_after": 10},
    {"text": "23%  Medical Literature Synthesis", "size": 16, "color": ACCENT_GREEN,
     "bold": True, "space_after": 4},
    {"text": "       Healthcare document processing", "size": 12, "color": LIGHT_GRAY,
     "space_after": 10},
    {"text": "77%  Use 2+ model families  |  59% use 3+", "size": 14,
     "color": ACCENT_TEAL, "bold": True, "space_after": 4},
    {"text": "       Multi-model strategy is the enterprise norm", "size": 12,
     "color": LIGHT_GRAY},
])


# ── SLIDE 10: ENTERPRISE CASE STUDIES ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Enterprise Case Studies from Databricks Summits",
             bar_color=ACCENT_PURPLE, bar_width=Inches(3))

companies = [
    ("adidas", "Agent Digital Twin for Governance & Cost",
     ["200+ serving endpoints, 300+ data & AI products",
      "98.5% token efficiency (200K → 3K input tokens)",
      "91.67% cost savings through LLM optimization",
      "Per-hop unit economics and cost leak forensics",
      "MLflow Tracing + Unity Catalog for audit evidence"],
     ACCENT_ORANGE),
    ("Walmart", "Autonomous AI Agents & Analytics",
     ["90% reduction in time-to-value with AI/BI Genie",
      "$5.6M annual savings from FTE hours",
      "AI chatbot 'Sparky' driving increased spend",
      "Conversational AI agents for store managers",
      "Autonomous agent architecture (Summit 2025)"],
     ACCENT_BLUE),
    ("AstraZeneca", "Clinical Trial Document Processing",
     ["400,000+ documents reviewed in under 60 minutes",
      "No coding required using Agent Bricks",
      "Surpassed open-source alternatives in accuracy",
      "Production-grade agent with enterprise data"],
     ACCENT_GREEN),
    ("7-Eleven", "GenAI Frontline Maintenance Assistant",
     ["AI agents for non-desk frontline workers",
      "Store maintenance and operations assistance",
      "Real-world deployment beyond knowledge workers",
      "Presented at Summit 2026"],
     ACCENT_PURPLE),
    ("Zillow", "Agentic AI & Enterprise Data Mesh",
     ["Production agentic AI strategy with Lakebase",
      "Enterprise data mesh for autonomous workloads",
      "Scaling real estate AI agents at production volume",
      "Presented at Summit 2026"],
     ACCENT_TEAL),
    ("Royal Bank of Canada", "AI Agents for Equity Research",
     ["Reduced research work from days to minutes",
      "Agents work with proprietary financial data",
      "Highlighted by CEO as key enterprise success",
      "Demonstrates value of agents + proprietary data"],
     ACCENT_RED),
]

card_w = Inches(3.7)
card_h = Inches(2.7)
for i, (company, subtitle, bullets, color) in enumerate(companies):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (card_w + Inches(0.2))
    y = Inches(1.5) + row * (card_h + Inches(0.15))
    add_card(slide, x, y, card_w, card_h)
    add_accent_bar(slide, x, y, card_w, Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3),
                 Inches(0.35), company, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.45), card_w - Inches(0.3),
                 Inches(0.3), subtitle, font_size=11, color=WHITE, bold=True)
    bullet_lines = [{"text": f"•  {b}", "size": 10, "color": LIGHT_GRAY,
                     "space_after": 2} for b in bullets]
    add_multi_text(slide, x + Inches(0.15), y + Inches(0.8), card_w - Inches(0.3),
                   card_h - Inches(1.0), bullet_lines)


# ── SLIDE 11: TOKEN ECONOMICS & COST ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Token Economics & Cost Optimization",
             bar_color=ACCENT_GREEN)

card = add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(5.5), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Databricks Inference Pricing (DBRX)", font_size=18, color=ACCENT_ORANGE,
             bold=True)
add_multi_text(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(1.8), [
    {"text": "Input Tokens:     $0.0005 / 1K tokens", "size": 15, "color": WHITE,
     "bold": True, "space_after": 8},
    {"text": "Output Tokens:   $0.0015 / 1K tokens", "size": 15, "color": WHITE,
     "bold": True, "space_after": 8},
    {"text": "Agent Evaluation:  $0.018 / judge request", "size": 15, "color": WHITE,
     "bold": True, "space_after": 8},
    {"text": "(5,000 tokens per judge request increment)", "size": 12,
     "color": LIGHT_GRAY},
])

card = add_card(slide, Inches(6.6), Inches(1.5), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(6.6), Inches(1.5), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(1.65), Inches(5), Inches(0.4),
             "Production Agent Costs (Monthly)", font_size=18, color=ACCENT_BLUE,
             bold=True)
add_multi_text(slide, Inches(6.9), Inches(2.1), Inches(5), Inches(1.8), [
    {"text": "Inference (1M convos):  $3,000 – $7,000", "size": 15, "color": WHITE,
     "bold": True, "space_after": 8},
    {"text": "Base Compute:                $5,000 – $15,000", "size": 15,
     "color": WHITE, "bold": True, "space_after": 8},
    {"text": "+ Storage, networking, and monitoring overhead", "size": 13,
     "color": LIGHT_GRAY, "space_after": 4},
    {"text": "AI Functions: 3x faster, 4x lower cost", "size": 13,
     "color": ACCENT_GREEN, "bold": True},
])

stats_token = [
    ("98.5%", "Token Efficiency Achieved\n(adidas: 200K → 3K tokens)", ACCENT_ORANGE),
    ("91.67%", "Cost Savings Through\nToken Optimization", ACCENT_BLUE),
    ("80%", "New Databases Created\nby AI Agents", ACCENT_GREEN),
    ("97%", "Test & Dev Environments\nCreated by Agents", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(stats_token):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(4.4), Inches(2.7),
                  Inches(1.6), val, label, color)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "Gartner predicts >40% of agentic AI projects will be canceled by "
             "end-2027 as costs rise and controls lag.",
             font_size=12, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
             "Sources: Databricks Summit Sessions, adidas Summit Presentation, "
             "Getmonetizely Analysis",
             font_size=10, color=MUTED)


# ── SLIDE 12: AGENTIC AI ARCHITECTURE PATTERNS ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Agentic AI Architecture Patterns",
             bar_color=ACCENT_RED, bar_width=Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.35), Inches(10), Inches(0.4),
             "Production architecture patterns highlighted across Databricks summits",
             font_size=13, color=LIGHT_GRAY)

card = add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(5.5), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.4),
             "Compound AI System (Databricks Pattern)", font_size=16,
             color=ACCENT_ORANGE, bold=True)

arch_lines = [
    "┌─────────────────────────────┐",
    "│       USER / TRIGGER        │",
    "└──────────────┬──────────────┘",
    "               │",
    "               ▼",
    "┌─────────────────────────────┐",
    "│    AGENT ORCHESTRATOR       │",
    "│ [Plan] [Memory] [Guardrail]│",
    "└──────────────┬──────────────┘",
    "         ┌─────┼─────┐",
    "         ▼     ▼     ▼",
    "      [RAG] [Tool] [Code]",
    "       ↕     ↕     ↕",
    "┌─────────────────────────────┐",
    "│  DATA LAYER (Lakehouse)     │",
    "│ Unity Catalog │ Lakebase    │",
    "│ Vector Search │ Delta Lake  │",
    "└─────────────────────────────┘",
    "         ┌─────┴─────┐",
    "    [MLflow 3.0] [AI Gateway]",
    "     Monitoring    Governance",
]
txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(5), Inches(4.2))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(arch_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = "Courier New"
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.CENTER

# Right side: patterns
card = add_card(slide, Inches(6.6), Inches(1.8), Inches(5.5), Inches(2.3))
add_accent_bar(slide, Inches(6.6), Inches(1.8), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(1.95), Inches(5), Inches(0.4),
             "Agent Orchestration Patterns", font_size=16, color=ACCENT_BLUE,
             bold=True)
add_multi_text(slide, Inches(6.9), Inches(2.4), Inches(5), Inches(1.5), [
    {"text": "•  Supervisor: Manager agent delegates to specialists (37% adoption)",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  Sequential: Agents work in defined pipeline order",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  Parallel: Concurrent tool calls via managed MCP",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  ReAct: Reasoning + Acting loop with self-correction",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  Digital Twin: Agent mirroring for governance (adidas)",
     "size": 12, "color": WHITE},
])

card = add_card(slide, Inches(6.6), Inches(4.3), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(6.6), Inches(4.3), Inches(5.5), Pt(4), ACCENT_GREEN)
add_text_box(slide, Inches(6.9), Inches(4.45), Inches(5), Inches(0.4),
             "Monitoring & Observability", font_size=16, color=ACCENT_GREEN,
             bold=True)
add_multi_text(slide, Inches(6.9), Inches(4.9), Inches(5), Inches(1.7), [
    {"text": "•  End-user feedback for iterative improvement",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  AI judges for continuous quality evaluation",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  Token usage tracking via AI Gateway system tables",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  Per-hop cost forensics for agent chains",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "•  MLflow 3.0 tracing across all deployment targets",
     "size": 12, "color": WHITE},
])


# ── SLIDE 13: COMPANY FILINGS — AI AGENT REVENUE & ADOPTION ───────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "AI Agent Usage in Company Filings (Part 1)",
             bar_color=ACCENT_TEAL, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Metrics from SEC filings, 10-K reports, and public earnings calls (2025–2026)",
             font_size=13, color=LIGHT_GRAY)

filing_companies_1 = [
    ("Salesforce — Agentforce", "Q4 FY2026 Earnings (Feb 2026)",
     [("$800M", "Agentforce ARR\n(+169% YoY)", ACCENT_ORANGE),
      ("29,000+", "Deals Closed\n(+50% QoQ)", ACCENT_BLUE),
      ("2.4B", "Agentic Work\nUnits Delivered", ACCENT_GREEN)],
     ["FY26 revenue $41.5B (+10% YoY); Agentforce + Data 360 ARR >$2.9B (+200% YoY)",
      "Introduced 'Agentic Work Units' (AWUs) as new metric — 2.4B delivered, +57% QoQ",
      "Processed 19T+ tokens (+5x YoY); agents in production +50% QoQ",
      'CEO Benioff: "Agentforce just became an $800 million business"',
      "Digital labor model: agents billed per conversation, not per seat"],
     ACCENT_ORANGE),
    ("Microsoft — Copilot & Agent 365", "FY2025 Annual Report / 10-K",
     [("90%+", "Fortune 500 Use\nM365 Copilot", ACCENT_BLUE),
      ("$281.7B", "Total Revenue\n(+15% YoY)", ACCENT_PURPLE),
      ("400+", "New Copilot\nFeatures Shipped", ACCENT_GREEN)],
     ['CEO Nadella: "Every customer I talk to is looking to reshape processes with AI agents"',
      '"Frontier Firm" concept: human-led, agent-operated organizations',
      "Agent 365 control plane for managing and securing agents at enterprise scale",
      "Multi-agent orchestration: agents collaborating on complex work together",
      "Copilot Studio Tuning: low-code model training with company data"],
     ACCENT_BLUE),
]

y_offset = Inches(1.8)
for comp_name, filing_ref, stat_cards, bullets, color in filing_companies_1:
    card = add_card(slide, Inches(0.8), y_offset, Inches(11.5), Inches(2.6))
    add_accent_bar(slide, Inches(0.8), y_offset, Pt(5), Inches(2.6), color)
    add_text_box(slide, Inches(1.2), y_offset + Inches(0.08), Inches(5), Inches(0.35),
                 comp_name, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(7.0), y_offset + Inches(0.08), Inches(5), Inches(0.35),
                 filing_ref, font_size=12, color=MUTED, bold=False,
                 alignment=PP_ALIGN.RIGHT)

    for j, (sv, sl_label, sc) in enumerate(stat_cards):
        sx = Inches(1.2) + j * Inches(1.85)
        add_stat_card(slide, sx, y_offset + Inches(0.45), Inches(1.7), Inches(0.9),
                      sv, sl_label, sc)

    bullet_lines = [{"text": f"•  {b}", "size": 10, "color": LIGHT_GRAY,
                     "space_after": 2} for b in bullets]
    add_multi_text(slide, Inches(6.9), y_offset + Inches(0.45), Inches(5.2),
                   Inches(2.0), bullet_lines)
    y_offset += Inches(2.75)


# ── SLIDE 14: COMPANY FILINGS — AI AGENT REVENUE & ADOPTION (Part 2) ─────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "AI Agent Usage in Company Filings (Part 2)",
             bar_color=ACCENT_TEAL, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Metrics from SEC filings, 10-K reports, and public earnings calls (2025–2026)",
             font_size=13, color=LIGHT_GRAY)

filing_companies_2 = [
    ("ServiceNow — Now Assist Agents", "Q4 2025 Earnings / 10-K",
     [("$600M+", "Now Assist ACV", ACCENT_GREEN),
      ("2x", "Net New ACV\nYoY Growth", ACCENT_BLUE),
      ("35", "Deals >$1M\nin Q4", ACCENT_PURPLE)],
     ["Q4 2025 subscription revenue $3.47B (+21% YoY); total revenue $3.57B (+20.5%)",
      "603 customers with >$5M ACV (+20% YoY); 244 deals >$1M net new ACV (+40% YoY)",
      "Launched thousands of pre-built AI agents + AI Agent Studio (no-code)",
      "AI Agent Orchestrator to manage agent fleets across the enterprise"],
     ACCENT_GREEN),
    ("AWS — Bedrock Agents", "Q4 2025 Earnings (Amazon)",
     [("$128.7B", "AWS Full-Year\nRevenue", ACCENT_ORANGE),
      ("100K+", "Bedrock\nOrganizations", ACCENT_TEAL),
      ("60%", "Customer Spend\nGrowth QoQ", ACCENT_BLUE)],
     ["AWS revenue $35.6B in Q4 (+24% YoY) — fastest growth in 13 quarters",
      "Bedrock: multi-billion dollar run-rate; customer spend +60% quarter-over-quarter",
      "Bedrock AgentCore launched: secure deployment, governance, monitoring at scale",
      "CEO Jassy: ~$200B planned capex for 2026, largely in AWS AI infrastructure"],
     ACCENT_TEAL),
    ("Palantir — AIP Agents", "Q4 2025 Earnings / 10-K",
     [("$4.5B", "Total Revenue\n(+56% YoY)", ACCENT_PURPLE),
      ("137%", "US Commercial\nQ4 Growth YoY", ACCENT_RED),
      ("50%", "Adjusted Op.\nMargin", ACCENT_ORANGE)],
     ["U.S. commercial revenue $1.47B (+109% YoY); Q4 alone $507M (+137% YoY)",
      "Record $4.26B total contract value (+138% YoY); Rule of 40 score: 127%",
      "Q4: 180 deals >$1M, 84 >$5M, 61 >$10M — AIP platform driving expansion",
      'CEO Karp: AI "has just put gasoline on all the tribal knowledge in our products"'],
     ACCENT_PURPLE),
]

y_offset = Inches(1.7)
for comp_name, filing_ref, stat_cards, bullets, color in filing_companies_2:
    ch = Inches(1.75)
    card = add_card(slide, Inches(0.8), y_offset, Inches(11.5), ch)
    add_accent_bar(slide, Inches(0.8), y_offset, Pt(5), ch, color)
    add_text_box(slide, Inches(1.2), y_offset + Inches(0.06), Inches(4.5), Inches(0.3),
                 comp_name, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(7.0), y_offset + Inches(0.06), Inches(5), Inches(0.3),
                 filing_ref, font_size=11, color=MUTED, alignment=PP_ALIGN.RIGHT)

    for j, (sv, sl_label, sc) in enumerate(stat_cards):
        sx = Inches(1.2) + j * Inches(1.75)
        add_stat_card(slide, sx, y_offset + Inches(0.38), Inches(1.6), Inches(0.75),
                      sv, sl_label, sc)

    bullet_lines = [{"text": f"•  {b}", "size": 9, "color": LIGHT_GRAY,
                     "space_after": 1} for b in bullets]
    add_multi_text(slide, Inches(6.6), y_offset + Inches(0.35), Inches(5.5),
                   Inches(1.3), bullet_lines)
    y_offset += Inches(1.85)

add_multi_text(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4), [
    {"text": "Additional filings: Google Cloud $70B+ run rate (+48% YoY in Q4 '25); "
             "NVIDIA $130.5B revenue (+114% YoY); Snowflake AI revenue $100M run rate",
     "size": 10, "color": MUTED},
])


# ── SLIDE 15: COMPETITIVE LANDSCAPE ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Competitive Landscape: Databricks vs. Peers",
             bar_color=ACCENT_RED, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Where Databricks positions in the agentic AI platform market "
             "(data from filings & earnings)",
             font_size=13, color=LIGHT_GRAY)

card = add_card(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.4))
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(11.5), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.35),
             "Data Platform Rivals: Databricks vs. Snowflake", font_size=16,
             color=ACCENT_ORANGE, bold=True)

col_data = [
    ("Dimension", "Databricks", "Snowflake"),
    ("Revenue Run Rate", "$5.4B ARR (~65% YoY growth)", "$4.47B product rev (29% growth)"),
    ("AI-Specific Revenue", "$1.4B AI ARR", "~$100M AI run rate"),
    ("Agent Capabilities", "Mosaic AI, Agent Bricks, MCP", "Cortex AI, Snowflake Intelligence"),
    ("Agent Customers", "2,500+ in Intelligence (3 mo)"),
    ("Valuation / Market Cap", "$134B (private)", "~$59B (NYSE: SNOW)"),
    ("Differentiation", "Lakehouse, ML/AI-native, open formats", "SQL warehouse, BI-optimized"),
]

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.35), Inches(10.8), Inches(1.7))
tf = txBox.text_frame
tf.word_wrap = True
comp_rows = [
    {"text": "Revenue          Databricks $5.4B ARR (~65% YoY)   vs   "
             "Snowflake $4.47B (29% YoY)", "size": 12, "color": WHITE, "space_after": 4},
    {"text": "AI Revenue      Databricks $1.4B AI ARR   vs   "
             "Snowflake ~$100M AI run rate", "size": 12, "color": WHITE, "space_after": 4},
    {"text": "Agent Tools      Mosaic AI + Agent Bricks + MCP   vs   "
             "Cortex AI + Snowflake Intelligence", "size": 12, "color": WHITE,
     "space_after": 4},
    {"text": "Agent Traction  20,000+ customers, 327% multi-agent growth   vs   "
             "2,500+ Intelligence customers (3 months)", "size": 12, "color": WHITE,
     "space_after": 4},
    {"text": "Valuation         $134B (private)   vs   ~$59B (NYSE: SNOW)",
     "size": 12, "color": WHITE, "space_after": 4},
    {"text": "Note: 74% of consulting firms work with both platforms in hybrid "
             "deployments", "size": 11, "color": ACCENT_ORANGE, "bold": True},
]
add_multi_text(slide, Inches(1.1), Inches(2.35), Inches(10.8), Inches(1.7), comp_rows)

# Cloud Hyperscalers
card = add_card(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(2.8))
add_accent_bar(slide, Inches(0.8), Inches(4.4), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(4.55), Inches(5), Inches(0.35),
             "Cloud Hyperscaler Agent Platforms", font_size=16, color=ACCENT_BLUE,
             bold=True)
add_multi_text(slide, Inches(1.1), Inches(4.95), Inches(5), Inches(2.1), [
    {"text": "AWS Bedrock Agents", "size": 13, "color": ACCENT_ORANGE, "bold": True,
     "space_after": 2},
    {"text": "  100K+ orgs; multi-B$ run rate; AgentCore for governance",
     "size": 11, "color": LIGHT_GRAY, "space_after": 6},
    {"text": "Google Cloud Vertex AI", "size": 13, "color": ACCENT_GREEN, "bold": True,
     "space_after": 2},
    {"text": "  $70B+ run rate (+48% YoY); 8M+ paid Gemini Enterprise seats",
     "size": 11, "color": LIGHT_GRAY, "space_after": 6},
    {"text": "Microsoft Azure AI / Copilot", "size": 13, "color": ACCENT_PURPLE,
     "bold": True, "space_after": 2},
    {"text": "  Azure >$75B rev; 90%+ Fortune 500 on M365 Copilot; Agent 365",
     "size": 11, "color": LIGHT_GRAY, "space_after": 6},
    {"text": "Databricks runs on all three clouds — complementary, not competing",
     "size": 11, "color": ACCENT_TEAL, "bold": True},
])

# Enterprise SaaS / Infra Agent Platforms
card = add_card(slide, Inches(6.6), Inches(4.4), Inches(5.5), Inches(2.8))
add_accent_bar(slide, Inches(6.6), Inches(4.4), Inches(5.5), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(6.9), Inches(4.55), Inches(5), Inches(0.35),
             "Enterprise SaaS & Infra Agents", font_size=16, color=ACCENT_PURPLE,
             bold=True)
add_multi_text(slide, Inches(6.9), Inches(4.95), Inches(5), Inches(2.1), [
    {"text": "Salesforce Agentforce", "size": 13, "color": ACCENT_ORANGE, "bold": True,
     "space_after": 2},
    {"text": "  $800M ARR (+169% YoY); 29K deals; digital labor model",
     "size": 11, "color": LIGHT_GRAY, "space_after": 6},
    {"text": "ServiceNow Now Assist", "size": 13, "color": ACCENT_GREEN, "bold": True,
     "space_after": 2},
    {"text": "  $600M+ ACV; 2x net new ACV growth; agent fleet orchestration",
     "size": 11, "color": LIGHT_GRAY, "space_after": 6},
    {"text": "Palantir AIP", "size": 13, "color": ACCENT_PURPLE, "bold": True,
     "space_after": 2},
    {"text": "  $4.5B revenue (+56% YoY); US comm +109%; enterprise AI agents",
     "size": 11, "color": LIGHT_GRAY, "space_after": 6},
    {"text": "NVIDIA Inference Infrastructure", "size": 13, "color": ACCENT_TEAL,
     "bold": True, "space_after": 2},
    {"text": "  $130.5B rev (+114%); Dynamo 40x token throughput; AI factories",
     "size": 11, "color": LIGHT_GRAY},
])

add_multi_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.3), [
    {"text": "Sources: Public SEC filings (10-K), earnings call transcripts, and "
             "investor presentations (2025–2026)", "size": 10, "color": MUTED},
])


# ── SLIDE 16: PRIVATE AI COMPANIES — FOUNDATION MODEL PROVIDERS ───────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Private AI Companies: Foundation Model Providers",
             bar_color=ACCENT_RED, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Privately held companies building agentic AI products and infrastructure",
             font_size=13, color=LIGHT_GRAY)

private_fm = [
    ("OpenAI", "$500B → $1T",
     "$25B ARR (Mar '26)", "Agents SDK, Operator, Swarm",
     [("$25B", "Annualized\nRevenue", ACCENT_ORANGE),
      ("$500B", "Valuation\n(Oct 2025)", ACCENT_BLUE),
      ("$175B", "Total\nFunding", ACCENT_GREEN)],
     ["ARR tripled from $6B (2024) to $20B (2025); hit $25B run rate by Mar 2026",
      "Agents SDK (Mar 2025): tool calls, handoffs, guardrails, built-in tracing",
      "Operator (Jan 2025): autonomous web browser agent with safety controls",
      "Swarm (Oct 2024): multi-agent handoff architecture for orchestration",
      "Compute: 0.2 GW (2023) → 0.6 GW (2024) → 1.9 GW (2025); targeting $600B by 2030"],
     ACCENT_ORANGE),
    ("Anthropic", "$380B",
     "$14B ARR (Feb '26)", "Claude Code, MCP, Computer Use",
     [("$14B", "Revenue\nRun Rate", ACCENT_PURPLE),
      ("$380B", "Valuation\n(Feb 2026)", ACCENT_RED),
      ("$30B", "Series G\nFunding", ACCENT_BLUE)],
     ["Three consecutive years of 10x growth; 80% revenue from enterprise customers",
      "Claude Code: $2.5B ARR, doubled since Jan 2026; 50%+ from enterprise accounts",
      "Created Model Context Protocol (MCP) — now adopted across industry (inc. Databricks)",
      "Computer Use: Claude agents that interact with desktop applications autonomously",
      "300,000+ business and enterprise customers; targeting $20–26B revenue in 2026"],
     ACCENT_PURPLE),
]

y_offset = Inches(1.7)
for comp_name, valuation, rev_label, products, stat_cards, bullets, color in private_fm:
    ch = Inches(2.7)
    card = add_card(slide, Inches(0.8), y_offset, Inches(11.5), ch)
    add_accent_bar(slide, Inches(0.8), y_offset, Pt(5), ch, color)
    add_text_box(slide, Inches(1.2), y_offset + Inches(0.08), Inches(3), Inches(0.3),
                 comp_name, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(4.5), y_offset + Inches(0.08), Inches(4), Inches(0.3),
                 products, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, Inches(9.5), y_offset + Inches(0.08), Inches(2.5), Inches(0.3),
                 f"Valuation: {valuation}", font_size=11, color=MUTED,
                 alignment=PP_ALIGN.RIGHT)

    for j, (sv, sl_label, sc) in enumerate(stat_cards):
        sx = Inches(1.2) + j * Inches(1.85)
        add_stat_card(slide, sx, y_offset + Inches(0.42), Inches(1.7), Inches(0.85),
                      sv, sl_label, sc)

    bullet_lines = [{"text": f"•  {b}", "size": 10, "color": LIGHT_GRAY,
                     "space_after": 2} for b in bullets]
    add_multi_text(slide, Inches(6.9), y_offset + Inches(0.42), Inches(5.2),
                   Inches(2.1), bullet_lines)
    y_offset += Inches(2.85)


# ── SLIDE 17: PRIVATE AI COMPANIES — AGENT-NATIVE PRODUCTS ───────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Private AI Companies: Agent-Native Products",
             bar_color=ACCENT_BLUE, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Companies building agentic AI products for coding, search, and enterprise",
             font_size=13, color=LIGHT_GRAY)

agent_native = [
    ("Cursor", "$29.3B valuation", "AI Coding Agent",
     [("$2B+", "ARR\n(Feb 2026)", ACCENT_ORANGE),
      ("1M+", "Daily Active\nUsers", ACCENT_BLUE),
      ("50%+", "Fortune 500\nAdoption", ACCENT_GREEN)],
     ["Fastest B2B company to $1B ARR (24 months); doubled to $2B in 3 more months",
      "Background agents, multi-file editing, autonomous coding workflows",
      "$2.3B Series D (Nov 2025) led by Accel/Coatue; NVIDIA & Google participated",
      "Enterprise 45-60% of revenue; ~100% of revenue spent on AI infrastructure"],
     ACCENT_ORANGE),
    ("Perplexity", "$20B valuation", "AI Search Agent",
     [("$148M", "ARR\n(Jun 2025)", ACCENT_PURPLE),
      ("45M", "Monthly Active\nUsers", ACCENT_TEAL),
      ("780M", "Monthly\nQueries", ACCENT_RED)],
     ["700% ARR growth in 2024; targeting $656M revenue by end of 2026",
      "Agentic search: multi-step research, source verification, and synthesis",
      "$1.5B+ total funding; investors include NVIDIA, Jeff Bezos, SoftBank",
      "60-70M daily queries globally; 50M+ mobile app downloads"],
     ACCENT_PURPLE),
    ("Cognition (Devin)", "$10.2B valuation", "Autonomous Coding Agent",
     [("$150M+", "Combined ARR\n(post-Windsurf)", ACCENT_BLUE),
      ("67%", "PR Merge\nRate", ACCENT_GREEN),
      ("$500", "Per Seat\n/ Month", ACCENT_ORANGE)],
     ["ARR grew from $1M (Sep 2024) to $73M (Jun 2025); 73x in 9 months",
      "Merged hundreds of thousands of PRs; Goldman Sachs, Citi, Dell, Palantir",
      "$400M raised (Sep 2025) at $10.2B; total net burn under $20M across history",
      "67% PR merge rate (+33pp YoY); 4x faster, 2x more efficient than prior year"],
     ACCENT_BLUE),
]

y_offset = Inches(1.7)
for comp_name, val_label, product, stat_cards, bullets, color in agent_native:
    ch = Inches(1.7)
    card = add_card(slide, Inches(0.8), y_offset, Inches(11.5), ch)
    add_accent_bar(slide, Inches(0.8), y_offset, Pt(5), ch, color)
    add_text_box(slide, Inches(1.2), y_offset + Inches(0.06), Inches(2.5), Inches(0.28),
                 comp_name, font_size=17, color=color, bold=True)
    add_text_box(slide, Inches(3.8), y_offset + Inches(0.06), Inches(2.5), Inches(0.28),
                 product, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, Inches(9.5), y_offset + Inches(0.06), Inches(2.5), Inches(0.28),
                 val_label, font_size=11, color=MUTED, alignment=PP_ALIGN.RIGHT)

    for j, (sv, sl_label, sc) in enumerate(stat_cards):
        sx = Inches(1.2) + j * Inches(1.75)
        add_stat_card(slide, sx, y_offset + Inches(0.36), Inches(1.6), Inches(0.75),
                      sv, sl_label, sc)

    bullet_lines = [{"text": f"•  {b}", "size": 9, "color": LIGHT_GRAY,
                     "space_after": 1} for b in bullets]
    add_multi_text(slide, Inches(6.6), y_offset + Inches(0.33), Inches(5.5),
                   Inches(1.3), bullet_lines)
    y_offset += Inches(1.82)


# ── SLIDE 18: PRIVATE AI COMPANIES — ENTERPRISE & INFRASTRUCTURE ─────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Private AI Companies: Enterprise & Infrastructure",
             bar_color=ACCENT_GREEN, bar_width=Inches(3))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Enterprise-focused and infrastructure private AI companies "
             "building the agentic stack",
             font_size=13, color=LIGHT_GRAY)

enterprise_private = [
    ("xAI (Grok)", "$250B valuation", "Grok Models & Enterprise API",
     [("$3.8B", "Annualized Rev.\n(incl. X, 2025)", ACCENT_ORANGE),
      ("600M", "Monthly Active\nUsers (X+Grok)", ACCENT_RED),
      ("$20B", "Series E\n(Jan 2026)", ACCENT_BLUE)],
     ["Standalone xAI ~$500M ARR run rate; 38x YoY growth (incl. X merger)",
      "500+ enterprise API clients; 30M Grok monthly active users (app)",
      "SpaceX acquired xAI at $250B (Feb 2026); burning ~$1B/month, target profit by 2027",
      "Investors: Valor, Fidelity, QIA, NVIDIA, Cisco; $175B+ total funding"],
     ACCENT_RED),
    ("Mistral AI", "$14B valuation", "Le Chat, Mistral Large, Enterprise API",
     [("€60M", "2025 Revenue\n(5x from '23)", ACCENT_GREEN),
      ("€1B", "2026 Revenue\nTarget", ACCENT_TEAL),
      ("20K+", "Companies\nUsing API", ACCENT_PURPLE)],
     ["Europe's largest AI company; €2B funding round at $14B valuation (Sep 2025)",
      "Clients: BNP Paribas, Airbus, Schneider Electric; targeting €1B revenue in 2026",
      "Le Chat hit 1M users within a month; Vibe 2.0 for enterprise agent workflows",
      "650 employees; building tiered professional and team agentic features"],
     ACCENT_GREEN),
    ("Cohere", "$7B valuation", "North Agent Platform, Command R",
     [("$240M", "ARR\n(2025)", ACCENT_BLUE),
      ("287%", "YoY ARR\nGrowth", ACCENT_ORANGE),
      ("70%", "Gross\nMargins", ACCENT_GREEN)],
     ["85% revenue from private deployments; enterprise clients: Oracle, RBC, LG, Notion",
      "'North' AI agent platform for knowledge workers — multi-step business processes",
      "$970M total funding; backed by NVIDIA, Salesforce Ventures, AMD Ventures",
      "Positioning for IPO; 50%+ quarter-over-quarter growth throughout 2025"],
     ACCENT_BLUE),
    ("Databricks (Private)", "$134B valuation", "Mosaic AI, Agent Bricks, Lakebase",
     [("$5.4B", "ARR\n(~65% YoY)", ACCENT_ORANGE),
      ("$1.4B", "AI-Specific\nARR", ACCENT_PURPLE),
      ("20K+", "Global\nCustomers", ACCENT_TEAL)],
     ["Largest private data AI company; 60%+ Fortune 500 on platform",
      "$7B equity + debt raise at $134B; approaching IPO readiness (FCF positive)",
      "80% of databases created by AI agents; 327% multi-agent workflow growth",
      "Open-source stack: MLflow, Unity Catalog; runs on AWS, GCP, Azure"],
     ACCENT_ORANGE),
]

y_offset = Inches(1.65)
for comp_name, val_label, product, stat_cards, bullets, color in enterprise_private:
    ch = Inches(1.32)
    card = add_card(slide, Inches(0.8), y_offset, Inches(11.5), ch)
    add_accent_bar(slide, Inches(0.8), y_offset, Pt(5), ch, color)
    add_text_box(slide, Inches(1.2), y_offset + Inches(0.04), Inches(2.3), Inches(0.25),
                 comp_name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(3.6), y_offset + Inches(0.04), Inches(3.5), Inches(0.25),
                 product, font_size=10, color=WHITE, bold=True)
    add_text_box(slide, Inches(9.5), y_offset + Inches(0.04), Inches(2.5), Inches(0.25),
                 val_label, font_size=10, color=MUTED, alignment=PP_ALIGN.RIGHT)

    for j, (sv, sl_label, sc) in enumerate(stat_cards):
        sx = Inches(1.2) + j * Inches(1.65)
        add_stat_card(slide, sx, y_offset + Inches(0.3), Inches(1.5), Inches(0.65),
                      sv, sl_label, sc)

    bullet_lines = [{"text": f"•  {b}", "size": 8, "color": LIGHT_GRAY,
                     "space_after": 0} for b in bullets]
    add_multi_text(slide, Inches(6.4), y_offset + Inches(0.28), Inches(5.7),
                   Inches(1.0), bullet_lines)
    y_offset += Inches(1.4)

add_multi_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.3), [
    {"text": "Sources: Press releases, funding announcements, investor disclosures, "
             "and media reports (2025–2026). Private company revenue figures are "
             "estimates from public reporting.",
     "size": 9, "color": MUTED},
])


# ── SLIDE 19: PRIVATE AI VALUATION LANDSCAPE ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Private AI Valuation Landscape",
             bar_color=ACCENT_PURPLE, bar_width=Inches(2.5))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Valuations and ARR of privately held agentic AI companies (2025–2026)",
             font_size=13, color=LIGHT_GRAY)

val_tiers = [
    ("$250B+", [
        ("OpenAI", "$500B", "$25B ARR", "GPT, Agents SDK, Operator", ACCENT_ORANGE),
        ("Anthropic", "$380B", "$14B ARR", "Claude, MCP, Computer Use", ACCENT_PURPLE),
        ("xAI", "$250B", "~$500M standalone", "Grok, Enterprise API", ACCENT_RED),
    ], ACCENT_ORANGE),
    ("$10B – $134B", [
        ("Databricks", "$134B", "$5.4B ARR", "Mosaic AI, Agent Bricks, Lakebase",
         ACCENT_BLUE),
        ("Cursor", "$29.3B", "$2B+ ARR", "AI coding agent, background agents",
         ACCENT_GREEN),
        ("Perplexity", "$20B", "$148M ARR", "AI search agent, research",
         ACCENT_TEAL),
        ("Mistral AI", "$14B", "€60M rev", "Le Chat, enterprise models",
         ACCENT_ORANGE),
        ("Cognition", "$10.2B", "$150M+ ARR", "Devin autonomous coding agent",
         ACCENT_PURPLE),
    ], ACCENT_BLUE),
    ("$5B – $10B", [
        ("Cohere", "$7B", "$240M ARR", "North agent platform, Command R",
         ACCENT_GREEN),
    ], ACCENT_GREEN),
]

y_offset = Inches(1.75)
for tier_label, companies, tier_color in val_tiers:
    add_accent_bar(slide, Inches(0.8), y_offset, Inches(11.5), Pt(3), tier_color)
    add_text_box(slide, Inches(0.8), y_offset + Inches(0.05), Inches(1.8), Inches(0.3),
                 tier_label, font_size=14, color=tier_color, bold=True)
    for j, (name, val, rev, prod, color) in enumerate(companies):
        x_start = Inches(0.8) + (j % 3) * Inches(3.9)
        row_y = y_offset + Inches(0.35) + (j // 3) * Inches(0.7)
        card = add_card(slide, x_start, row_y, Inches(3.7), Inches(0.6))
        add_accent_bar(slide, x_start, row_y, Pt(4), Inches(0.6), color)
        add_text_box(slide, x_start + Inches(0.15), row_y + Inches(0.02),
                     Inches(1.5), Inches(0.25), name, font_size=13, color=color,
                     bold=True)
        add_text_box(slide, x_start + Inches(1.7), row_y + Inches(0.02),
                     Inches(1.8), Inches(0.25), f"{val}  |  {rev}", font_size=10,
                     color=WHITE)
        add_text_box(slide, x_start + Inches(0.15), row_y + Inches(0.28),
                     Inches(3.4), Inches(0.25), prod, font_size=9, color=LIGHT_GRAY)
    n_rows = (len(companies) + 2) // 3
    y_offset += Inches(0.4) + n_rows * Inches(0.7) + Inches(0.1)

add_multi_text(slide, Inches(0.8), y_offset + Inches(0.1), Inches(11), Inches(1.0), [
    {"text": "Key Trend: Combined private AI company valuations exceed $1.3 trillion. "
             "Agentic AI capabilities are the primary growth driver across all tiers.",
     "size": 13, "color": ACCENT_ORANGE, "bold": True, "space_after": 6},
    {"text": "Databricks ($134B) sits at the center of this landscape — providing "
             "the enterprise data platform that both foundation model providers and "
             "agent-native companies rely on for production deployments.",
     "size": 12, "color": LIGHT_GRAY, "space_after": 4},
    {"text": "Sources: Press releases, funding announcements, and media reports. "
             "Private valuations and revenue figures are estimates from public reporting.",
     "size": 10, "color": MUTED},
])


# ── SLIDE 20: CEO QUOTES ON AGENTIC AI ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "CEO Quotes on Agentic AI — From Filings & Earnings",
             bar_color=ACCENT_ORANGE, bar_width=Inches(3))

ceo_quotes = [
    ("Ali Ghodsi", "Databricks CEO", "Fortune Brainstorm AI, Dec 2025",
     '"It\'s very elusive to get AI that really works and understands that proprietary '
     'data that\'s inside enterprise... this could maybe get all the way to a trillion."',
     ACCENT_ORANGE),
    ("Marc Benioff", "Salesforce CEO", "Q4 FY2026 Earnings, Feb 2026",
     '"Agentforce just became an $800 million business. '
     'Agentic AI is a tailwind for our business... this is the third wave of AI."',
     ACCENT_BLUE),
    ("Satya Nadella", "Microsoft CEO", "FY2025 Annual Report / Ignite 2025",
     '"Every customer I talk to is looking to reshape their business processes '
     'with AI agents. We are building the Frontier Firm — human-led, agent-operated."',
     ACCENT_PURPLE),
    ("Jensen Huang", "NVIDIA CEO", "GTC 2025 Keynote",
     '"The amount of compute we need for AI inference is easily 100x more than we '
     'thought we needed this time last year. Agents reason, plan, and act."',
     ACCENT_GREEN),
    ("Alex Karp", "Palantir CEO", "Q4 2025 Earnings, Feb 2026",
     '"AI has just put gasoline on all the tribal knowledge we have in our products. '
     'The demand for AI agents that can do things in the real world is unlike '
     'anything we\'ve seen."',
     ACCENT_RED),
    ("Sridhar Ramaswamy", "Snowflake CEO", "Q4 FY2026 Earnings, Feb 2026",
     '"Winners will be platforms that combine trusted data with secure execution. '
     'We are the control plane for the agentic era."',
     ACCENT_TEAL),
]

for i, (name, role, source, quote, color) in enumerate(ceo_quotes):
    y = Inches(1.5) + i * Inches(0.95)
    card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(0.85))
    add_accent_bar(slide, Inches(0.8), y, Pt(5), Inches(0.85), color)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(2.2), Inches(0.3),
                 f"{name}, {role}", font_size=12, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.3), Inches(2.2), Inches(0.5),
                 source, font_size=9, color=MUTED)
    add_text_box(slide, Inches(3.7), y + Inches(0.08), Inches(8.3), Inches(0.7),
                 quote, font_size=12, color=LIGHT_GRAY)

add_multi_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.3), [
    {"text": "Sources: Public earnings call transcripts, SEC filings, and conference "
             "keynotes (2025–2026)", "size": 10, "color": MUTED},
])


# ── SLIDE 21: MARKET CONTEXT ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Enterprise AI Agent Market Context",
             bar_color=ACCENT_ORANGE)

stats_row1 = [
    ("$58B → $474B", "Enterprise AI Market\n2025 → 2030 (52% CAGR)", ACCENT_ORANGE),
    ("$644B", "GenAI Spend\nForecast 2025", ACCENT_BLUE),
    ("$1.5T", "Total AI Spend\nForecast 2025", ACCENT_GREEN),
    ("75%", "C-Suite Executives Rank\nAI as Top-3 Priority", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(stats_row1):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(1.5), Inches(2.7),
                  Inches(1.6), val, label, color)

stats_row2 = [
    ("11x", "More AI Models\nDeployed YoY", ACCENT_TEAL),
    ("3x", "More Efficient at\nProduction Deployment", ACCENT_RED),
    ("377%", "Vector Database\nUsage Growth YoY", ACCENT_BLUE),
    ("76%", "Companies Using LLMs\nChoose Open Source", ACCENT_ORANGE),
]
for i, (val, label, color) in enumerate(stats_row2):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(3.5), Inches(2.7),
                  Inches(1.6), val, label, color)

card = add_card(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5))
add_accent_bar(slide, Inches(0.8), Inches(5.5), Pt(5), Inches(1.5), ACCENT_ORANGE)
add_multi_text(slide, Inches(1.3), Inches(5.6), Inches(10.5), Inches(1.3), [
    {"text": "Key Insight: Financial Services leads GPU usage growth at 88% over "
             "6 months,", "size": 14, "color": ACCENT_ORANGE, "bold": True,
     "space_after": 4},
    {"text": "followed by Healthcare & Life Sciences — highly regulated industries "
             "are the surprise early adopters of GenAI.", "size": 14,
     "color": LIGHT_GRAY, "space_after": 8},
    {"text": "Source: Databricks State of Data + AI Report 2024 — 10,000+ global "
             "customers including 300+ Fortune 500 companies", "size": 11,
     "color": MUTED},
])


# ── SLIDE 22: KEY TAKEAWAYS ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Key Takeaways", bar_color=ACCENT_ORANGE, bar_width=Inches(1.5))

takeaways = [
    ("AI Agents Are Reshaping Enterprise Data",
     "80% of databases on Databricks are now created by AI agents. Multi-agent "
     "workflows grew 327% in just 4 months. The shift from single chatbots to "
     "autonomous multi-agent systems is accelerating rapidly.",
     ACCENT_ORANGE),
    ("Token Optimization Is Critical for ROI",
     "Companies like adidas achieved 98.5% token efficiency (200K → 3K tokens) and "
     "91.67% cost savings. Without governance, Gartner predicts >40% of agentic AI "
     "projects will be canceled by 2027.",
     ACCENT_BLUE),
    ("Governance Multiplies Production Success",
     "Companies using AI governance tools deploy 12x more projects to production. "
     "Evaluation tools yield 6x more deployments. AI Gateway governance usage grew "
     "7x in 9 months.",
     ACCENT_GREEN),
    ("Enterprise Case Studies Show Massive Value",
     "Walmart saved $5.6M annually; AstraZeneca reviewed 400K documents in 60 minutes; "
     "Royal Bank of Canada reduced research from days to minutes; 7-Eleven deployed "
     "frontline AI agents.",
     ACCENT_PURPLE),
    ("The Technology Stack Is Maturing",
     "Mosaic AI for compound systems, MLflow 3.0 for observability, Agent Bricks for "
     "no-code building, MCP for tool-calling, and Lakebase for agent databases — the "
     "enterprise agent platform is production-ready.",
     ACCENT_TEAL),
]

for i, (title, desc, color) in enumerate(takeaways):
    y = Inches(1.5) + i * Inches(1.15)
    card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.0))
    add_accent_bar(slide, Inches(0.8), y, Pt(5), Inches(1.0), color)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.4), Inches(10.8), Inches(0.55),
                 desc, font_size=12, color=LIGHT_GRAY)


# ── SLIDE 23: SOURCES & REFERENCES ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Sources & References", bar_color=ACCENT_ORANGE,
             bar_width=Inches(1.8))

sources = [
    "Databricks Data + AI Summit 2024, 2025, 2026 — Keynotes & Agent Sessions",
    "Databricks State of AI Agents Report 2026 / State of Data + AI Report 2024",
    "Databricks Blog: AI Agent Trends, Mosaic AI, MCP Integration",
    "Ali Ghodsi (Databricks CEO): Fortune Brainstorm AI (Dec '25) & CNBC (Feb '26)",
    "Salesforce Q4 FY2026 Earnings — Agentforce $800M ARR (salesforce.com, Feb '26)",
    "Microsoft FY2025 10-K — Copilot & Agent 365 (sec.gov)",
    "ServiceNow Q4 2025 / 10-K — Now Assist $600M+ ACV (sec.gov)",
    "Amazon Q4 2025 — AWS Bedrock multi-B$ run rate (sec.gov)",
    "Palantir Q4 2025 / 10-K — AIP $4.5B revenue (sec.gov)",
    "Snowflake Q4 FY2026 — Cortex AI $100M run rate (sec.gov)",
    "NVIDIA FY2025 / GTC 2025 — $130.5B revenue, token economics (sec.gov)",
    "Google / Alphabet Q4 2025 / 10-K — Cloud $70B+ run rate (sec.gov)",
    "OpenAI — $25B ARR, $500B valuation (Reuters, Jan–Mar 2026)",
    "Anthropic — $14B ARR, $380B valuation, Series G (Reuters, Feb 2026)",
    "Cursor — $2B+ ARR, $29.3B valuation (Series D, Nov 2025)",
    "Perplexity — $148M ARR, $20B valuation (Sacra, Sep 2025)",
    "xAI / Grok — $250B valuation, $20B Series E (TechCrunch, Jan 2026)",
    "Mistral AI — $14B valuation, €1B 2026 target (Davos, Jan 2026)",
    "Cohere — $240M ARR, $7B valuation (Reuters, Aug 2025)",
    "Cognition / Devin — $150M+ ARR, $10.2B valuation (Sep 2025)",
    "adidas / Walmart / AstraZeneca / 7-Eleven / Zillow / RBC — Summit presentations",
]

source_lines = [{"text": f"{i+1}.  {s}", "size": 11, "color": LIGHT_GRAY,
                 "space_after": 4} for i, s in enumerate(sources)]
add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), source_lines)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "All data sourced from publicly available Databricks summit sessions, "
             "blog posts, and reports.",
             font_size=11, color=MUTED)


# ── SAVE ───────────────────────────────────────────────────────────────────────
output_path = "/workspace/agentic_ai_databricks_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
