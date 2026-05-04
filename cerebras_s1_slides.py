"""
Generate a polished PowerPoint deck summarising the best slides from
Cerebras Systems' S-1 filing (April 2026).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0xA8, 0x6B)
RED = RGBColor(0xE8, 0x3E, 0x3E)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_paragraph(tf, text, font_size=16, bold=False, color=DARK,
                   alignment=PP_ALIGN.LEFT, space_before=Pt(6),
                   font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def _add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6),
                    color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_kpi_card(slide, left, top, width, height, label, value,
                  sub_text="", value_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = "Calibri"
    p.font.bold = False

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = value_color
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    if sub_text:
        p3 = tf.add_paragraph()
        p3.text = sub_text
        p3.font.size = Pt(10)
        p3.font.color.rgb = MEDIUM_GRAY
        p3.font.name = "Calibri"
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(2)


def _add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Build the deck ──────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, DARK)
_add_accent_bar(sl, Inches(1), Inches(2.4), Inches(0.1), Inches(2.2), ACCENT)

_add_textbox(sl, Inches(1.4), Inches(2.3), Inches(10), Inches(1),
             "CEREBRAS SYSTEMS", 44, True, WHITE, font_name="Calibri")
_add_textbox(sl, Inches(1.4), Inches(3.1), Inches(10), Inches(0.6),
             "S-1 Filing Analysis  |  April 2026  |  Nasdaq: CBRS", 22, False,
             RGBColor(0xAA, 0xCC, 0xEE), font_name="Calibri")
_add_textbox(sl, Inches(1.4), Inches(3.8), Inches(10), Inches(0.5),
             "IPO Target: ~$2B raise  ·  Valuation: $22 – 25B  ·  28M Class A shares at $115 – $125/share",
             16, False, RGBColor(0x88, 0xAA, 0xCC), font_name="Calibri")
_add_textbox(sl, Inches(1.4), Inches(5.4), Inches(6), Inches(0.4),
             "Source: SEC EDGAR, S-1 Registration Statement filed 17-Apr-2026",
             11, False, MEDIUM_GRAY, font_name="Calibri")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Company Overview
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, LIGHT_GRAY)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Company Snapshot", 32, True, DARK)

items = [
    ("Founded", "2016 in Sunnyvale, CA"),
    ("CEO", "Andrew D. Feldman (co-founder)"),
    ("Mission", "Build the world's largest and fastest AI chips to accelerate AI compute"),
    ("Core Product", "Wafer-Scale Engine 3 (WSE-3) — the largest chip ever built"),
    ("Business Model", "Chip sales + cloud inference/training services via own data centers"),
    ("Filing Status", "Emerging growth company; non-accelerated filer"),
    ("Legal Counsel", "Latham & Watkins LLP / Davis Polk & Wardwell LLP"),
    ("Underwriters", "Citi (lead), Barclays, UBS, Wells Fargo, Mizuho, TD Cowen"),
]

y = Inches(1.5)
for label, value in items:
    _add_textbox(sl, Inches(1.0), y, Inches(2.3), Inches(0.4),
                 label, 14, True, ACCENT)
    _add_textbox(sl, Inches(3.4), y, Inches(9), Inches(0.4),
                 value, 14, False, DARK)
    y += Inches(0.58)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WSE-3 Technology
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, WHITE)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "WSE-3: The World's Largest AI Chip", 32, True, DARK)
_add_textbox(sl, Inches(1.0), Inches(1.15), Inches(10), Inches(0.4),
             "Manufactured by TSMC on 5 nm process  ·  Unveiled March 2024",
             14, False, MEDIUM_GRAY)

card_w = Inches(2.5)
card_h = Inches(1.6)
gap = Inches(0.35)
start_x = Inches(0.8)
start_y = Inches(2.0)

specs = [
    ("Transistors", "4 Trillion", "19× NVIDIA B200"),
    ("AI Cores", "900,000", "52× more than largest GPU"),
    ("Silicon Area", "46,225 mm²", "57× largest GPU"),
    ("Peak Performance", "125 PFLOPS", "28× NVIDIA B200"),
    ("On-Chip SRAM", "44 GB", "vs. off-chip HBM in GPUs"),
    ("Memory BW", "21 PB/s", "7,000× GPU mem BW"),
    ("Fabric BW", "214 PB/s", "On-wafer interconnect"),
    ("Inference Speed", "Up to 15×", "faster than GPU solutions"),
]

for i, (label, value, sub) in enumerate(specs):
    col = i % 4
    row = i // 4
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    _add_kpi_card(sl, x, y, card_w, card_h, label, value, sub)

_add_textbox(sl, Inches(0.8), Inches(5.7), Inches(11), Inches(0.6),
             "Cerebras's wafer-scale approach places the entire model on a single chip, eliminating "
             "inter-chip communication bottlenecks that limit GPU cluster performance.",
             13, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Financial Highlights (KPI dashboard)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, LIGHT_GRAY)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Financial Highlights — FY 2025", 32, True, DARK)

kw = Inches(3.5)
kh = Inches(1.8)
kg = Inches(0.5)
kx0 = Inches(0.8)
ky = Inches(1.6)

kpis = [
    ("Revenue", "$510M", "+76% YoY", ACCENT),
    ("Operating Loss", "($75.7M)", "vs. ($21.8M) in FY24", RED),
    ("Non-GAAP Net Income", "$237.8M", "Includes one-time items", GREEN),
]

for i, (lbl, val, sub, clr) in enumerate(kpis):
    _add_kpi_card(sl, kx0 + i * (kw + kg), ky, kw, kh, lbl, val, sub, clr)

ky2 = Inches(3.8)
kpis2 = [
    ("Order Backlog", "$24.6B", "Driven by OpenAI deal", ACCENT),
    ("OpenAI Loan", "$1.0B", "Advance for compute buildout", ORANGE),
    ("G42 Paper Gain", "$363M", "Restructured deal (2025)", GREEN),
]
for i, (lbl, val, sub, clr) in enumerate(kpis2):
    _add_kpi_card(sl, kx0 + i * (kw + kg), ky2, kw, kh, lbl, val, sub, clr)

_add_textbox(sl, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
             "Note: GAAP net income was aided by $363M non-cash gain from G42 deal restructuring. "
             "Underlying operating loss widened YoY.",
             11, False, MEDIUM_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Revenue Trajectory
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, WHITE)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Revenue Growth Trajectory", 32, True, DARK)

bar_base_y = Inches(5.0)
bar_w = Inches(2.2)
bar_gap = Inches(1.5)
max_h = Inches(3.2)

rev_data = [
    ("FY 2024", 289, "$289M"),
    ("FY 2025", 510, "$510M"),
    ("Backlog", 800, "$24.6B backlog"),
]

for i, (year_label, val, display) in enumerate(rev_data):
    x = Inches(2.0) + i * (bar_w + bar_gap)
    bar_height = Emu(int(int(max_h) * val / 800))
    bar_top = Emu(int(bar_base_y) - int(bar_height))

    clr = ACCENT if i < 2 else GREEN
    _add_bar(sl, x, bar_top, bar_w, bar_height, clr)

    _add_textbox(sl, x, Emu(int(bar_top) - Inches(0.5)), bar_w, Inches(0.45),
                 display, 18, True, clr, PP_ALIGN.CENTER)
    _add_textbox(sl, x, Emu(int(bar_base_y) + Inches(0.15)), bar_w, Inches(0.4),
                 year_label, 14, True, DARK, PP_ALIGN.CENTER)

_add_textbox(sl, Inches(7.5), Inches(1.5), Inches(5), Inches(0.4),
             "+76% YoY Revenue Growth", 20, True, GREEN)
tf = _add_textbox(sl, Inches(7.5), Inches(2.1), Inches(5), Inches(2.5),
                  "", 14, False, DARK)
for line in [
    "Revenue nearly doubled from FY24 to FY25",
    "Massive $24.6B backlog provides forward visibility",
    "OpenAI deal: 750 MW AI compute through 2028",
    "Option for ~3 GW more by 2030",
    "Business model shift: chip vendor → compute provider",
]:
    _add_paragraph(tf, "•  " + line, 13, False, DARK, space_before=Pt(8))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Customer Concentration
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, WHITE)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Customer & Revenue Concentration (FY 2025)", 32, True, DARK)

segments = [
    ("MBZUAI (UAE)", 62, ACCENT),
    ("G42 (UAE)", 24, RGBColor(0x00, 0xB4, 0xD8)),
    ("US Customers", 14, RGBColor(0x90, 0xBE, 0x6D)),
]

bar_x = Inches(1.0)
bar_y_start = Inches(2.0)
total_bar_w = Inches(11)
bar_h = Inches(0.9)

cum_x = bar_x
for label, pct, clr in segments:
    seg_w = Emu(int(total_bar_w * pct / 100))
    _add_bar(sl, cum_x, bar_y_start, seg_w, bar_h, clr)
    if pct > 10:
        _add_textbox(sl, cum_x, bar_y_start, seg_w, bar_h,
                     f"{label}  —  {pct}%", 13, True, WHITE, PP_ALIGN.CENTER)
    cum_x = Emu(int(cum_x) + int(seg_w))

_add_textbox(sl, Inches(1.0), Inches(3.2), Inches(11), Inches(0.4),
             "86% of FY 2025 revenue from UAE-based entities", 16, True, RED,
             PP_ALIGN.LEFT)

detail_items = [
    ("MBZUAI", "62% of total revenue — Mohamed bin Zayed University of AI"),
    ("G42", "24% of total revenue — Abu Dhabi tech conglomerate"),
    ("US Revenue", "$187.6M — down 34% YoY (from $282.7M in FY24)"),
    ("OpenAI (new)", "$10B contract + $1B advance loan — future revenue driver"),
    ("Geographic Mix", "Customers on 4 continents, but UAE dominates near-term"),
]

y = Inches(4.0)
for lbl, desc in detail_items:
    _add_textbox(sl, Inches(1.0), y, Inches(2.5), Inches(0.4), lbl, 14, True, ACCENT)
    _add_textbox(sl, Inches(3.6), y, Inches(8.5), Inches(0.4), desc, 14, False, DARK)
    y += Inches(0.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — OpenAI Partnership Deep Dive
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, DARK)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55), GREEN)
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "OpenAI Partnership — Transformational Deal", 32, True, WHITE)

card_data = [
    ("Contract Value", "$10B+", "Multi-year compute supply agreement"),
    ("Compute Commitment", "750 MW", "AI compute capacity through 2028"),
    ("Expansion Option", "~3 GW", "Additional capacity by 2030"),
    ("Advance Loan", "$1.0B", "To fund data center buildout"),
    ("Equity Warrants", "33M shares", "Near-free shares issued to OpenAI"),
    ("Backlog Impact", "$24.6B", "Majority driven by this single deal"),
]

cw = Inches(3.5)
ch = Inches(1.5)
cg = Inches(0.4)
cx0 = Inches(0.8)
cy0 = Inches(1.6)

for i, (lbl, val, sub) in enumerate(card_data):
    col = i % 3
    row = i // 3
    x = cx0 + col * (cw + cg)
    y = cy0 + row * (ch + cg)
    _add_kpi_card(sl, x, y, cw, ch, lbl, val, sub, GREEN)

_add_textbox(sl, Inches(0.8), Inches(5.3), Inches(11), Inches(1.5),
             "This deal fundamentally repositions Cerebras from a chip vendor to an AI compute "
             "infrastructure provider. OpenAI's commitment validates the WSE-3 architecture for "
             "large-scale inference and training workloads, and provides Cerebras with the capital "
             "and demand certainty to scale its data center operations.",
             13, False, RGBColor(0xBB, 0xCC, 0xDD), PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Risks
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, WHITE)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55), RED)
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Key Risk Factors", 32, True, DARK)

risks = [
    ("Customer Concentration",
     "86% of FY25 revenue from two UAE customers (MBZUAI & G42). Loss of either "
     "would materially impact financials. OpenAI contract mitigates long-term but "
     "adds single-customer execution risk."),
    ("Geopolitical / Regulatory",
     "CFIUS review of G42 stake delayed the original 2024 IPO. Export controls and "
     "national security scrutiny on UAE/China tech transfers remain an ongoing risk."),
    ("Profitability Gap",
     "Operating losses widened to $75.7M in FY25. GAAP profitability relied on a "
     "$363M non-cash gain from G42 restructuring — not sustainable operations."),
    ("TSMC Sole-Source Dependency",
     "WSE-3 fabricated exclusively by TSMC on 5 nm. Any supply disruption, "
     "geopolitical event, or capacity allocation change directly impacts production."),
    ("Competitive Pressure",
     "NVIDIA dominates AI accelerator market with entrenched ecosystem (CUDA). "
     "Cerebras must prove WSE-3 superiority at scale against B200/Blackwell and "
     "custom silicon from hyperscalers (Google TPU, Amazon Trainium, etc.)."),
    ("Capital Intensity",
     "Transition to data-center operator requires massive capex. $1B OpenAI loan "
     "creates financial obligations. Executing on $24.6B backlog demands flawless "
     "infrastructure buildout."),
]

y = Inches(1.4)
for title, desc in risks:
    _add_accent_bar(sl, Inches(1.0), y, Inches(0.06), Inches(0.8), RED)
    _add_textbox(sl, Inches(1.2), y, Inches(3.0), Inches(0.4),
                 title, 14, True, RED)
    _add_textbox(sl, Inches(1.2), Emu(int(y) + Inches(0.35)), Inches(11), Inches(0.55),
                 desc, 12, False, DARK)
    y = Emu(int(y) + Inches(0.95))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Bull vs Bear Case
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, LIGHT_GRAY)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Investment Thesis: Bull vs. Bear", 32, True, DARK)

col_w = Inches(5.5)
col_h = Inches(5.0)
col_y = Inches(1.5)

bull_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.8), col_y, col_w, col_h)
bull_shape.fill.solid()
bull_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
bull_shape.line.color.rgb = GREEN
bull_shape.line.width = Pt(2)

_add_textbox(sl, Inches(1.2), Inches(1.7), Inches(5), Inches(0.4),
             "BULL CASE", 20, True, GREEN, PP_ALIGN.CENTER)

bull_points = [
    "$24.6B backlog = exceptional revenue visibility",
    "OpenAI validation of WSE-3 architecture",
    "15× inference speed advantage vs. GPUs",
    "Chip-to-cloud model captures full value chain",
    "76% revenue growth with potential to accelerate",
    "AI compute TAM expanding rapidly ($100B+ market)",
]

tf = _add_textbox(sl, Inches(1.2), Inches(2.3), Inches(5), Inches(3.5), "", 13, False, GREEN)
for bp in bull_points:
    _add_paragraph(tf, "✓  " + bp, 13, False, RGBColor(0x2E, 0x7D, 0x32),
                   space_before=Pt(10))

bear_shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(7.0), col_y, col_w, col_h)
bear_shape.fill.solid()
bear_shape.fill.fore_color.rgb = RGBColor(0xFD, 0xE8, 0xE8)
bear_shape.line.color.rgb = RED
bear_shape.line.width = Pt(2)

_add_textbox(sl, Inches(7.4), Inches(1.7), Inches(5), Inches(0.4),
             "BEAR CASE", 20, True, RED, PP_ALIGN.CENTER)

bear_points = [
    "Extreme customer concentration (86% UAE)",
    "Operating losses widening, not narrowing",
    "GAAP profit entirely from one-time paper gain",
    "NVIDIA ecosystem moat (CUDA) remains deep",
    "TSMC single-source fab risk",
    "CFIUS / export control regulatory overhang",
]

tf = _add_textbox(sl, Inches(7.4), Inches(2.3), Inches(5), Inches(3.5), "", 13, False, RED)
for bp in bear_points:
    _add_paragraph(tf, "✗  " + bp, 13, False, RGBColor(0xC6, 0x28, 0x28),
                   space_before=Pt(10))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — IPO Terms Summary
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, WHITE)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55))
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "IPO Structure & Terms", 32, True, DARK)

ipo_cards = [
    ("Exchange", "Nasdaq", "Global Market"),
    ("Ticker", "CBRS", "Class A Common Stock"),
    ("Price Range", "$115 – $125", "Per share"),
    ("Shares Offered", "28M", "+4.2M overallotment"),
    ("Target Raise", "~$2B", "Gross proceeds"),
    ("Implied Valuation", "$22 – $25B", "Fully diluted"),
]

cw2 = Inches(3.5)
ch2 = Inches(1.5)
cg2 = Inches(0.5)
cx2 = Inches(0.8)
cy2 = Inches(1.5)

for i, (lbl, val, sub) in enumerate(ipo_cards):
    col = i % 3
    row = i // 3
    x = cx2 + col * (cw2 + cg2)
    y = cy2 + row * (ch2 + cg2)
    _add_kpi_card(sl, x, y, cw2, ch2, lbl, val, sub)

_add_textbox(sl, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
             "Dual-class share structure: Class A (1 vote/share) offered to public; Class B (10 votes/share) retained by insiders",
             13, False, MEDIUM_GRAY)

uw_text = ("Lead Book-Runners: Citi, Barclays  |  Book-Runners: UBS Investment Bank, "
           "Wells Fargo Securities, Mizuho, TD Cowen")
_add_textbox(sl, Inches(0.8), Inches(5.3), Inches(11), Inches(0.4),
             uw_text, 13, False, MEDIUM_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Takeaways / Summary
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, DARK)

_add_accent_bar(sl, Inches(0.8), Inches(0.5), Inches(0.07), Inches(0.55), ACCENT)
_add_textbox(sl, Inches(1.0), Inches(0.5), Inches(10), Inches(0.6),
             "Key Takeaways", 32, True, WHITE)

takeaways = [
    "Cerebras is the first credible challenger to NVIDIA's AI accelerator dominance, "
    "with a fundamentally differentiated wafer-scale architecture (WSE-3).",

    "Revenue is scaling rapidly ($510M, +76% YoY), but the company is not yet "
    "operationally profitable — operating losses widened to $75.7M in FY25.",

    "The $10B+ OpenAI partnership is transformational: it validates the technology, "
    "provides demand certainty ($24.6B backlog), and funds the transition to a "
    "compute-as-a-service model.",

    "Customer concentration is the single biggest near-term risk — 86% of revenue "
    "from UAE entities; US revenue actually declined 34% YoY.",

    "At $22–25B valuation (~45–50× FY25 revenue), the market is pricing in "
    "successful execution of the OpenAI deal and broader customer diversification.",

    "Investors should monitor: (1) OpenAI contract execution milestones, "
    "(2) US customer growth, (3) path to operating profitability, "
    "(4) CFIUS / export control developments.",
]

y = Inches(1.4)
for i, t in enumerate(takeaways):
    _add_accent_bar(sl, Inches(1.0), y, Inches(0.06), Inches(0.7), ACCENT)
    num_tf = _add_textbox(sl, Inches(1.2), y, Inches(0.4), Inches(0.4),
                          str(i + 1), 16, True, ACCENT)
    _add_textbox(sl, Inches(1.6), y, Inches(10.5), Inches(0.8),
                 t, 13, False, RGBColor(0xDD, 0xDD, 0xDD))
    y = Emu(int(y) + Inches(0.88))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Disclaimer
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
_set_slide_bg(sl, LIGHT_GRAY)

_add_textbox(sl, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "Disclaimer", 28, True, DARK, PP_ALIGN.CENTER)
_add_textbox(sl, Inches(1.5), Inches(3.3), Inches(10), Inches(2.5),
             "This presentation is for informational purposes only and does not constitute "
             "investment advice, a recommendation, or an offer to buy or sell securities. "
             "All data is sourced from Cerebras Systems' S-1 registration statement filed "
             "with the SEC on April 17, 2026, and publicly available analyses. "
             "Financial figures are subject to change and should be verified against the "
             "official filing. Past performance is not indicative of future results.",
             14, False, MEDIUM_GRAY, PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
output_path = "/workspace/Cerebras_S1_Filing_Analysis.pptx"
prs.save(output_path)
print(f"Saved → {output_path}")
