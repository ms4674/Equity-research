#!/usr/bin/env python3
"""Builds `AI_Equity_Research_Contributions.pptx`.

A 16-slide, 16:9 capabilities deck showing the breadth of an equity research
department's contribution across AI, LLMs and the datacenter buildout.
All quantitative figures in the deck are illustrative.

Usage:  python3 build_deck.py
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ----------------------------------------------------------------------------
# Design system
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)       # primary dark
NAVY_2 = RGBColor(0x13, 0x2E, 0x52)     # card-on-dark
BLUE = RGBColor(0x1F, 0x6F, 0xEB)       # accent
SKY = RGBColor(0x7F, 0xB3, 0xE8)        # secondary accent
GOLD = RGBColor(0xD9, 0xA4, 0x41)       # highlight
TEAL = RGBColor(0x2C, 0xA8, 0x9A)
SLATE = RGBColor(0x33, 0x47, 0x5B)      # body text
GREY = RGBColor(0x6B, 0x7A, 0x8D)       # muted text
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)      # light panel background
LINE = RGBColor(0xD7, 0xE0, 0xEB)       # hairline
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_W = Inches(12.233)

CHART_SERIES_COLORS = [BLUE, GOLD, TEAL, SKY, NAVY_2, GREY]

DEPT = "Global Technology Equity Research"
FOOTER = f"{DEPT}  |  AI, LLMs & the Datacenter Buildout"


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=None,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(0.75)
    return sp


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=Pt(0), line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of
    (text, size, color, bold, italic) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        for text, size, color, bold, italic in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = size
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
    return tb


def para(text, size, color, bold=False, italic=False):
    return [(text, size, color, bold, italic)]


def bullet_block(slide, x, y, w, h, items, size=Pt(12.5), color=SLATE,
                 head_color=NAVY, gap=Pt(7), line_spacing=1.06):
    """items: list of (level, text) or (level, text, bold). Level 0 gets a
    square accent marker, level 1 an en-dash."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in items:
        level, text = item[0], item[1]
        bold = item[2] if len(item) > 2 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = gap
        p.line_spacing = line_spacing
        if level == 0:
            marker = p.add_run()
            marker.text = "\u25aa "
            marker.font.name = FONT
            marker.font.size = size
            marker.font.color.rgb = BLUE
            marker.font.bold = True
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = size
            r.font.color.rgb = head_color if bold else color
            r.font.bold = bold
        else:
            marker = p.add_run()
            marker.text = "      \u2013  "
            marker.font.name = FONT
            marker.font.size = size - Pt(1)
            marker.font.color.rgb = GREY
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = size - Pt(1)
            r.font.color.rgb = color
    return tb


def slide_header(slide, kicker, title, page_no):
    """Standard light-slide header: kicker, title, accent rule, footer."""
    textbox(slide, MARGIN, Inches(0.42), CONTENT_W, Inches(0.30),
            [para(kicker.upper(), Pt(11), BLUE, bold=True)])
    textbox(slide, MARGIN, Inches(0.70), CONTENT_W, Inches(0.62),
            [para(title, Pt(25), NAVY, bold=True)])
    rect(slide, MARGIN, Inches(1.36), Inches(0.85), Pt(3.2), fill=GOLD)
    footer(slide, page_no)


def footer(slide, page_no, dark=False):
    color = SKY if dark else GREY
    textbox(slide, MARGIN, Inches(7.10), Inches(8.5), Inches(0.3),
            [para(FOOTER, Pt(8.5), color)])
    textbox(slide, SLIDE_W - Inches(1.15), Inches(7.10), Inches(0.6), Inches(0.3),
            [para(str(page_no), Pt(8.5), color)], align=PP_ALIGN.RIGHT)
    textbox(slide, Inches(9.0), Inches(7.10), Inches(3.15), Inches(0.3),
            [para("All figures illustrative", Pt(8.5), color, italic=True)],
            align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, value, label, value_color=BLUE,
              bg=LIGHT, label_color=SLATE, value_size=Pt(30)):
    rect(slide, x, y, w, h, fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.055)
    rect(slide, x, y + Inches(0.14), Pt(3.2), h - Inches(0.28), fill=GOLD)
    textbox(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4),
            Inches(0.55), [para(value, value_size, value_color, bold=True)])
    textbox(slide, x + Inches(0.22), y + h - Inches(0.78), w - Inches(0.4),
            Inches(0.68), [para(label, Pt(10.5), label_color)],
            line_spacing=1.03)


def chart_title(slide, x, y, w, text, sub=None):
    textbox(slide, x, y, w, Inches(0.28),
            [para(text, Pt(12.5), NAVY, bold=True)])
    if sub:
        textbox(slide, x, y + Inches(0.27), w, Inches(0.24),
                [para(sub, Pt(9.5), GREY, italic=True)])


def style_axis(axis, size=Pt(9), keep_gridlines=False):
    axis.tick_labels.font.size = size
    axis.tick_labels.font.name = FONT
    axis.tick_labels.font.color.rgb = GREY
    axis.format.line.color.rgb = LINE
    axis.major_tick_mark = XL_TICK_MARK.NONE
    axis.minor_tick_mark = XL_TICK_MARK.NONE
    if not keep_gridlines:
        axis.has_major_gridlines = False
    else:
        axis.has_major_gridlines = True
        axis.major_gridlines.format.line.color.rgb = LINE
        axis.major_gridlines.format.line.width = Pt(0.5)


def base_chart(slide, kind, x, y, w, h, categories, series, legend=True):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    frame = slide.shapes.add_chart(kind, x, y, w, h, data)
    chart = frame.chart
    chart.has_title = False
    chart.font.size = Pt(9)
    chart.font.name = FONT
    chart.font.color.rgb = SLATE
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    else:
        chart.has_legend = False
    return chart


def color_series(chart, colors=CHART_SERIES_COLORS, lines=False):
    for i, s in enumerate(chart.plots[0].series):
        c = colors[i % len(colors)]
        if lines:
            s.format.line.color.rgb = c
            s.format.line.width = Pt(2.5)
            s.smooth = False
        else:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = c
            s.format.line.fill.background()


# ----------------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------------

def s01_title(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    rect(s, 0, Inches(7.34), SLIDE_W, Inches(0.16), fill=GOLD)
    # Faint oversized backdrop glyph
    textbox(s, Inches(8.1), Inches(0.4), Inches(5.2), Inches(5.0),
            [para("AI", Pt(300), NAVY_2, bold=True)], wrap=False)

    textbox(s, MARGIN, Inches(1.15), Inches(10.5), Inches(0.35),
            [para(DEPT.upper(), Pt(13), GOLD, bold=True)])
    textbox(s, MARGIN, Inches(1.75), Inches(11.6), Inches(2.2),
            [para("From Silicon to Software:", Pt(43), WHITE, bold=True),
             para("The Breadth of Our AI Research Contribution", Pt(43), WHITE,
                  bold=True)],
            line_spacing=1.02)
    textbox(s, MARGIN, Inches(3.75), Inches(10.8), Inches(0.75),
            [para("Coverage, data and analysis across artificial intelligence, "
                  "large language models and the global datacenter buildout",
                  Pt(16), SKY)], line_spacing=1.1)
    textbox(s, MARGIN, Inches(4.55), Inches(6.0), Inches(0.3),
            [para("Department Capabilities Review  |  August 2026", Pt(12),
                  GREY)])

    stats = [("160+", "Publications across the\nAI stack in the last year"),
             ("45+", "Companies covered from\nsilicon to applications"),
             ("25+", "Proprietary data trackers\nmaintained and shared"),
             ("30+", "Financial & scenario\nmodels built for clients")]
    w = Inches(2.85)
    gap = Inches(0.27)
    x = MARGIN
    for value, label in stats:
        rect(s, x, Inches(5.35), w, Inches(1.45), fill=NAVY_2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
        textbox(s, x + Inches(0.22), Inches(5.52), w - Inches(0.4),
                Inches(0.5), [para(value, Pt(26), GOLD, bold=True)])
        textbox(s, x + Inches(0.22), Inches(6.10), w - Inches(0.4),
                Inches(0.65), [para(label.replace("\n", " "), Pt(10), SKY)],
                line_spacing=1.05)
        x += w + gap


def s02_exec_summary(prs):
    s = blank_slide(prs)
    slide_header(s, "Executive summary", "The AI research franchise at a glance", 2)

    cards = [("6", "Research pillars spanning the full AI value chain"),
             ("160+", "Notes, deep dives, primers and daily memos published"),
             ("25+", "Proprietary trackers: tokens, capex, capacity, pricing"),
             ("30+", "Company and scenario models, from S-1s to 3-statements"),
             ("45+", "Companies covered across semis, infra, cloud & software"),
             ("1,200+", "Client interactions on AI themes in the last year")]
    cw, ch = Inches(2.30), Inches(1.55)
    gx, gy = Inches(0.22), Inches(0.24)
    x0, y0 = MARGIN, Inches(1.75)
    for i, (v, l) in enumerate(cards):
        x = x0 + (i % 3) * (cw + gx)
        y = y0 + (i // 3) * (ch + gy)
        stat_card(s, x, y, cw, ch, v, l, value_size=Pt(26))

    rx = Inches(8.30)
    rw = SLIDE_W - rx - MARGIN
    rect(s, rx, Inches(1.75), rw, Inches(3.34), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
    textbox(s, rx + Inches(0.3), Inches(1.98), rw - Inches(0.6), Inches(0.3),
            [para("HOW WE CONTRIBUTE", Pt(11), GOLD, bold=True)])
    lines = [
        "One connected view of the AI trade — from wafer starts and megawatts "
        "to tokens and software seats.",
        "Proprietary data first: we measure the buildout rather than repeat "
        "consensus narratives.",
        "Cross-sector by design: semis, hardware, power, real estate, cloud, "
        "internet and software analysts publish together.",
    ]
    yy = Inches(2.38)
    for t in lines:
        textbox(s, rx + Inches(0.3), yy, rw - Inches(0.6), Inches(0.85),
                [[("\u25aa ", Pt(11.5), GOLD, True, False),
                  (t, Pt(11.5), WHITE, False, False)]], line_spacing=1.1)
        yy += Inches(0.92)

    rect(s, MARGIN, Inches(5.40), CONTENT_W, Inches(1.35), fill=LIGHT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    textbox(s, MARGIN + Inches(0.3), Inches(5.58), CONTENT_W - Inches(0.6),
            Inches(0.3), [para("THE MANDATE", Pt(10.5), BLUE, bold=True)])
    textbox(s, MARGIN + Inches(0.3), Inches(5.90), CONTENT_W - Inches(0.6),
            Inches(0.75),
            [para("Help investors size, price and risk-manage the largest "
                  "capital cycle in technology history: an AI infrastructure "
                  "buildout measured in hundreds of billions of dollars of "
                  "annual capex, gigawatts of power and quadrillions of tokens.",
                  Pt(13), SLATE)], line_spacing=1.12)


def s03_value_chain(prs):
    s = blank_slide(prs)
    slide_header(s, "Coverage universe",
                 "One research platform across the full AI value chain", 3)

    layers = [
        ("Semiconductors & hardware",
         "GPUs & custom accelerators \u00b7 HBM & memory \u00b7 AI networking "
         "\u00b7 optics & silicon photonics \u00b7 foundries & fab economics",
         NAVY),
        ("Datacenter physical infrastructure",
         "Hyperscale & colocation buildouts \u00b7 project-level GW pipeline "
         "\u00b7 cooling, racks & electrical gear \u00b7 land and shell supply",
         NAVY_2),
        ("Power & energy",
         "Power procurement & PPAs \u00b7 grid interconnection bottlenecks "
         "\u00b7 gas, nuclear & renewables \u00b7 energy cost pass-through",
         RGBColor(0x1B, 0x4A, 0x7A)),
        ("Cloud & hyperscalers",
         "Capex & depreciation \u00b7 owned vs leased capacity \u00b7 neoclouds "
         "& GPU rental economics \u00b7 AI financing vehicles",
         BLUE),
        ("Foundation models & LLMs",
         "Token consumption & pricing \u00b7 model economics & scaling "
         "\u00b7 open vs closed ecosystems \u00b7 enterprise LLM market share",
         RGBColor(0x41, 0x8B, 0xD8)),
        ("Software & applications",
         "Agentic AI & SaaS business models \u00b7 coding agents \u00b7 "
         "cybersecurity in the AI era \u00b7 internet platforms & ad tech",
         SKY),
    ]
    y = Inches(1.62)
    row_h = Inches(0.80)
    gap = Inches(0.075)
    label_w = Inches(3.55)
    for name, desc, color in layers:
        rect(s, MARGIN, y, label_w, row_h, fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        textbox(s, MARGIN + Inches(0.22), y, label_w - Inches(0.4), row_h,
                [para(name, Pt(13), WHITE, bold=True)],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        bx = MARGIN + label_w + Inches(0.15)
        bw = CONTENT_W - label_w - Inches(0.15)
        rect(s, bx, y, bw, row_h, fill=LIGHT, line_color=LINE, line_w=Pt(0.5))
        textbox(s, bx + Inches(0.22), y, bw - Inches(0.4), row_h,
                [para(desc, Pt(11), SLATE)], anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.03)
        y += row_h + gap

    textbox(s, MARGIN, y + Inches(0.06), CONTENT_W, Inches(0.35),
            [para("Every layer is covered by dedicated analysts publishing "
                  "into one coordinated AI research agenda \u2014 the stack is "
                  "analyzed as a single system, not six silos.",
                  Pt(11.5), GREY, italic=True)])


def s04_pillars(prs):
    s = blank_slide(prs)
    slide_header(s, "Research agenda",
                 "Six pillars organize our contribution to the AI debate", 4)

    pillars = [
        ("01", "LLM & token economics",
         ["Token consumption time series by model & workload",
          "LLM pricing index and $/M-token cost curves",
          "Returns to intelligence: unit economics of inference"]),
        ("02", "Datacenter buildout & capacity",
         ["Project-level global GW pipeline tracker",
          "Supply / demand capacity model through 2030",
          "Content-per-GW: memory, networking, optics, power gear"]),
        ("03", "Hyperscaler capex & AI financing",
         ["Capex, leases and depreciation deep dives",
          "Owned vs leased compute capacity analysis",
          "Neocloud commitments and AI compute deal comparisons"]),
        ("04", "Semis & hardware supply chain",
         ["GPU / TPU market share and performance frontier",
          "AI networking and optical supply chain models",
          "Fab economics: costs, subsidies, geographic risk"]),
        ("05", "Power, energy & sustainability",
         ["Power procurement strategies of hyperscalers",
          "Energy cost sensitivity of datacenter economics",
          "Grid bottlenecks, climate exposure, energy M&A"]),
        ("06", "Software & security in the AI era",
         ["Agentic AI impact on SaaS pricing & seats",
          "Cybersecurity winners from AI-driven threats",
          "Internet platforms: ads, engagement, regulation"]),
    ]
    cw, ch = Inches(3.95), Inches(2.48)
    gx, gy = Inches(0.19), Inches(0.22)
    x0, y0 = MARGIN, Inches(1.62)
    for i, (num, title, items) in enumerate(pillars):
        x = x0 + (i % 3) * (cw + gx)
        y = y0 + (i // 3) * (ch + gy)
        rect(s, x, y, cw, ch, fill=LIGHT, line_color=LINE, line_w=Pt(0.5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
        rect(s, x, y, cw, Inches(0.10), fill=GOLD)
        textbox(s, x + Inches(0.22), y + Inches(0.20), Inches(0.85),
                Inches(0.4), [para(num, Pt(17), SKY, bold=True)])
        textbox(s, x + Inches(0.22), y + Inches(0.50), cw - Inches(0.44),
                Inches(0.55), [para(title, Pt(14.5), NAVY, bold=True)],
                line_spacing=1.0)
        bullet_block(s, x + Inches(0.22), y + Inches(1.06), cw - Inches(0.44),
                     ch - Inches(1.2),
                     [(0, t) for t in items], size=Pt(10.5), gap=Pt(4),
                     line_spacing=1.03)


def s05_tokens(prs):
    s = blank_slide(prs)
    slide_header(s, "Pillar 01 \u2014 LLM & token economics",
                 "We measure the token economy so clients can price it", 5)

    bullet_block(s, MARGIN, Inches(1.66), Inches(5.9), Inches(5.2), [
        (0, "Token consumption database", True),
        (1, "Time series of usage across frontier and open-source models, "
            "coding agents vs copilots, and API aggregators"),
        (0, "LLM pricing & cost curves", True),
        (1, "Blended $/M-token pricing index; inference cost decomposition "
            "across compute, memory and context length"),
        (0, "Model economics & architecture", True),
        (1, "Pre-training vs fine-tuning spend, KV-cache & prefill costs, "
            "diffusion vs autoregressive trade-offs, parametric memory scaling"),
        (0, "Market structure", True),
        (1, "Enterprise LLM market share, open-source model adoption, "
            "foundation-model revenue and ARR benchmarking"),
        (0, "Why it matters", True),
        (1, "Token volumes and pricing are the demand signal that "
            "underwrites the entire datacenter and semiconductor capex cycle"),
    ], size=Pt(12))

    cx = Inches(6.85)
    cw = SLIDE_W - cx - MARGIN
    chart_title(s, cx, Inches(1.62), cw,
                "Tokens get cheaper, usage explodes",
                "Indexed to 100 in 2022 \u2014 illustrative")

    textbox(s, cx, Inches(2.18), cw, Inches(0.24),
            [para("Industry token consumption (index)", Pt(10), BLUE,
                  bold=True)])
    chart = base_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED, cx, Inches(2.46), cw, Inches(1.95),
        ["2022", "2023", "2024", "2025", "2026E"],
        [("Consumption", (100, 340, 1250, 4300, 11000))], legend=False)
    color_series(chart, [BLUE])
    chart.plots[0].gap_width = 60
    style_axis(chart.category_axis)
    style_axis(chart.value_axis, keep_gridlines=True)
    chart.value_axis.tick_labels.number_format = "#,##0"

    textbox(s, cx, Inches(4.58), cw, Inches(0.24),
            [para("Blended price per M tokens (index)", Pt(10), GOLD,
                  bold=True)])
    chart = base_chart(
        s, XL_CHART_TYPE.LINE, cx, Inches(4.86), cw, Inches(1.85),
        ["2022", "2023", "2024", "2025", "2026E"],
        [("Price", (100, 62, 30, 14, 8))], legend=False)
    color_series(chart, [GOLD], lines=True)
    style_axis(chart.category_axis)
    style_axis(chart.value_axis, keep_gridlines=True)
    chart.value_axis.tick_labels.number_format = "#,##0"


def s06_datacenter(prs):
    s = blank_slide(prs)
    slide_header(s, "Pillar 02 \u2014 Datacenter buildout & capacity",
                 "A project-level map of the global buildout", 6)

    bullet_block(s, MARGIN, Inches(1.66), Inches(5.9), Inches(5.2), [
        (0, "Global buildout tracker", True),
        (1, "Named-project database of hyperscale and neocloud campuses: "
            "MW/GW, location, developer, power source, timeline"),
        (0, "Supply / demand capacity model", True),
        (1, "Reconciles announced capacity against accelerator shipments, "
            "power availability and demand from token growth"),
        (0, "Content per gigawatt", True),
        (1, "Bill-of-materials analysis: memory, networking, optics, cooling "
            "and electrical content per GW of new capacity"),
        (0, "Frontier topics", True),
        (1, "Orbital datacenters, Gulf-region mega-projects, China/HK "
            "capacity, climate and water constraints on siting"),
        (0, "Why it matters", True),
        (1, "Capacity in the ground is the hardest constraint on the AI "
            "roadmap \u2014 our tracker turns press releases into a forecast"),
    ], size=Pt(12))

    cx = Inches(6.85)
    cw = SLIDE_W - cx - MARGIN
    chart_title(s, cx, Inches(1.62), cw,
                "Tracked datacenter capacity additions (GW)",
                "By region, 2023\u20132027E \u2014 illustrative")
    chart = base_chart(
        s, XL_CHART_TYPE.COLUMN_STACKED, cx, Inches(2.20), cw, Inches(4.45),
        ["2023", "2024", "2025", "2026E", "2027E"],
        [("United States", (4.5, 7.8, 13.5, 21.0, 29.0)),
         ("Rest of world", (2.6, 4.4, 7.5, 12.0, 17.5))])
    color_series(chart, [BLUE, SKY])
    chart.plots[0].gap_width = 60
    style_axis(chart.category_axis)
    style_axis(chart.value_axis, keep_gridlines=True)
    chart.value_axis.tick_labels.number_format = "#,##0"


def s07_capex(prs):
    s = blank_slide(prs)
    slide_header(s, "Pillar 03 \u2014 Hyperscaler capex & AI financing",
                 "Following the money behind the buildout", 7)

    bullet_block(s, MARGIN, Inches(1.66), Inches(5.9), Inches(5.2), [
        (0, "Capex & depreciation analytics", True),
        (1, "Quarterly hyperscaler capex tracker with GPU vs non-GPU split, "
            "useful-life assumptions and margin sensitivity"),
        (0, "Owned vs leased capacity", True),
        (1, "Lease accounting deep dives quantifying off-balance-sheet "
            "compute commitments and future rent obligations"),
        (0, "AI financing structures", True),
        (1, "Neocloud take-or-pay contracts, GPU-backed lending, "
            "vendor financing and circular revenue analysis"),
        (0, "Deal comparisons", True),
        (1, "Side-by-side economics of landmark AI compute deals, memory "
            "supply agreements and sovereign AI partnerships"),
        (0, "Why it matters", True),
        (1, "The durability of AI capex is the single biggest earnings "
            "question in global tech \u2014 we track it line by line"),
    ], size=Pt(12))

    cx = Inches(6.85)
    cw = SLIDE_W - cx - MARGIN
    chart_title(s, cx, Inches(1.62), cw,
                "Big-four hyperscaler capex ($B)",
                "Calendar years \u2014 illustrative")
    chart = base_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED, cx, Inches(2.20), cw, Inches(4.45),
        ["2022", "2023", "2024", "2025", "2026E"],
        [("Reported / consensus capex", (150, 148, 230, 365, 480)),
         ("of which: AI infrastructure", (35, 55, 130, 250, 360))])
    color_series(chart, [NAVY_2, GOLD])
    chart.plots[0].gap_width = 80
    chart.plots[0].overlap = -10
    style_axis(chart.category_axis)
    style_axis(chart.value_axis, keep_gridlines=True)
    chart.value_axis.tick_labels.number_format = "$#,##0"


def s08_semis(prs):
    s = blank_slide(prs)
    slide_header(s, "Pillar 04 \u2014 Semis & hardware supply chain",
                 "From accelerator roadmaps to fab economics", 8)

    cols = [
        ("Compute", [
            "GPU vs TPU vs custom ASIC market share and the "
            "performance-per-dollar frontier",
            "Accelerator roadmap analysis and cloud capacity implications",
            "S-1 and IPO work on emerging AI silicon vendors"]),
        ("Connectivity & memory", [
            "AI networking market share: Ethernet vs InfiniBand vs optical "
            "circuit switching",
            "Silicon photonics and co-packaged optics supply chain model",
            "HBM roadmaps and memory supply agreements with hyperscalers "
            "and neoclouds"]),
        ("Manufacturing & materials", [
            "Fab cost comparisons across geographies; new fab location "
            "tracking and subsidy analysis",
            "EU Chips Act and export-control dependency mapping",
            "Niche input risk: helium, sulphur and specialty materials "
            "exposure of leading-edge fabs"]),
    ]
    cw = Inches(3.95)
    gx = Inches(0.19)
    y = Inches(1.66)
    ch = Inches(4.05)
    for i, (title, items) in enumerate(cols):
        x = MARGIN + i * (cw + gx)
        rect(s, x, y, cw, ch, fill=LIGHT, line_color=LINE, line_w=Pt(0.5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        rect(s, x, y, cw, Inches(0.52), fill=NAVY)
        textbox(s, x + Inches(0.22), y, cw - Inches(0.44), Inches(0.52),
                [para(title, Pt(13.5), WHITE, bold=True)],
                anchor=MSO_ANCHOR.MIDDLE)
        bullet_block(s, x + Inches(0.22), y + Inches(0.72), cw - Inches(0.44),
                     ch - Inches(0.9), [(0, t) for t in items],
                     size=Pt(11), gap=Pt(8), line_spacing=1.08)

    rect(s, MARGIN, Inches(5.95), CONTENT_W, Inches(0.85), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    textbox(s, MARGIN + Inches(0.3), Inches(5.95), CONTENT_W - Inches(0.6),
            Inches(0.85),
            [[("Why it matters:  ", Pt(12), GOLD, True, False),
              ("hardware supply is the rate limiter of the AI cycle; our "
               "supply chain models convert component lead times into "
               "revenue visibility for the whole coverage universe.",
               Pt(12), WHITE, False, False)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)


def s09_power(prs):
    s = blank_slide(prs)
    slide_header(s, "Pillar 05 \u2014 Power, energy & sustainability",
                 "Electrons are the new constraint on compute", 9)

    bullet_block(s, MARGIN, Inches(1.66), Inches(5.9), Inches(5.2), [
        (0, "Power procurement strategies", True),
        (1, "How hyperscalers secure gigawatts: PPAs, behind-the-meter gas, "
            "nuclear restarts, SMRs and grid-interconnection queues"),
        (0, "Datacenter energy economics", True),
        (1, "Electricity cost sensitivity of cloud gross margins and "
            "GPU-hour pricing across regions"),
        (0, "Bottleneck analysis", True),
        (1, "Transformers, turbines, switchgear and transmission: where "
            "power equipment lead times gate the buildout"),
        (0, "Scenario & event work", True),
        (1, "Oil price shocks and cloud margins, tech-energy alliances and "
            "M&A, climate and water risk to datacenter siting"),
        (0, "Why it matters", True),
        (1, "Power availability now determines datacenter timelines more "
            "than chip supply \u2014 energy analysis is core tech research"),
    ], size=Pt(12))

    cx = Inches(6.85)
    cw = SLIDE_W - cx - MARGIN
    chart_title(s, cx, Inches(1.62), cw,
                "US datacenter power demand (GW)",
                "2023\u20132027E \u2014 illustrative")
    chart = base_chart(
        s, XL_CHART_TYPE.COLUMN_CLUSTERED, cx, Inches(2.20), cw, Inches(4.45),
        ["2023", "2024", "2025", "2026E", "2027E"],
        [("Datacenter power demand (GW)", (19, 25, 34, 46, 60))])
    color_series(chart, [TEAL])
    chart.plots[0].gap_width = 70
    style_axis(chart.category_axis)
    style_axis(chart.value_axis, keep_gridlines=True)
    chart.value_axis.tick_labels.number_format = "#,##0"
    chart.has_legend = False


def s10_software(prs):
    s = blank_slide(prs)
    slide_header(s, "Pillar 06 \u2014 Software & security in the AI era",
                 "Where AI value lands: seats, agents and security budgets", 10)

    cols = [
        ("Agentic AI & SaaS", [
            "Business-model impact of agents on per-seat pricing and the "
            "future of application software",
            "Coding agents vs copilots: token intensity, pricing and share",
            "Labor spend vs IT spend: sizing the automation opportunity"]),
        ("Cybersecurity", [
            "Vendor-by-vendor AI tailwind scorecards across endpoint, "
            "identity, data and network security",
            "Securing AI agents: new attack surfaces, new budgets",
            "Post-quantum and nation-state scenarios for security spend"]),
        ("Internet & platforms", [
            "AI-driven ad-tech gains and engagement economics at scale",
            "LLM licensing and data-rights value for content platforms",
            "Regulatory exposure: youth-safety rules and teen revenue at risk"]),
    ]
    cw = Inches(3.95)
    gx = Inches(0.19)
    y = Inches(1.66)
    ch = Inches(4.05)
    accents = [BLUE, GOLD, TEAL]
    for i, (title, items) in enumerate(cols):
        x = MARGIN + i * (cw + gx)
        rect(s, x, y, cw, ch, fill=LIGHT, line_color=LINE, line_w=Pt(0.5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        rect(s, x, y, cw, Inches(0.10), fill=accents[i])
        textbox(s, x + Inches(0.22), y + Inches(0.22), cw - Inches(0.44),
                Inches(0.4), [para(title, Pt(14), NAVY, bold=True)])
        bullet_block(s, x + Inches(0.22), y + Inches(0.72), cw - Inches(0.44),
                     ch - Inches(0.9), [(0, t) for t in items],
                     size=Pt(11), gap=Pt(8), line_spacing=1.08)

    rect(s, MARGIN, Inches(5.95), CONTENT_W, Inches(0.85), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    textbox(s, MARGIN + Inches(0.3), Inches(5.95), CONTENT_W - Inches(0.6),
            Inches(0.85),
            [[("Why it matters:  ", Pt(12), GOLD, True, False),
              ("software is where AI capex must ultimately be monetized; we "
               "connect infrastructure spend to application-layer revenue so "
               "clients can test the bull case end to end.",
               Pt(12), WHITE, False, False)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)


def s11_companies(prs):
    s = blank_slide(prs)
    slide_header(s, "Company-level depth",
                 "Single-name work products behind the themes", 11)

    cards = [
        ("Mega-cap platforms",
         "Capex & lease deep dives, owned-vs-leased compute, restructuring "
         "and opex analysis, ad-business AI sensitivity"),
        ("Cloud & neoclouds",
         "GPU rental economics, take-or-pay commitments, cloud COGS "
         "comparisons and AI-era margin bridges"),
        ("AI silicon & hardware",
         "S-1 breakdowns for AI chip IPOs, server OEM share shifts, "
         "networking & optics vendor models"),
        ("Memory & foundry",
         "HBM supply agreements, fab cost benchmarking, new-fab trackers "
         "and geographic risk analysis"),
        ("Software & security",
         "Earnings-quality work, AI-agent competitive assessments and "
         "growth durability frameworks for leading vendors"),
        ("Internet & mobility",
         "Full 3-statement models, investment memos, AV-alliance maps and "
         "sum-of-the-parts work including private AI exposure"),
        ("Private & pre-IPO",
         "Financial models for frontier labs and space/AI adjacencies "
         "built from disclosures, deals and supply chain signals"),
        ("Event-driven",
         "Layoff and restructuring impact notes, earnings previews and "
         "rapid-response analysis around AI announcements"),
    ]
    cw, ch = Inches(2.95), Inches(2.30)
    gx, gy = Inches(0.145), Inches(0.22)
    x0, y0 = MARGIN, Inches(1.66)
    for i, (title, desc) in enumerate(cards):
        x = x0 + (i % 4) * (cw + gx)
        y = y0 + (i // 4) * (ch + gy)
        rect(s, x, y, cw, ch, fill=LIGHT, line_color=LINE, line_w=Pt(0.5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        rect(s, x, y + Inches(0.16), Pt(3.2), Inches(0.42), fill=GOLD)
        textbox(s, x + Inches(0.2), y + Inches(0.16), cw - Inches(0.38),
                Inches(0.62), [para(title, Pt(12.5), NAVY, bold=True)],
                line_spacing=1.0)
        textbox(s, x + Inches(0.2), y + Inches(0.80), cw - Inches(0.38),
                ch - Inches(0.95), [para(desc, Pt(10.5), SLATE)],
                line_spacing=1.1)

    textbox(s, MARGIN, Inches(6.55), CONTENT_W, Inches(0.4),
            [para("45+ covered companies, plus systematic work on private AI "
                  "leaders where public-market exposure is indirect.",
                  Pt(11.5), GREY, italic=True)])


def s12_data_assets(prs):
    s = blank_slide(prs)
    slide_header(s, "Proprietary data assets",
                 "Trackers clients cannot get anywhere else", 12)

    left = [
        ("Token consumption series", "Usage by model family, workload and "
         "aggregator, updated monthly"),
        ("LLM pricing index", "Blended $/M-token price across vendors, "
         "weighted by observed traffic"),
        ("Datacenter capacity database", "Project-level GW pipeline with "
         "power source, developer and timing"),
        ("Hyperscaler capex tracker", "Quarterly capex, leases and "
         "commitments with GPU split estimates"),
    ]
    right = [
        ("Developer activity monitor", "Quarterly commit and repository "
         "trends as a proxy for AI coding adoption"),
        ("Corporate AI adoption panel", "Enterprise deployment, spend "
         "intentions and use-case maturity data"),
        ("IT vs labor spend series", "Long-run substitution data framing "
         "the AI automation opportunity"),
        ("Agent ecosystem telemetry", "MCP server downloads, agent token "
         "usage and tool-call growth metrics"),
    ]
    for col, items in enumerate((left, right)):
        x = MARGIN + col * Inches(4.12)
        y = Inches(1.66)
        for title, desc in items:
            rect(s, x, y, Inches(3.95), Inches(1.12), fill=LIGHT,
                 line_color=LINE, line_w=Pt(0.5),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
            textbox(s, x + Inches(0.2), y + Inches(0.12), Inches(3.6),
                    Inches(0.3), [para(title, Pt(12), NAVY, bold=True)])
            textbox(s, x + Inches(0.2), y + Inches(0.44), Inches(3.6),
                    Inches(0.62), [para(desc, Pt(10), SLATE)],
                    line_spacing=1.06)
            y += Inches(1.26)

    cx = Inches(9.05)
    cw = SLIDE_W - cx - MARGIN
    chart_title(s, cx, Inches(1.62), cw, "Data asset mix",
                "25+ maintained trackers \u2014 illustrative")
    chart = base_chart(
        s, XL_CHART_TYPE.DOUGHNUT, cx, Inches(2.25), cw, Inches(3.9),
        ["Usage & tokens", "Capex & capacity", "Pricing", "Adoption surveys",
         "Supply chain"],
        [("Trackers", (8, 7, 4, 3, 3))])
    pts = chart.plots[0].series[0].points
    colors = [BLUE, NAVY_2, GOLD, TEAL, SKY]
    for i, pt in enumerate(pts):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
    chart.legend.font.size = Pt(8.5)

    textbox(s, cx, Inches(6.30), cw, Inches(0.6),
            [para("Every tracker ships as a living spreadsheet or dataset "
                  "clients can interrogate directly.", Pt(10), GREY,
                  italic=True)], line_spacing=1.08)


def s13_formats(prs):
    s = blank_slide(prs)
    slide_header(s, "Deliverables & cadence",
                 "The right format for every investment question", 13)

    cx = MARGIN
    cw = Inches(4.6)
    chart_title(s, cx, Inches(1.62), cw, "Publication mix, last 12 months",
                "160+ deliverables \u2014 illustrative")
    chart = base_chart(
        s, XL_CHART_TYPE.DOUGHNUT, cx, Inches(2.25), cw, Inches(4.2),
        ["Company notes", "Thematic deep dives", "Data trackers & sheets",
         "Financial models", "Primers & explainers", "Daily / weekly memos"],
        [("Mix", (34, 22, 18, 12, 8, 6))])
    pts = chart.plots[0].series[0].points
    colors = [NAVY, BLUE, GOLD, TEAL, SKY, GREY]
    for i, pt in enumerate(pts):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
    chart.legend.font.size = Pt(9)

    rx = Inches(5.75)
    rw = SLIDE_W - rx - MARGIN
    rows = [
        ("Daily", "Morning software & AI stock memo: what moved, what "
         "matters, what to do about it"),
        ("Weekly", "Token and datacenter tracker refreshes with delta "
         "commentary"),
        ("Quarterly", "Earnings previews and reviews across the AI complex; "
         "capex model updates within hours of prints"),
        ("Event-driven", "Rapid-response notes on model launches, compute "
         "deals, export controls and regulation"),
        ("Foundational", "Primers on quantum computing, LLM architecture "
         "innovation and agentic protocols for generalist PMs"),
    ]
    y = Inches(1.66)
    for cadence, desc in rows:
        rect(s, rx, y, Inches(1.55), Inches(0.88), fill=NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        textbox(s, rx, y, Inches(1.55), Inches(0.88),
                [para(cadence, Pt(11.5), WHITE, bold=True)],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bx = rx + Inches(1.70)
        rect(s, bx, y, rw - Inches(1.70), Inches(0.88), fill=LIGHT,
             line_color=LINE, line_w=Pt(0.5))
        textbox(s, bx + Inches(0.2), y, rw - Inches(2.1), Inches(0.88),
                [para(desc, Pt(11), SLATE)], anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.06)
        y += Inches(1.02)


def s14_thematic(prs):
    s = blank_slide(prs)
    slide_header(s, "Thematic & event-driven research",
                 "Connecting AI to geopolitics, regulation and macro", 14)

    cols = [
        ("Geopolitics & industrial policy", [
            "EU Chips Act 2.0: dependencies, subsidies and beneficiaries",
            "China 2030 AI ambitions and the semiconductor supply response",
            "Gulf sovereign AI programs and hyperscaler exposure mapping",
            "Export controls: scenario trees for accelerator supply"]),
        ("Regulation & society", [
            "Under-16 social app bans: revenue-at-risk by platform",
            "Quantum & post-quantum policy: security-spend implications",
            "AI safety and model-governance frameworks for investors",
            "Data rights and LLM licensing economics for content owners"]),
        ("Macro & cross-asset scenarios", [
            "Oil price shock scenarios for cloud margins and gig platforms",
            "Regional conflict stress tests of the AI supply chain",
            "Rates and the financing cost of the datacenter buildout",
            "AI capex sustainability under bear-case token pricing"]),
    ]
    cw = Inches(3.95)
    gx = Inches(0.19)
    y = Inches(1.66)
    ch = Inches(4.35)
    for i, (title, items) in enumerate(cols):
        x = MARGIN + i * (cw + gx)
        rect(s, x, y, cw, ch, fill=LIGHT, line_color=LINE, line_w=Pt(0.5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        rect(s, x, y, cw, Inches(0.52), fill=NAVY_2)
        textbox(s, x + Inches(0.2), y, cw - Inches(0.4), Inches(0.52),
                [para(title, Pt(12.5), WHITE, bold=True)],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        bullet_block(s, x + Inches(0.2), y + Inches(0.70), cw - Inches(0.4),
                     ch - Inches(0.85), [(0, t) for t in items],
                     size=Pt(10.5), gap=Pt(7), line_spacing=1.08)

    textbox(s, MARGIN, Inches(6.30), CONTENT_W, Inches(0.45),
            [para("The AI trade does not live inside a sector box \u2014 our "
                  "thematic work gives clients the cross-asset context that "
                  "single-stock notes cannot.", Pt(11.5), GREY, italic=True)])


def s15_impact(prs):
    s = blank_slide(prs)
    slide_header(s, "Client impact",
                 "Engagement that compounds the research", 15)

    stats = [("1,200+", "Client meetings, calls and video sessions on AI "
              "topics"),
             ("300+", "Bespoke data pulls and custom scenario requests "
              "fulfilled"),
             ("40+", "Conference presentations, expert panels and teach-ins"),
             ("15", "Cross-sector collaborations with energy, real estate, "
              "utilities and macro teams")]
    cw, ch = Inches(2.90), Inches(1.70)
    gx = Inches(0.21)
    x = MARGIN
    for v, l in stats:
        stat_card(s, x, Inches(1.66), cw, ch, v, l, value_size=Pt(27))
        x += cw + gx

    bullet_block(s, MARGIN, Inches(3.75), Inches(7.4), Inches(3.0), [
        (0, "Advisory, not just publishing", True),
        (1, "Standing AI strategy dialogues with CIOs and portfolio "
            "managers; models handed over, not just PDFs"),
        (0, "Corporate access with substance", True),
        (1, "Curated meetings with hyperscaler infrastructure leaders, "
            "power developers, model labs and supply chain executives"),
        (0, "Fast where it counts", True),
        (1, "Same-day reaction frameworks for model releases, compute "
            "deals and capex guidance changes"),
    ], size=Pt(12))

    rx = Inches(8.30)
    rw = SLIDE_W - rx - MARGIN
    rect(s, rx, Inches(3.75), rw, Inches(2.95), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    textbox(s, rx + Inches(0.3), Inches(4.00), rw - Inches(0.6), Inches(0.3),
            [para("WHAT CLIENTS SAY THEY VALUE", Pt(10.5), GOLD, bold=True)])
    textbox(s, rx + Inches(0.3), Inches(4.40), rw - Inches(0.6), Inches(2.2),
            [para("\u201cOne team that can argue the token bull case and the "
                  "power bottleneck bear case with the same dataset.\u201d",
                  Pt(13), WHITE, italic=True),
             para("", Pt(6), WHITE),
             para("\u2014 Composite of client feedback themes, 2026 review",
                  Pt(10), SKY)], line_spacing=1.15)


def s16_close(prs):
    s = blank_slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    rect(s, 0, Inches(7.34), SLIDE_W, Inches(0.16), fill=GOLD)

    textbox(s, MARGIN, Inches(1.3), Inches(11.0), Inches(0.35),
            [para(DEPT.upper(), Pt(12), GOLD, bold=True)])
    textbox(s, MARGIN, Inches(1.85), Inches(12.0), Inches(1.7),
            [para("Partner with us on the defining", Pt(38), WHITE, bold=True),
             para("investment debate of the decade.", Pt(38), WHITE,
                  bold=True)], line_spacing=1.05)

    items = [
        ("Breadth", "Six pillars, one platform: silicon, datacenters, power, "
         "cloud, models and software covered as a single system."),
        ("Depth", "Proprietary trackers and models that measure the buildout "
         "instead of narrating it."),
        ("Access", "Analysts, data and corporate connections available on "
         "demand for client-specific questions."),
    ]
    x = MARGIN
    w = Inches(3.85)
    for title, desc in items:
        rect(s, x, Inches(3.85), w, Inches(1.95), fill=NAVY_2,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        textbox(s, x + Inches(0.28), Inches(4.08), w - Inches(0.56),
                Inches(0.4), [para(title, Pt(15), GOLD, bold=True)])
        textbox(s, x + Inches(0.28), Inches(4.55), w - Inches(0.56),
                Inches(1.15), [para(desc, Pt(11.5), WHITE)],
                line_spacing=1.12)
        x += w + Inches(0.27)

    textbox(s, MARGIN, Inches(6.30), CONTENT_W, Inches(0.7),
            [para("Contact your research sales representative for analyst "
                  "access, data subscriptions and bespoke project work.",
                  Pt(12), SKY),
             para("This deck is a capabilities overview; all statistics and "
                  "chart data are illustrative and not investment advice.",
                  Pt(9), GREY)], line_spacing=1.3)


# ----------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [s01_title, s02_exec_summary, s03_value_chain, s04_pillars,
                s05_tokens, s06_datacenter, s07_capex, s08_semis, s09_power,
                s10_software, s11_companies, s12_data_assets, s13_formats,
                s14_thematic, s15_impact, s16_close]
    for build in builders:
        build(prs)

    out = "AI_Equity_Research_Contributions.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
