"""
Generate a PowerPoint deck cataloging AI Agent slides from company presentations,
filings, and webinars. Styled to match the Anthropic/Asana "Should I build an agent?"
slide aesthetic: off-white background, bold titles, clean table layout with
alternating beige/white rows.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette (matches original slide) ──────────────────────────────────
BG_CREAM   = RGBColor(0xF7, 0xF5, 0xF0)
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ROW_BEIGE  = RGBColor(0xEB, 0xE7, 0xDE)
ROW_WHITE  = RGBColor(0xF7, 0xF5, 0xF0)
TEXT_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GRAY   = RGBColor(0x55, 0x55, 0x55)
TEXT_MUTED  = RGBColor(0x88, 0x88, 0x88)
ACCENT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
HIGHLIGHT   = RGBColor(0xC0, 0x39, 0x2B)
BLUE_ACCENT = RGBColor(0x2C, 0x5F, 0x8A)
BLACK_BAR   = RGBColor(0x00, 0x00, 0x00)
SECTION_BG  = RGBColor(0x2D, 0x2D, 0x2D)
SECTION_TXT = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_black_bars(slide):
    """Add top and bottom black bars like the original slide."""
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.25)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BLACK_BAR
    top_bar.line.fill.background()

    bot_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25)
    )
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = BLACK_BAR
    bot_bar.line.fill.background()


def add_title(slide, title_text, left=Inches(0.8), top=Inches(0.5), width=Inches(11.5)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_BLACK
    p.font.name = "Arial"
    return txBox


def add_subtitle(slide, text, left=Inches(0.8), top=Inches(1.3), width=Inches(11.5)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_GRAY
    p.font.name = "Arial"
    return txBox


def add_source_footer(slide, source_text):
    txBox = slide.shapes.add_textbox(
        Inches(0.8), SLIDE_H - Inches(0.75), Inches(11.5), Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = source_text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Arial"
    p.font.italic = True


def _remove_cell_borders(cell):
    """Remove all borders from a cell via XML manipulation."""
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        existing = tcPr.find(qn(border_tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border_tag), w="0", cap="flat", cmpd="sng", algn="ctr")
        noFill = etree.SubElement(ln, qn('a:noFill'))


def _set_cell_style(cell, text, font_size=Pt(14), bold=False, color=TEXT_BLACK, bg=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.15)
    cell.margin_right = Inches(0.15)
    cell.margin_top = Inches(0.08)
    cell.margin_bottom = Inches(0.08)
    _remove_cell_borders(cell)


def add_decision_table(slide, rows_data, left=Inches(0.8), top=Inches(2.0),
                       width=Inches(11.5), col_widths=None):
    """
    rows_data: list of (question_text, answer_text) tuples
    Renders an alternating-row table matching the original slide style.
    """
    n_rows = len(rows_data)
    table_shape = slide.shapes.add_table(n_rows, 2, left, top, width, Inches(n_rows * 0.65))
    table = table_shape.table

    if col_widths is None:
        col_widths = [Inches(5.75), Inches(5.75)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for idx, (q, a) in enumerate(rows_data):
        bg = ROW_BEIGE if idx % 2 == 0 else ROW_WHITE
        _set_cell_style(table.cell(idx, 0), q, bold=True, bg=bg)
        _set_cell_style(table.cell(idx, 1), a, bg=bg)

    return table_shape


def add_framework_table(slide, header, rows_data, left=Inches(0.8), top=Inches(2.0),
                        width=Inches(11.5)):
    """Table with a header row and content rows."""
    n_cols = len(header)
    n_rows = len(rows_data) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                         Inches(n_rows * 0.55))
    table = table_shape.table

    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w

    for ci, h in enumerate(header):
        _set_cell_style(table.cell(0, ci), h, font_size=Pt(13), bold=True,
                        color=SECTION_TXT, bg=ACCENT_DARK)

    for ri, row in enumerate(rows_data):
        bg = ROW_BEIGE if ri % 2 == 0 else ROW_WHITE
        for ci, val in enumerate(row):
            _set_cell_style(table.cell(ri + 1, ci), val, font_size=Pt(12), bg=bg)

    return table_shape


def add_bullet_slide(slide, title, bullets, source, subtitle=None):
    set_slide_bg(slide, BG_CREAM)
    add_black_bars(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)
    start_top = Inches(2.0) if not subtitle else Inches(1.9)

    txBox = slide.shapes.add_textbox(Inches(0.8), start_top, Inches(11.5),
                                     Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_BLACK
        p.font.name = "Arial"
        p.space_after = Pt(8)
        p.level = 0
    add_source_footer(slide, source)


# ═══════════════════════════════════════════════════════════════════════════════
#  BUILD THE PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank

# ── SLIDE 1: Cover ───────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)

txBox = sl.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AI Agent Slides"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = TEXT_BLACK
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "From Company Presentations, Filings & Webinars"
p2.font.size = Pt(28)
p2.font.color.rgb = TEXT_GRAY
p2.font.name = "Arial"
p2.space_before = Pt(12)

txBox2 = sl.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
desc = tf2.paragraphs[0]
desc.text = (
    "Curated collection of publicly available decision frameworks, architecture "
    "diagrams, and strategy slides about AI agents from 30+ major technology "
    "companies, enterprise SaaS vendors, consulting firms, and venture capital investors."
)
desc.font.size = Pt(16)
desc.font.color.rgb = TEXT_GRAY
desc.font.name = "Arial"

p3 = tf2.add_paragraph()
p3.text = "February 2025"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED
p3.font.name = "Arial"
p3.space_before = Pt(20)


# ── SLIDE 2: Table of Contents ───────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Companies Covered")
companies = [
    ["Anthropic", "OpenAI", "Google / DeepMind", "Microsoft", "Salesforce"],
    ["ServiceNow", "AWS", "NVIDIA", "SAP", "Workday"],
    ["HubSpot", "Palantir", "Snowflake", "Databricks", "Atlassian"],
    ["Cognition AI", "LangChain", "CrewAI", "Sequoia Capital", "a16z"],
    ["McKinsey", "Gartner", "Accenture", "Asana", "Cisco"],
    ["Box", "Intuit", "Meta", "Apple", "Klarna"],
]
add_framework_table(
    sl,
    header=["Model Providers & Infra", "Enterprise SaaS", "Enterprise SaaS (cont.)",
            "Agent Frameworks", "Investors & Analysts"],
    rows_data=[
        ["Anthropic", "Salesforce", "HubSpot", "Cognition AI (Devin)", "Sequoia Capital"],
        ["OpenAI", "ServiceNow", "Palantir", "LangChain / LangGraph", "a16z"],
        ["Google / DeepMind", "AWS", "Snowflake", "CrewAI", "McKinsey"],
        ["Microsoft", "NVIDIA", "Databricks", "", "Gartner"],
        ["Meta", "SAP", "Atlassian", "", "Accenture"],
        ["Apple", "Workday", "Box / Intuit", "", ""],
    ],
    top=Inches(1.8),
)
add_source_footer(sl, "30+ companies across technology, enterprise SaaS, AI infrastructure, and advisory firms")


# ── SLIDE 3: Section Divider — Model Providers ──────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, SECTION_BG)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "MODEL PROVIDERS & AI INFRASTRUCTURE"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = SECTION_TXT
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "Anthropic  |  OpenAI  |  Google  |  Microsoft  |  AWS  |  NVIDIA  |  Meta  |  Apple"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.font.name = "Arial"
p2.space_before = Pt(16)


# ── SLIDE 4: Anthropic / Asana — "Should I build an agent?" ─────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, '"Should I build an agent?"')
add_decision_table(sl, [
    ("Is the task complex enough?",
     "No \u2192 Workflows\nYes \u2192 Agents"),
    ("Is the task valuable enough?",
     "<$0.1 \u2192 Workflows\n>$1 \u2192 Agents"),
    ("Are all parts of the task doable?",
     "No \u2192 Reduce scope\nYes \u2192 Agents"),
    ("What is the cost of error/error discovery?",
     "High \u2192 Human-in-the-loop\nLow \u2192 Agents"),
])
add_source_footer(sl, "Source: Anthropic / Asana joint presentation (2024-2025)  \u2014  anthropic.com/research/building-effective-agents")


# ── SLIDE 5: Anthropic — Building Effective Agents ──────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Anthropic: Building Effective Agents")
add_subtitle(sl, "Key frameworks from the canonical agent design guide (Dec 2024)")
add_decision_table(sl, [
    ("Agents vs. Workflows spectrum",
     "Workflows: LLMs orchestrated through predefined code paths\n"
     "Agents: LLMs dynamically direct their own processes & tool usage"),
    ("Agent loop",
     "Environment \u2192 LLM call \u2192 Tool use \u2192 Environment \u2192 (repeat)"),
    ("Augmented LLM",
     "LLM + Retrieval + Tools + Memory = foundational building block"),
    ("Common workflow patterns",
     "Prompt chaining, Routing, Parallelization,\n"
     "Orchestrator-workers, Evaluator-optimizer"),
    ("When to use agents",
     "Tasks requiring flexibility & model-driven decisions \u2192 Agents\n"
     "Tasks needing predictability & consistency \u2192 Workflows"),
], top=Inches(2.2))
add_source_footer(sl, "Source: anthropic.com/research/building-effective-agents (Dec 2024)")


# ── SLIDE 6: OpenAI — Practical Guide to Building Agents ───────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "OpenAI: A Practical Guide to Building Agents")
add_subtitle(sl, "Dev Day 2024 keynote & Agents SDK (March 2025)")
add_decision_table(sl, [
    ("Agent definition",
     '"An agent is a system that independently accomplishes\ntasks on behalf of a user."'),
    ("Core components",
     "Agent = Model + Tools + Instructions + Knowledge"),
    ("When to go agentic",
     "Multi-step reasoning + Tool use + Real-world actions \u2192 Agent\n"
     "Single-step, no tools needed \u2192 Simple API call"),
    ("Agent orchestration spectrum",
     "Function calling \u2192 Assistants API \u2192 Custom agent loops\n"
     "\u2192 Multi-agent systems (Swarm)"),
    ("Agent evaluation",
     "Task completion rate, Cost per task,\nLatency, Safety metrics"),
], top=Inches(2.2))
add_source_footer(sl, "Source: OpenAI Dev Day (2024)  |  cdn.openai.com/business-guides-and-resources/a-practical-guide-to-building-agents.pdf  |  Agents SDK (Mar 2025)")


# ── SLIDE 7: OpenAI — Swarm & Operator ──────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "OpenAI: Multi-Agent & Computer-Using Agents")
add_decision_table(sl, [
    ("Swarm framework (Oct 2024)",
     "Multi-agent handoff architecture\n"
     'Core primitives: "Routines" and "Handoffs"\n'
     "Agent-to-agent communication patterns"),
    ("Operator / CUA (Jan 2025)",
     "Agent interacting with web browsers autonomously\n"
     "Task completion workflow with oversight guardrails\n"
     "Safety controls for autonomous web navigation"),
    ("Agents SDK (Mar 2025)",
     "Agent loop: tool calls \u2192 handoffs \u2192 guardrails\n"
     "Built-in tracing and observability\n"
     "Multi-agent orchestration primitives"),
], top=Inches(1.8))
add_source_footer(sl, "Source: OpenAI blog  |  openai.com/index/new-tools-for-building-agents/")


# ── SLIDE 8: Google / DeepMind ──────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Google / DeepMind: The Age of Agents")
add_subtitle(sl, "Cloud Next 2024-2025  |  Google I/O 2025")
add_decision_table(sl, [
    ("AI maturity model",
     "Chatbots \u2192 Assistants \u2192 Agents \u2192 Multi-agent systems"),
    ("Vertex AI Agent Builder",
     "Gemini models + Tools + Data stores\n"
     "Customer agent vs. Employee agent deployment patterns"),
    ("Agent grounding",
     "Google Search grounding + Enterprise data\n"
     "Real-time data connection for agent accuracy"),
    ("Google I/O 2025 launches",
     'Project Mariner (Chrome browser agent)\n'
     'Jules (AI coding agent)\n'
     "Agent-to-Agent protocol (A2A) specification"),
    ("Agent maturity ladder",
     "Retrieval \u2192 Extensions \u2192 Function calling\n"
     "\u2192 Data agents \u2192 Multi-agent"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Google Cloud Next (2024-2025)  |  Google I/O 2025  |  developers.googleblog.com/en/a2a-a-new-era-of-agent-interoperability/")


# ── SLIDE 9: Microsoft ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Microsoft: Copilot as Agent Orchestrator")
add_subtitle(sl, "Build 2024-2025  |  Ignite 2024  |  AutoGen Research")
add_decision_table(sl, [
    ('"AI agents will reshape every\nsoftware category" — Nadella',
     "SaaS \u2192 AI agents as the new application paradigm\n"
     "Every business process rebuilt around agents"),
    ("Copilot agent types",
     "Prompt-response \u2192 RAG-augmented \u2192 Autonomous agents\n"
     "Built in Copilot Studio (low-code)"),
    ("AutoGen multi-agent framework",
     "Multiple specialized agents conversing to solve tasks\n"
     '"Conversable agents" as composable building blocks'),
    ("Agent memory & planning",
     "Agents maintain state across sessions\n"
     "Long-term memory for personalized experiences"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Microsoft Build (2024-2025)  |  Ignite 2024  |  microsoft.github.io/autogen/  |  10-K FY2024-2025")


# ── SLIDE 10: Section Divider — Enterprise SaaS ────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, SECTION_BG)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ENTERPRISE SaaS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = SECTION_TXT
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "Salesforce  |  ServiceNow  |  AWS  |  NVIDIA  |  SAP  |  Workday  |  HubSpot  |  Palantir  |  and more"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.font.name = "Arial"
p2.space_before = Pt(16)


# ── SLIDE 11: Salesforce — Agentforce ───────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, 'Salesforce: "The Third Wave of AI"')
add_subtitle(sl, "Dreamforce 2024  |  Agentforce Launch  |  Marc Benioff keynote")
add_decision_table(sl, [
    ("Three waves of AI",
     "Predictive AI \u2192 Generative AI \u2192 Agentic AI\n"
     "(Autonomous Agents)"),
    ("Agentforce architecture",
     "Atlas Reasoning Engine + Data Cloud\n"
     "+ Trust Layer + Tools"),
    ("Digital Labor model",
     'Agents as "digital workers"\n'
     "Billed per conversation, not per seat"),
    ("Agent builder (low-code)",
     "Topics \u2192 Instructions \u2192 Actions \u2192 Guardrails"),
    ("Agent types",
     "Service Agent, Sales Agent, Marketing Agent,\n"
     "Commerce Agent, Custom agents"),
    ("Trust & safety",
     "Einstein Trust Layer: toxicity detection,\n"
     "PII masking, prompt injection defense"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Dreamforce 2024  |  salesforce.com/agentforce/  |  Q3/Q4 FY2025 Earnings")


# ── SLIDE 12: ServiceNow ───────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "ServiceNow: AI Agents for Enterprise Workflows")
add_subtitle(sl, "Knowledge 2024-2025  |  Investor Day")
add_decision_table(sl, [
    ("Now Assist Agents",
     "Agentic AI built into IT, HR, Customer Service,\nSecurity workflows"),
    ("Evolution",
     "Rule-based virtual agents \u2192 LLM-powered\nautonomous AI agents"),
    ("Agent orchestration",
     "Agents interact with ServiceNow's workflow engine\nfor multi-step task resolution"),
    ("Key metric",
     "Now Assist agent attach rates in new deals\n"
     '"Agentic workflows" as future of enterprise automation'),
], top=Inches(2.2))
add_source_footer(sl, "Source: ServiceNow Knowledge (2024-2025)  |  servicenow.com/now-platform/ai-agents.html  |  10-K / Earnings 2024-2025")


# ── SLIDE 13: AWS ──────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "AWS: Agents for Amazon Bedrock")
add_subtitle(sl, "re:Invent 2024  |  Matt Garman & Swami Sivasubramanian keynotes")
add_decision_table(sl, [
    ("Bedrock Agents architecture",
     "Foundation model + Action groups\n"
     "+ Knowledge bases + Guardrails"),
    ("Agent orchestration flow",
     "User request \u2192 Agent reasoning \u2192 Action execution\n"
     "\u2192 Response synthesis"),
    ("Multi-agent collaboration",
     "Supervisor agent delegating to specialist agents"),
    ("Automated Reasoning",
     "Formal verification layer to check agent outputs\n"
     "Mathematical proofs for correctness guarantees"),
    ("Amazon Q",
     "AI agent for developers and business users\n"
     "Integrated across AWS console, IDE, and business apps"),
], top=Inches(2.2))
add_source_footer(sl, "Source: AWS re:Invent 2024  |  aws.amazon.com/bedrock/agents/  |  Amazon Shareholder Letters 2024-2025")


# ── SLIDE 14: NVIDIA ──────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, 'NVIDIA: "Agentic AI Is the Next Platform"')
add_subtitle(sl, "GTC 2025  |  Jensen Huang keynote  |  Investor Day")
add_decision_table(sl, [
    ('"Every company will have AI agent\ndepartments" — Jensen Huang',
     "Agentic AI as the next computing platform\n"
     "Data centers as \"AI Factories\""),
    ("Agent infrastructure stack",
     "Training \u2192 Inference \u2192 Agent frameworks\n"
     "\u2192 Deployment (NIM, NeMo, ACE)"),
    ("Token economics",
     "Agents consume 10-100x more tokens\nthan single-turn queries"),
    ("Physical AI agents",
     "Robots and autonomous systems as\nembodied agents (Isaac, Cosmos)"),
    ("TAM expansion",
     "Inference \u2192 Agentic reasoning workloads\n"
     "Agent workloads drive next wave of GPU demand"),
], top=Inches(2.2))
add_source_footer(sl, "Source: NVIDIA GTC 2025  |  Investor Day 2025  |  Q4 FY2025 Earnings")


# ── SLIDE 15: SAP ─────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "SAP: Joule AI Agent")
add_subtitle(sl, "SAP Sapphire / TechEd 2024-2025")
add_decision_table(sl, [
    ("Joule as collaborative AI agent",
     "Embedded across all SAP applications"),
    ("Business AI agents",
     "Agents for procurement, finance, HR,\nsupply chain"),
    ("Agent-to-agent orchestration",
     "Joule coordinating with third-party agents"),
    ("Business context grounding",
     "Agents grounded in SAP business data\nand processes"),
], top=Inches(2.2))
add_source_footer(sl, "Source: SAP Sapphire / TechEd (2024-2025)  |  sap.com/products/artificial-intelligence/ai-assistant.html")


# ── SLIDE 16: Workday, HubSpot, Atlassian ──────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Enterprise SaaS: Agent Announcements")
add_subtitle(sl, "Workday  |  HubSpot  |  Atlassian  |  Box  |  Intuit")
add_decision_table(sl, [
    ("Workday AI Agents",
     "Recruiter Agent, Expenses Agent, Succession Agent\n"
     "Agent-driven automation in HCM & financials"),
    ("HubSpot Breeze AI Agents",
     "Prospecting Agent, Content Agent, Social Agent,\n"
     "Customer Agent — GTM automation"),
    ("Atlassian Rovo Agents",
     "AI agents for Jira, Confluence, and teamwork\n"
     'Agent marketplace + "virtual teammate" concept'),
    ("Box AI Agents",
     "Content-grounded agents with enterprise permissions\n"
     "Document extraction, summarization, and Q&A"),
    ("Intuit Assist",
     "AI agent for tax, accounting, personal finance\n"
     "Powered by GenOS platform"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Workday Rising  |  HubSpot INBOUND  |  Atlassian Team '24  |  BoxWorks  |  Intuit Investor Day (2024-2025)")


# ── SLIDE 17: Palantir ────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Palantir: AIP Agent Platform")
add_subtitle(sl, "AIPCon events  |  Investor presentations (2024-2025)")
add_decision_table(sl, [
    ("AIP boot camp \u2192 production pipeline",
     "Prototyping to deployed agents in weeks\n"
     "Enterprise adoption playbook"),
    ("Agent orchestration in AIP",
     "LLM + Ontology + Actions + Guardrails"),
    ('"Warp Speed" defense agents',
     "AI agents for military/defense decision support"),
    ("Agent ROI case studies",
     "Specific customer metrics on agent deployments\n"
     "AIP agent adoption as leading growth driver"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Palantir AIPCon  |  palantir.com/platforms/aip/  |  10-K / Earnings 2024-2025")


# ── SLIDE 18: Snowflake & Databricks ──────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Data Platforms: Agents for Analytics")
add_subtitle(sl, "Snowflake Summit  |  Databricks Data + AI Summit (2024-2025)")
add_decision_table(sl, [
    ("Snowflake Cortex Agents",
     "AI agents that operate on Snowflake data\n"
     "Agentic RAG: search and analyze data autonomously"),
    ("Snowflake Intelligence",
     "Natural language data agent\n"
     "Tool-use: SQL generation, chart creation, pipeline orchestration"),
    ('Databricks: "Compound AI Systems"',
     "Agents as compound systems: multiple models,\n"
     "retrieval, tools, code"),
    ("Mosaic AI Agent Framework",
     "Build, evaluate, deploy agents on Databricks\n"
     "Automated evaluation of quality, latency, cost"),
    ("Maturity curve",
     "Fine-tuned models \u2192 RAG \u2192 Agents \u2192 Multi-agent"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Snowflake Summit  |  Databricks Data + AI Summit (2024-2025)  |  snowflake.com  |  databricks.com")


# ── SLIDE 19: Section Divider — Agent Frameworks & Startups ────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, SECTION_BG)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "AGENT FRAMEWORKS & STARTUPS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = SECTION_TXT
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "Cognition AI (Devin)  |  LangChain / LangGraph  |  CrewAI"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.font.name = "Arial"
p2.space_before = Pt(16)


# ── SLIDE 20: Cognition AI (Devin) ────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Cognition AI: Devin — The First AI Software Engineer")
add_subtitle(sl, "Launch presentation (March 2024)")
add_decision_table(sl, [
    ("Autonomous coding agent",
     "Fully autonomous software engineering agent\n"
     "Plans, codes, debugs, and deploys"),
    ("Agent workspace",
     "Terminal, browser, code editor — all controlled\n"
     "by the agent simultaneously"),
    ("Self-correction",
     "Multi-step task execution with planning,\n"
     "error detection, and self-repair"),
    ("Benchmark results",
     "SWE-bench: resolved real-world GitHub issues\n"
     "autonomously"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Cognition AI launch demo (March 2024)  |  cognition.ai")


# ── SLIDE 21: LangChain / LangGraph ──────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "LangChain / LangGraph: Agent Architecture Patterns")
add_subtitle(sl, '"State of AI Agents" reports  |  "In the Loop" survey (2024-2025)')
add_decision_table(sl, [
    ("Agent architecture spectrum",
     "Simple chains \u2192 Routers \u2192 Tool-using agents\n"
     "\u2192 Multi-agent \u2192 Autonomous agents"),
    ("LangGraph patterns",
     "State machines for agent orchestration:\n"
     "nodes, edges, conditional routing"),
    ('"Cognitive architecture"',
     "How an agent thinks: planning, memory, tool use\n"
     "ReAct pattern = Reasoning + Acting loop"),
    ("Human-in-the-loop",
     "Interrupt, approve, edit, resume patterns\n"
     "in agent workflows"),
    ('"In the Loop" survey (2025)',
     "51% of companies have agents in production\n"
     "Top challenges: reliability, evaluation, latency"),
], top=Inches(2.2))
add_source_footer(sl, "Source: blog.langchain.dev  |  LangGraph documentation  |  \"In the Loop\" survey 2025")


# ── SLIDE 22: CrewAI ─────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "CrewAI: Multi-Agent Orchestration")
add_subtitle(sl, "Conference presentations and documentation (2024-2025)")
add_decision_table(sl, [
    ("Role-based agent design",
     "Each agent has: Role, Goal, Backstory, Tools"),
    ("Crew patterns",
     "Sequential: agents work one after another\n"
     "Hierarchical: manager agent delegates to workers"),
    ('"Crews" as the unit',
     "Multi-agent collaboration organized as crews\n"
     "Composable, reusable agent teams"),
], top=Inches(2.2))
add_source_footer(sl, "Source: crewai.com  |  CrewAI conference talks (2024-2025)")


# ── SLIDE 23: Section Divider — Investors & Analysts ──────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, SECTION_BG)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "INVESTORS & ANALYSTS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = SECTION_TXT
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "Sequoia Capital  |  a16z  |  McKinsey  |  Gartner  |  Accenture"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.font.name = "Arial"
p2.space_before = Pt(16)


# ── SLIDE 24: Sequoia Capital ────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, 'Sequoia Capital: "Software That Works for You"')
add_subtitle(sl, "Sonya Huang — AI Agents presentation  |  Arc conference (2024)")
add_decision_table(sl, [
    ("New paradigm",
     "From tools you use \u2192 tools that use themselves\n"
     "Agents as the new software paradigm"),
    ("Agent market map",
     "Landscape of agent startups across verticals"),
    ("Value chain",
     "Model providers \u2192 Agent frameworks\n"
     "\u2192 Vertical agents \u2192 Agent infrastructure"),
    ("Adoption curve",
     "Copilots (now) \u2192 Task agents (near-term)\n"
     "\u2192 Autonomous agents (future)"),
    ("Agent economics (2025)",
     'Agent-native vs. agent-added companies\n'
     "Cost per task, margin structure benchmarks"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Sequoia Arc conference (2024)  |  sequoiacap.com/article/ai-agents-part-i/  |  \"AI in the Real World\" (2025)")


# ── SLIDE 25: a16z ──────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, 'a16z: "AI Agents Are the New Apps"')
add_subtitle(sl, '"Big Ideas 2025"  |  "Emerging Architectures for LLM Agents" (2024)')
add_decision_table(sl, [
    ("Agents as application layer",
     "Agents are the application layer of the AI stack\n"
     "The \"agent economy\" — agents transact on behalf of users"),
    ("Agent infrastructure stack",
     "Observability, Evaluation, Memory,\n"
     "Tool integration, Orchestration"),
    ("Enterprise adoption curve",
     "Internal tools \u2192 Customer-facing\n"
     "\u2192 Fully autonomous"),
    ("Agent memory patterns",
     "Short-term (context window)\n"
     "Long-term (vector DB) | Episodic"),
    ("Reference architecture",
     "Agent Core (LLM + Planning + Memory)\n"
     "+ Tools + Environment"),
], top=Inches(2.2))
add_source_footer(sl, "Source: a16z.com  |  \"Emerging Architectures for LLM Applications\"  |  \"Big Ideas 2025\"")


# ── SLIDE 26: McKinsey ──────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "McKinsey: Economic Potential of AI Agents")
add_subtitle(sl, "McKinsey Global Institute reports (2024-2025)")
add_decision_table(sl, [
    ("Impact estimate",
     "Agentic AI could automate 25% of work tasks\nby 2030"),
    ("Automation potential by industry",
     "Heat map of agentic AI potential across sectors\n"
     "Highest: customer ops, software, finance"),
    ("Deployment maturity model",
     "Pilot \u2192 Scaled \u2192 Embedded \u2192 Autonomous"),
    ('"From copilot to autopilot"',
     "Evolution of AI assistance in enterprise\n"
     "Human-augmented \u2192 Human-supervised \u2192 Autonomous"),
    ("ROI framework",
     "Cost savings + Revenue uplift + Risk reduction"),
], top=Inches(2.2))
add_source_footer(sl, "Source: McKinsey Global Institute  |  mckinsey.com/capabilities/mckinsey-digital/our-insights/")


# ── SLIDE 27: Gartner ──────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Gartner: AI Agent Hype Cycle & Forecasts")
add_subtitle(sl, "IT Symposium / Hype Cycle (2024-2025)")
add_decision_table(sl, [
    ("Hype Cycle positioning",
     'AI agents on "Peak of Inflated Expectations" (2024)\n'
     '\u2192 "Trough of Disillusionment" (2025)'),
    ("Market forecast",
     "By 2028, 33% of enterprise software\nwill include agentic AI"),
    ("Agent design patterns",
     "Gartner taxonomy: Reactive, Deliberative,\n"
     "Hybrid, Multi-agent"),
    ("Agent governance framework",
     "Trust, Accountability, Oversight,\n"
     "Compliance for enterprise agents"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Gartner IT Symposium  |  gartner.com/en/articles/what-s-new-in-artificial-intelligence-from-the-2024-gartner-hype-cycle")


# ── SLIDE 28: Accenture ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Accenture: AI Agents — The New Digital Workforce")
add_subtitle(sl, "Technology Vision 2025")
add_decision_table(sl, [
    ("Agent ecosystems",
     "Multiple agents working together across\nenterprise functions"),
    ("From automation to autonomy",
     "Maturity model: RPA \u2192 Intelligent automation\n"
     "\u2192 Copilots \u2192 Autonomous agents"),
    ("Agent trust framework",
     "Transparency, Explainability, Accountability"),
    ("Industry case studies",
     "Industry-specific agent deployment examples\nacross verticals"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Accenture Technology Vision 2025  |  accenture.com")


# ── SLIDE 29: Section Divider — Case Studies ─────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, SECTION_BG)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CASE STUDIES & ADDITIONAL COMPANIES"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = SECTION_TXT
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "Klarna  |  Asana  |  Cisco  |  Meta  |  Apple"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.font.name = "Arial"
p2.space_before = Pt(16)


# ── SLIDE 30: Klarna ────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Klarna: AI Agent Case Study")
add_subtitle(sl, "Customer service AI agent deployment (2024)")
add_decision_table(sl, [
    ("Conversations handled by AI agent",
     "2.3 million in first month"),
    ("Equivalent human workforce",
     "~700 full-time customer service agents replaced"),
    ("Resolution time",
     "Reduced from 11 minutes to 2 minutes"),
    ("Repeat inquiries",
     "Dropped 25%"),
    ("Customer satisfaction",
     "On par with human agents"),
    ("Estimated annual savings",
     "$40 million"),
], top=Inches(2.2))
add_source_footer(sl, "Source: Klarna blog  |  klarna.com/international/press/klarna-ai-assistant/  |  Investor presentations (2024)")


# ── SLIDE 31: Asana, Cisco, Meta, Apple ────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Additional Company Agent Frameworks")
add_decision_table(sl, [
    ('Asana: "AI Teammates"',
     "Agents as team members in project management\n"
     '"Work Graph" provides context layer for agents'),
    ("Cisco: Enterprise Agent Security",
     "Agents for network ops, security ops, CX\n"
     "Enterprise deployment with security guardrails"),
    ("Meta: Consumer AI Agents",
     "AI agents as social companions in AR/VR\n"
     "Business AI agents for WhatsApp & Messenger"),
    ("Apple: Privacy-First Agents",
     "Siri as cross-app orchestration agent\n"
     "On-device + Private Cloud Compute architecture"),
], top=Inches(1.8))
add_source_footer(sl, "Source: Asana Forward  |  Cisco Live  |  Meta Connect  |  WWDC 2024-2025")


# ── SLIDE 32: Section Divider — Cross-Company Analysis ──────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, SECTION_BG)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "CROSS-COMPANY ANALYSIS"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = SECTION_TXT
p.font.name = "Arial"
p2 = tf.add_paragraph()
p2.text = "Common frameworks, architecture patterns, and maturity models"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
p2.font.name = "Arial"
p2.space_before = Pt(16)


# ── SLIDE 33: Common Decision Framework ─────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Cross-Company: When to Deploy Agents")
add_subtitle(sl, "Common decision dimensions across 30+ company presentations")
add_framework_table(
    sl,
    header=["Dimension", "Threshold for Agents", "Alternative"],
    rows_data=[
        ["Task Complexity", "Multi-step, requires reasoning", "Simple automation / workflows"],
        ["Task Value", "High value per task (>$1-10)", "Low-value \u2192 batch automation"],
        ["Error Tolerance", "Low cost of error, easy to verify", "High-stakes \u2192 human-in-the-loop"],
        ["Data Availability", "Rich context available", "Poor data \u2192 improve data first"],
        ["Repeatability", "Frequent, recurring tasks", "One-off \u2192 manual"],
        ["Tool Integration", "APIs and tools available", "No APIs \u2192 build integrations first"],
    ],
    top=Inches(2.2),
)
add_source_footer(sl, "Composite analysis across Anthropic, OpenAI, Google, Microsoft, Salesforce, AWS, and others")


# ── SLIDE 34: Agent Maturity Model (Composite) ─────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Composite Agent Maturity Model")
add_subtitle(sl, "Synthesized from Google, Databricks, McKinsey, Sequoia, Accenture, and others")
add_framework_table(
    sl,
    header=["Level", "Stage", "Description", "Examples"],
    rows_data=[
        ["Level 0", "Chatbot", "Single-turn Q&A, no tools", "Basic ChatGPT, FAQ bots"],
        ["Level 1", "Copilot", "Human-directed, tool-assisted", "GitHub Copilot, Copilot for M365"],
        ["Level 2", "Task Agent", "Autonomous single-task execution", "Coding agents, CS agents"],
        ["Level 3", "Multi-Agent", "Multiple agents collaborating", "AutoGen, CrewAI orchestrations"],
        ["Level 4", "Autonomous", "Fully autonomous, long-running", "AI employees, digital workers"],
    ],
    top=Inches(2.2),
)
add_source_footer(sl, "Composite model: Google Cloud Next  |  Databricks Summit  |  McKinsey  |  Sequoia  |  Accenture Technology Vision")


# ── SLIDE 35: Agent Architecture Patterns ──────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Common Agent Architecture Patterns")
add_subtitle(sl, "Converging design from Anthropic, OpenAI, AWS, LangChain, a16z")

txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
    "\u2502              USER / TRIGGER                  \u2502",
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
    "                        \u2502",
    "                        \u25bc",
    "\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
    "\u2502           AGENT ORCHESTRATOR                 \u2502",
    "\u2502  [ Planning ]  [ Memory ]  [ Guardrails ]    \u2502",
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
    "                        \u2502",
    "          \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u253c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
    "          \u25bc            \u25bc            \u25bc",
    "    [ Tool A ]    [ Tool B ]    [ Tool C ]",
    "    ( API  )      ( Search )    (  Code  )",
]
for i, line in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_BLACK
    p.font.name = "Courier New"
    p.space_after = Pt(1)
    p.alignment = PP_ALIGN.CENTER

add_source_footer(sl, "Architecture pattern converges across: Anthropic, OpenAI, AWS Bedrock, LangChain, a16z")


# ── SLIDE 36: Key CEO Quotes ────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Key CEO Quotes on AI Agents")
add_subtitle(sl, "From earnings calls and keynotes (2024-2025)")
add_decision_table(sl, [
    ("Marc Benioff\nSalesforce, Q3 FY2025",
     '"Agentforce is the biggest thing we\'ve ever done.\n'
     'This is the third wave of AI."'),
    ("Satya Nadella\nMicrosoft, Q2 FY2025",
     '"Every customer I talk to is looking to reshape\n'
     'their business processes with AI agents."'),
    ("Jensen Huang\nNVIDIA, Q4 FY2025",
     '"Agents require 10x to 100x more compute.\n'
     'They reason, plan, and act."'),
    ("Sundar Pichai\nGoogle, Q4 2024",
     '"Agents are the killer app for\nlarge language models."'),
    ("Alex Karp\nPalantir, Q4 2024",
     '"The demand for AI agents that can do things\n'
     'in the real world is unlike anything we\'ve seen."'),
    ("Bill McDermott\nServiceNow, Q4 2024",
     '"AI agents will become the standard way\n'
     'enterprises interact with platforms."'),
], top=Inches(2.2))
add_source_footer(sl, "Source: Public earnings call transcripts  |  Conference keynotes (2024-2025)")


# ── SLIDE 37: Sources & References ──────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
add_title(sl, "Sources & References")
add_framework_table(
    sl,
    header=["Company", "Source Type", "Reference"],
    rows_data=[
        ["Anthropic", "Blog", "anthropic.com/research/building-effective-agents"],
        ["OpenAI", "Guide / SDK", "cdn.openai.com/.../a-practical-guide-to-building-agents.pdf"],
        ["Google", "Cloud Next / I/O", "cloud.google.com/products/agent-builder"],
        ["Microsoft", "Build / Research", "microsoft.github.io/autogen/"],
        ["Salesforce", "Dreamforce", "salesforce.com/agentforce/"],
        ["ServiceNow", "Knowledge", "servicenow.com/now-platform/ai-agents.html"],
        ["AWS", "re:Invent", "aws.amazon.com/bedrock/agents/"],
        ["NVIDIA", "GTC", "nvidia.com/en-us/ai/"],
        ["SAP", "Sapphire", "sap.com/.../ai-assistant.html"],
        ["Palantir", "AIPCon", "palantir.com/platforms/aip/"],
        ["Snowflake", "Summit", "snowflake.com/en/data-cloud/cortex/"],
        ["Databricks", "Summit", "databricks.com/.../build-genai-apps"],
        ["LangChain", "Blog", "blog.langchain.dev"],
        ["Sequoia", "Blog", "sequoiacap.com/article/ai-agents-part-i/"],
        ["a16z", "Blog", "a16z.com/emerging-architectures-for-llm-applications/"],
        ["McKinsey", "Report", "mckinsey.com/.../our-insights/"],
        ["Gartner", "Hype Cycle", "gartner.com/.../2024-gartner-hype-cycle"],
        ["Klarna", "Blog", "klarna.com/international/press/klarna-ai-assistant/"],
    ],
    top=Inches(1.7),
)
add_source_footer(sl, "All sources are from publicly available presentations, filings, earnings calls, and blog posts. Links may change over time.")


# ── SLIDE 38: Disclaimer / Back Cover ──────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_CREAM)
add_black_bars(sl)
txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AI Agent Slides"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_BLACK
p.font.name = "Arial"

p2 = tf.add_paragraph()
p2.text = "From Company Presentations, Filings & Webinars"
p2.font.size = Pt(22)
p2.font.color.rgb = TEXT_GRAY
p2.font.name = "Arial"
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = (
    "This document compiles publicly available information from company "
    "presentations, earnings calls, blog posts, and conference materials. "
    "All trademarks and logos belong to their respective owners. "
    "Content is for informational purposes only."
)
p3.font.size = Pt(12)
p3.font.color.rgb = TEXT_MUTED
p3.font.name = "Arial"
p3.space_before = Pt(30)

p4 = tf.add_paragraph()
p4.text = "February 2025"
p4.font.size = Pt(14)
p4.font.color.rgb = TEXT_MUTED
p4.font.name = "Arial"
p4.space_before = Pt(16)


# ── SAVE ────────────────────────────────────────────────────────────────
output_path = "/workspace/AI_Agent_Slides_From_Company_Presentations.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
