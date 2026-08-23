#!/usr/bin/env python3
"""Generate 'The Agentic Shift' — a PowerPoint deck on the state of the
agentic AI market and the use of AI agents in finance (August 2026 data).

Usage:  python generate_presentation.py
Output: agentic_ai_market_finance.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK

# ---------------------------------------------------------------- palette --
NAVY = RGBColor(0x0B, 0x1E, 0x3D)        # deep navy (dark slides)
NAVY_2 = RGBColor(0x14, 0x2E, 0x5C)      # lighter navy (cards on dark)
TEAL = RGBColor(0x18, 0xA9, 0x99)        # primary accent
GOLD = RGBColor(0xE8, 0xA3, 0x3D)        # secondary accent
INK = RGBColor(0x1B, 0x2A, 0x41)         # body text on light
MUTED = RGBColor(0x5A, 0x6B, 0x85)       # secondary text on light
CLOUD = RGBColor(0xF4, 0xF6, 0xFA)       # light card fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_LINE = RGBColor(0xD8, 0xDF, 0xEA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)

FONT = "Calibri"
FONT_LIGHT = "Calibri Light"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers --

def new_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, rounded=False, radius=0.06,
             line_color=None, line_w=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs_or_text, size=14, color=INK, bold=False,
             font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, space_after=0, wrap=True):
    """Add a simple text box. `runs_or_text` may be a string or a list of
    (text, {overrides}) tuples for mixed formatting in one paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    if space_after:
        p.space_after = Pt(space_after)
    if isinstance(runs_or_text, str):
        runs_or_text = [(runs_or_text, {})]
    for text, over in runs_or_text:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        if over.get("italic"):
            f.italic = True
    return box


def add_bullets(slide, x, y, w, h, items, size=14, color=INK, bullet_color=TEAL,
                line_spacing=1.08, space_after=8, lead_bold=True):
    """Bulleted list. Each item is either a string or a (lead, rest) tuple —
    the lead is rendered bold."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        b = p.add_run()
        b.text = "\u25AA  "
        b.font.name = FONT
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        b.font.bold = True
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = lead_bold
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return box


PAGE_NO = [0]


def content_header(slide, kicker, title, dark=False):
    """Standard header for light content slides: kicker, title, accent rule."""
    txt = WHITE if dark else INK
    kick_color = GOLD if dark else TEAL
    add_text(slide, MARGIN, Inches(0.42), SLIDE_W - 2 * MARGIN, Inches(0.3),
             kicker.upper(), size=12, color=kick_color, bold=True)
    add_text(slide, MARGIN, Inches(0.72), SLIDE_W - 2 * MARGIN, Inches(0.75),
             title, size=30, color=txt, bold=True, font=FONT_LIGHT)
    add_rect(slide, MARGIN, Inches(1.42), Inches(0.85), Pt(3.2), GOLD)


def footer(slide, dark=False):
    PAGE_NO[0] += 1
    color = RGBColor(0x8C, 0x9B, 0xB5) if dark else MUTED
    add_text(slide, MARGIN, Inches(7.08), Inches(7), Inches(0.3),
             "The Agentic Shift  ·  Agentic AI Market & Finance  ·  August 2026",
             size=9, color=color)
    add_text(slide, SLIDE_W - Inches(1.2), Inches(7.08), Inches(0.6), Inches(0.3),
             str(PAGE_NO[0]), size=9, color=color, align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, value, label, value_color=TEAL, dark=False,
              value_size=30, label_size=11.5):
    fill = NAVY_2 if dark else CLOUD
    card = add_rect(slide, x, y, w, h, fill, rounded=True, radius=0.09)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.08)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r = p1.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(value_size)
    r.font.bold = True
    r.font.color.rgb = value_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    p2.line_spacing = 1.0
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = FONT
    r2.font.size = Pt(label_size)
    r2.font.color.rgb = RGBColor(0xC5, 0xD0, 0xE0) if dark else MUTED
    return card


def info_card(slide, x, y, w, h, heading, body, heading_color=NAVY,
              body_size=11.5, heading_size=14, fill=CLOUD, body_color=None,
              accent=True):
    card = add_rect(slide, x, y, w, h, fill, rounded=True, radius=0.07)
    if accent:
        add_rect(slide, x, y + Inches(0.14), Inches(0.055), h - Inches(0.28), TEAL)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.1)
    p1 = tf.paragraphs[0]
    r = p1.add_run()
    r.text = heading
    r.font.name = FONT
    r.font.size = Pt(heading_size)
    r.font.bold = True
    r.font.color.rgb = heading_color
    p2 = tf.add_paragraph()
    p2.space_before = Pt(5)
    p2.line_spacing = 1.08
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = FONT
    r2.font.size = Pt(body_size)
    r2.font.color.rgb = body_color or MUTED if body_color is None else body_color
    return card


# ================================================================ SLIDE 1 ==
# Title
s = new_slide()
fill_bg(s, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(5), TEAL)
add_rect(s, Inches(0), SLIDE_H - Pt(5), SLIDE_W, Pt(5), GOLD)
# oversized watermark glyph
add_text(s, Inches(8.6), Inches(0.7), Inches(4.6), Inches(6),
         "AI", size=340, color=NAVY_2, bold=True, font=FONT_LIGHT, wrap=False)

add_text(s, MARGIN, Inches(1.7), Inches(10.5), Inches(0.4),
         "MARKET BRIEFING  ·  AUGUST 2026", size=14, color=GOLD, bold=True)
add_text(s, MARGIN, Inches(2.25), Inches(11.6), Inches(1.4),
         "The Agentic Shift", size=66, color=WHITE, bold=True, font=FONT_LIGHT)
add_text(s, MARGIN, Inches(3.55), Inches(11.2), Inches(1.1),
         "State of the agentic AI market — and how AI agents are\nrewiring financial services",
         size=22, color=RGBColor(0xC5, 0xD0, 0xE0), line_spacing=1.15)
add_rect(s, MARGIN, Inches(4.85), Inches(1.1), Pt(3.5), TEAL)
add_text(s, MARGIN, Inches(5.15), Inches(11), Inches(0.4),
         "From copilots to autonomous execution: sizing the market, mapping adoption, and locating real ROI in finance",
         size=13, color=RGBColor(0x8C, 0x9B, 0xB5))
footer(s, dark=True)

# ================================================================ SLIDE 2 ==
# Executive summary
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "Executive summary", "Five things to know in one slide")

cols = [
    ("$19B \u2192 $206B", "Global agentic AI market: $19.3B in 2026, projected "
     "$205.9B by 2033 (40.2% CAGR). Gartner counts $206.5B of 2026 spend on "
     "purpose-built agent software — up 139% YoY."),
    ("83% vs 23%", "The adoption gap is the story: 83% of enterprises plan to "
     "deploy agents, but only ~23% are scaling them in even one function. "
     "Pilots stall on integration and governance, not on models."),
    ("52%", "Financial services is the leading vertical: 52% of institutions "
     "are actively adopting agentic AI, and BFSI captures ~72% of specialized "
     "agent-platform spend. Fintechs lead incumbents 57% to 45%."),
    ("30\u201360%", "Where agents reach production in finance, returns are "
     "measurable: 40\u201370% fewer false positives, 50%+ faster SAR cycles, "
     "30\u201360% analyst productivity gains, payback in 8\u201318 months."),
    ("Aug 2026", "Governance is now the gating factor: the EU AI Act's "
     "high-risk regime covers credit scoring and fraud detection this month, "
     "with fines up to \u20AC35M or 7% of global turnover."),
]
cw = Inches(2.32)
gap = Inches(0.12)
x0 = MARGIN
for i, (big, body) in enumerate(cols):
    x = x0 + i * (cw + gap)
    card = add_rect(s, x, Inches(1.75), cw, Inches(4.9), CLOUD, rounded=True,
                    radius=0.05)
    add_rect(s, x, Inches(1.75), cw, Pt(4), TEAL if i % 2 == 0 else GOLD)
    add_text(s, x + Inches(0.16), Inches(2.0), cw - Inches(0.32), Inches(0.8),
             big, size=23, color=NAVY, bold=True)
    add_text(s, x + Inches(0.16), Inches(2.75), cw - Inches(0.32), Inches(3.7),
             body, size=11.5, color=INK, line_spacing=1.14)

add_text(s, MARGIN, Inches(6.72), SLIDE_W - 2 * MARGIN, Inches(0.35),
         [("Bottom line:  ", {"bold": True, "color": NAVY}),
          ("the market has moved from \u201Cwhat is an agent?\u201D to \u201Cwho can run one under audit?\u201D "
           "— and finance is running that test first.", {"color": MUTED})],
         size=12.5)
footer(s)

# ================================================================ SLIDE 3 ==
# What is agentic AI
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "Definitions", "From automation to autonomy")

add_text(s, MARGIN, Inches(1.65), Inches(12.1), Inches(0.6),
         "Agentic AI systems plan, reason, use tools, and act across multi-step workflows with limited human "
         "instruction — a step change from rule-based automation and prompt-response copilots.",
         size=14.5, color=MUTED, line_spacing=1.15)

stages = [
    ("2000s", "RPA & rules", "Scripted clicks and macros. Fast but brittle; "
     "breaks when the process changes."),
    ("2016+", "ML & chatbots", "Predictions and FAQ-style answers. Insight, "
     "but a human still executes."),
    ("2023+", "GenAI copilots", "Drafts, summaries, code suggestions. Human "
     "prompts every step."),
    ("2025+", "AI agents", "Goal in, outcome out: plans tasks, calls tools "
     "and APIs, acts under policy."),
    ("2026\u21922030", "Multi-agent systems", "Specialized agents orchestrated "
     "end-to-end; humans supervise by exception."),
]
cw = Inches(2.32)
gap = Inches(0.12)
for i, (era, name, body) in enumerate(stages):
    x = MARGIN + i * (cw + gap)
    dark_card = i >= 3
    fill = NAVY if dark_card else CLOUD
    add_rect(s, x, Inches(2.55), cw, Inches(2.5), fill, rounded=True, radius=0.06)
    add_text(s, x + Inches(0.16), Inches(2.72), cw - Inches(0.32), Inches(0.3),
             era, size=11, color=GOLD if dark_card else TEAL, bold=True)
    add_text(s, x + Inches(0.16), Inches(3.02), cw - Inches(0.32), Inches(0.4),
             name, size=15.5, color=WHITE if dark_card else NAVY, bold=True)
    add_text(s, x + Inches(0.16), Inches(3.5), cw - Inches(0.32), Inches(1.45),
             body, size=11, color=RGBColor(0xC5, 0xD0, 0xE0) if dark_card else MUTED,
             line_spacing=1.12)
    if i < len(stages) - 1:
        add_text(s, x + cw - Inches(0.06), Inches(3.55), Inches(0.3), Inches(0.4),
                 "\u203A", size=22, color=GOLD, bold=True, wrap=False)

add_text(s, MARGIN, Inches(5.35), Inches(12.1), Inches(0.35),
         "What makes an agent different", size=15, color=NAVY, bold=True)
add_bullets(s, MARGIN, Inches(5.75), Inches(12.1), Inches(1.3), [
    ("Autonomy with guardrails —", " pursues a goal across many steps, escalating to humans only at defined checkpoints."),
    ("Tool use —", " reads systems of record, calls APIs, writes back results; the agent executes work rather than describing it."),
    ("Memory and adaptation —", " retains context across a workflow and adjusts to new information mid-task."),
], size=12.5, space_after=5)
footer(s)

# ================================================================ SLIDE 4 ==
# Market sizing
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "State of the market", "Market sizing: it depends who's counting")

add_text(s, MARGIN, Inches(1.62), Inches(12.1), Inches(0.55),
         "Analyst estimates for 2026 span $9B to $206B because each firm draws the perimeter differently — "
         "pure-play platforms vs. all agent software embedded in the enterprise stack.",
         size=13.5, color=MUTED, line_spacing=1.12)

rows = [
    ("Gartner", "All purpose-built AI agent software", "$206.5B in 2026 (+139% YoY from $86.4B)",
     "~62.7% CAGR \u2192 ~$985B by 2030"),
    ("MarketsandMarkets", "Agentic platforms, orchestration & services", "$19.3B in 2026",
     "40.2% CAGR \u2192 $205.9B by 2033"),
    ("Fortune Business Insights", "Broad pure-play agentic AI market", "$9.1B in 2026",
     "40.5% CAGR \u2192 $139.2B by 2034"),
    ("Grand View Research", "Enterprise-specific agentic AI", "$5.4B in 2026",
     "46.2% CAGR \u2192 $24.5B by 2030"),
]
ty = Inches(2.35)
row_h = Inches(0.78)
col_x = [MARGIN, Inches(3.3), Inches(6.55), Inches(10.0)]
col_w = [Inches(2.55), Inches(3.1), Inches(3.3), Inches(2.7)]
hdr = ["Source", "Scope measured", "2026 market size", "Trajectory"]
add_rect(s, MARGIN, ty, SLIDE_W - 2 * MARGIN, Inches(0.42), NAVY, rounded=True, radius=0.12)
for j, htxt in enumerate(hdr):
    add_text(s, col_x[j] + Inches(0.14), ty + Inches(0.07), col_w[j], Inches(0.3),
             htxt, size=12, color=WHITE, bold=True)
y = ty + Inches(0.5)
for i, (src, scope, size_2026, traj) in enumerate(rows):
    if i % 2 == 0:
        add_rect(s, MARGIN, y, SLIDE_W - 2 * MARGIN, row_h, CLOUD)
    add_text(s, col_x[0] + Inches(0.14), y + Inches(0.12), col_w[0], Inches(0.55),
             src, size=12.5, color=NAVY, bold=True)
    add_text(s, col_x[1] + Inches(0.14), y + Inches(0.12), col_w[1], Inches(0.6),
             scope, size=11.5, color=INK, line_spacing=1.05)
    add_text(s, col_x[2] + Inches(0.14), y + Inches(0.12), col_w[2], Inches(0.6),
             size_2026, size=12, color=TEAL, bold=True, line_spacing=1.05)
    add_text(s, col_x[3] + Inches(0.14), y + Inches(0.12), col_w[3], Inches(0.6),
             traj, size=11.5, color=MUTED, line_spacing=1.05)
    y += row_h

add_text(s, MARGIN, Inches(6.35), SLIDE_W - 2 * MARGIN, Inches(0.55),
         [("The consensus that matters:  ", {"bold": True, "color": NAVY}),
          ("whatever the perimeter, every major forecast lands in a 40\u201355% CAGR band through 2030 — "
           "the fastest-growing category inside a $2.59T global AI spend.", {"color": MUTED})],
         size=13, line_spacing=1.12)
footer(s)

# ================================================================ SLIDE 5 ==
# Growth chart
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "State of the market", "Pilots are converting to production spend")

chart_data = CategoryChartData()
chart_data.categories = ["2024", "2025", "2026", "2027", "2028", "2029", "2030"]
chart_data.add_series("Enterprise agentic AI market ($B)",
                      (2.58, 3.67, 5.37, 7.85, 11.48, 16.79, 24.50))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            MARGIN, Inches(1.85), Inches(7.9), Inches(4.7),
                            chart_data)
chart = gframe.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 68
series = plot.series[0]
for i, pt in enumerate(series.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = GOLD if i == 2 else TEAL
plot.has_data_labels = True
dl = plot.data_labels
dl.number_format = '"$"0.0"B"'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(11)
dl.font.bold = True
dl.font.color.rgb = NAVY
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(12)
cat_ax.tick_labels.font.color.rgb = INK
cat_ax.major_tick_mark = XL_TICK_MARK.NONE
cat_ax.format.line.color.rgb = SOFT_LINE
val_ax = chart.value_axis
val_ax.visible = False
val_ax.has_major_gridlines = False

add_text(s, MARGIN, Inches(6.55), Inches(7.9), Inches(0.35),
         "Enterprise-specific agentic AI revenue, Grand View Research (2026 est. onward). 2026 highlighted.",
         size=10, color=MUTED)

rx = Inches(8.95)
rw = Inches(3.75)
stat_card(s, rx, Inches(1.85), rw, Inches(1.28), "46%",
          "compound annual growth 2025\u20132030 — pilots become core infrastructure by 2027")
stat_card(s, rx, Inches(3.25), rw, Inches(1.28), "$2.59T",
          "total worldwide AI spend forecast for 2026 (+47% YoY); agent software is its fastest-growing slice",
          value_color=GOLD)
stat_card(s, rx, Inches(4.65), rw, Inches(1.28), "72%",
          "of the 2026 agentic market is software — platforms, orchestration and governance layers")
add_text(s, rx, Inches(6.1), rw, Inches(0.7),
         "Services grow even faster than software: once platforms land, spend shifts to integration, "
         "governance and workflow design.", size=10.5, color=MUTED, line_spacing=1.12)
footer(s)

# ================================================================ SLIDE 6 ==
# Adoption reality gap
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "State of the market", "The adoption reality gap: intent far outruns production")

bar_data = CategoryChartData()
bar_data.categories = [
    "Agentic projects predicted\nto be scrapped by 2027",
    "Scaling agents in \u22651\nbusiness function",
    "Actively piloting today",
    "Plan to deploy agents",
]
bar_data.add_series("Share of enterprises (%)", (40, 23, 52, 83))
gframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            MARGIN, Inches(1.8), Inches(6.9), Inches(4.4),
                            bar_data)
chart = gframe.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 60
series = plot.series[0]
colors = [RGBColor(0xC0, 0x53, 0x4F), TEAL, TEAL, NAVY]
for i, pt in enumerate(series.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = colors[i]
plot.has_data_labels = True
dl = plot.data_labels
dl.number_format = '0"%"'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(12)
dl.font.bold = True
dl.font.color.rgb = NAVY
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(10.5)
cat_ax.tick_labels.font.color.rgb = INK
cat_ax.major_tick_mark = XL_TICK_MARK.NONE
cat_ax.format.line.color.rgb = SOFT_LINE
val_ax = chart.value_axis
val_ax.visible = False
val_ax.has_major_gridlines = False
val_ax.maximum_scale = 100

add_text(s, MARGIN, Inches(6.3), Inches(6.9), Inches(0.5),
         "Sources: Cisco (Oct 2025), McKinsey State of AI (2025), Gartner (2025\u201326).",
         size=10, color=MUTED)

rx = Inches(7.95)
rw = Inches(4.75)
info_card(s, rx, Inches(1.8), rw, Inches(1.45), "Only ~1 in 4 converts",
          "The deployment-reality index (scaling \u00F7 planning) sits at ~28: "
          "for every four enterprises that intend to run agents, one has them in scaled production.")
info_card(s, rx, Inches(3.4), rw, Inches(1.45), "Why pilots stall",
          "The bottleneck is rarely the model. Legacy integration, data controls, "
          "security reviews and unclear ownership of agent decisions stop pilots at the gate.")
info_card(s, rx, Inches(5.0), rw, Inches(1.45), "The shake-out is priced in",
          "Gartner expects over 40% of agentic projects to be scrapped by 2027 on cost and unclear ROI — "
          "spend keeps growing anyway, concentrated on survivors with governed workflows.")
footer(s)

# ================================================================ SLIDE 7 ==
# Competitive landscape
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "State of the market", "Landscape: platforms, standards, and capital")

vendors = [
    ("Microsoft", "Copilot Studio embedded across the enterprise estate; agents woven into M365, Dynamics, GitHub."),
    ("Salesforce", "Agentforce: 29,000+ deals and ~$800M ARR by Q4 FY26; ARR up 330% YoY — fastest enterprise ramp."),
    ("Google", "Vertex AI + Agentspace; launched the A2A interoperability protocol with 50+ partners, donated to Linux Foundation."),
    ("AWS", "Bedrock ecosystem with sandboxed agent compute; the default infrastructure play."),
    ("OpenAI / Anthropic", "Frontier models and APIs powering most agent stacks; co-creators of open agent standards (AGENTS.md, MCP)."),
    ("ServiceNow", "ITSM-native agents plus a cross-vendor orchestration platform to unify agents from multiple providers."),
]
cw2 = Inches(4.0)
ch2 = Inches(1.32)
gx = Inches(0.14)
gy = Inches(0.14)
for i, (name, body) in enumerate(vendors):
    col = i % 3
    row = i // 3
    x = MARGIN + col * (cw2 + gx)
    y = Inches(1.72) + row * (ch2 + gy)
    info_card(s, x, y, cw2, ch2, name, body, heading_size=13, body_size=10.5)

add_rect(s, MARGIN, Inches(4.85), SLIDE_W - 2 * MARGIN, Pt(1.2), SOFT_LINE)

add_text(s, MARGIN, Inches(5.02), Inches(6), Inches(0.3),
         "Standards & business model", size=14, color=NAVY, bold=True)
add_bullets(s, MARGIN, Inches(5.4), Inches(6.4), Inches(1.5), [
    ("Interoperability is arriving: ", "MCP for tool access, A2A for agent-to-agent, AGENTS.md guiding 60k+ projects."),
    ("Pricing is shifting ", "from seats to outcomes — customers increasingly pay per completed task."),
], size=11.5, space_after=5)

add_text(s, Inches(7.35), Inches(5.02), Inches(5.3), Inches(0.3),
         "Where the capital goes", size=14, color=NAVY, bold=True)
add_bullets(s, Inches(7.35), Inches(5.4), Inches(5.35), Inches(1.5), [
    ("21% of the CB Insights AI 100 ", "(2025) are agent or agent-infrastructure companies."),
    ("Horizontal agents raised $1.6B; ", "infrastructure and vertical players ~$1.2B each — investors fund the integration layer where pilots stall."),
], size=11.5, space_after=5)
footer(s)

# ================================================================ SLIDE 8 ==
# Section divider — finance
s = new_slide()
fill_bg(s, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Pt(5), GOLD)
add_text(s, Inches(9.0), Inches(1.2), Inches(4.5), Inches(5.5),
         "02", size=300, color=NAVY_2, bold=True, font=FONT_LIGHT, wrap=False)
add_text(s, MARGIN, Inches(2.5), Inches(9), Inches(0.4),
         "PART II", size=14, color=GOLD, bold=True)
add_text(s, MARGIN, Inches(3.0), Inches(10), Inches(1.2),
         "AI agents in finance", size=54, color=WHITE, bold=True, font=FONT_LIGHT)
add_rect(s, MARGIN, Inches(4.15), Inches(1.1), Pt(3.5), TEAL)
add_text(s, MARGIN, Inches(4.45), Inches(9.5), Inches(1.2),
         "The most regulated industry is also the fastest adopter — because "
         "compliance forced the data discipline agents need to work.",
         size=18, color=RGBColor(0xC5, 0xD0, 0xE0), line_spacing=1.2)
footer(s, dark=True)

# ================================================================ SLIDE 9 ==
# Why finance leads
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "AI agents in finance", "Why finance leads every other vertical")

stat_card(s, MARGIN, Inches(1.75), Inches(3.9), Inches(1.5), "#1",
          "BFSI is the largest agentic AI vertical — ~19% of enterprise spend "
          "and ~72% of specialized agent-platform spend")
stat_card(s, Inches(4.72), Inches(1.75), Inches(3.9), Inches(1.5), "$1.03B",
          "BFSI agentic AI spend in 2026 — doubled from $493M in 2024",
          value_color=GOLD)
stat_card(s, Inches(8.82), Inches(1.75), Inches(3.9), Inches(1.5), "66%",
          "of corporate finance functions expected to run meaningful agentic AI "
          "by 2030, up from 20% in 2025")

add_text(s, MARGIN, Inches(3.55), Inches(12), Inches(0.35),
         "Four structural reasons", size=15, color=NAVY, bold=True)

reasons = [
    ("Data discipline was already paid for", "Decades of GRC requirements forced clean, governed, auditable data — "
     "exactly what agents need to act reliably. The constraint that slowed cloud adoption now accelerates agent ROI."),
    ("Workflows are agent-shaped", "High-volume, rules-bounded, multi-step processes (KYC reviews, claims, "
     "reconciliations, alert triage) are the ideal first workloads for autonomous execution."),
    ("Cost pressure meets digital rails", "Persistent margin pressure plus mature digital infrastructure means the "
     "business case writes itself where headcount scales linearly with volume."),
    ("Fraud is an arms race", "~50% of fraud already involves AI and deepfake scams are up 2,000%+ in three years — "
     "defense at machine speed requires agents, not queues of analysts."),
]
cw3 = Inches(6.0)
ch3 = Inches(1.35)
for i, (h, b) in enumerate(reasons):
    x = MARGIN + (i % 2) * (cw3 + Inches(0.14))
    y = Inches(3.95) + (i // 2) * (ch3 + Inches(0.14))
    info_card(s, x, y, cw3, ch3, h, b, heading_size=13, body_size=10.8)
footer(s)

# =============================================================== SLIDE 10 ==
# Adoption in financial services
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "AI agents in finance", "Adoption: half in motion, a quarter scaling")

bar_data = CategoryChartData()
bar_data.categories = [
    "Scaling or transforming today",
    "Traditional FIs actively adopting",
    "Fintechs actively adopting",
    "All FIs actively adopting",
    "Expect meaningful deployment by 2030",
]
bar_data.add_series("Share of institutions (%)", (23, 45, 57, 52, 81))
gframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            MARGIN, Inches(1.8), Inches(7.3), Inches(4.4),
                            bar_data)
chart = gframe.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 55
series = plot.series[0]
colors = [MUTED, NAVY, TEAL, TEAL, GOLD]
for i, pt in enumerate(series.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = colors[i]
plot.has_data_labels = True
dl = plot.data_labels
dl.number_format = '0"%"'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(12)
dl.font.bold = True
dl.font.color.rgb = NAVY
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(10.5)
cat_ax.tick_labels.font.color.rgb = INK
cat_ax.major_tick_mark = XL_TICK_MARK.NONE
cat_ax.format.line.color.rgb = SOFT_LINE
val_ax = chart.value_axis
val_ax.visible = False
val_ax.has_major_gridlines = False
val_ax.maximum_scale = 100

add_text(s, MARGIN, Inches(6.3), Inches(7.3), Inches(0.4),
         "Source: Cambridge Centre for Alternative Finance, 2026 Global AI in Financial Services Report "
         "(survey of industry firms and 130 regulators).",
         size=9.5, color=MUTED, line_spacing=1.1)

rx = Inches(8.3)
rw = Inches(4.4)
info_card(s, rx, Inches(1.8), rw, Inches(1.42), "Fastest uptake on record",
          "Agentic AI reached 52% active adoption within ~2 years of category emergence — "
          "faster than classical ML or GenAI at the same age.")
info_card(s, rx, Inches(3.36), rw, Inches(1.42), "Fintechs set the pace",
          "57% vs 45% for incumbents — and 19% vs 6% at the fully transforming stage. "
          "Legacy infrastructure, not ambition, is the drag.")
info_card(s, rx, Inches(4.92), rw, Inches(1.42), "Banking's execution gap",
          "Only 16% of banking executives report use cases actually deployed; "
          "52% remain in piloting. Regulators trail further — half are still exploring.")
footer(s)

# =============================================================== SLIDE 11 ==
# Use cases
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "AI agents in finance", "Where agents work today: six proven workloads")

cases = [
    ("Fraud & financial crime", "Agents investigate alerts end-to-end: pull transaction history, check device "
     "fingerprints, cross-reference watchlists, and act under policy — at machine speed.", TEAL),
    ("AML / KYC compliance", "Autonomous evidence gathering, sanctions & PEP screening, risk profiling and "
     "SAR drafting, with a human validating before filing.", NAVY),
    ("Credit & underwriting", "Agents assemble borrower files, verify documents, run policy-bound risk "
     "assessments and route exceptions — compressing days into hours.", TEAL),
    ("Wealth & advisory", "Meeting prep, portfolio review drafts, personalized agendas and client-ready "
     "summaries — one firm saves ~20,000 advisor hours a year.", NAVY),
    ("Treasury & finance ops", "Cash forecasting, reconciliations, payment-instruction validation and "
     "close automation inside the CFO stack.", TEAL),
    ("Agentic payments & commerce", "Agents that transact on the customer's behalf — Mastercard Agent Pay, "
     "Visa Intelligent Commerce, PayPal's agent toolkit build the rails.", NAVY),
]
cw4 = Inches(4.0)
ch4 = Inches(2.28)
for i, (h, b, acc) in enumerate(cases):
    col = i % 3
    row = i // 3
    x = MARGIN + col * (cw4 + Inches(0.14))
    y = Inches(1.78) + row * (ch4 + Inches(0.16))
    card = add_rect(s, x, y, cw4, ch4, CLOUD, rounded=True, radius=0.06)
    add_rect(s, x, y, cw4, Pt(4), acc)
    add_text(s, x + Inches(0.2), y + Inches(0.18), cw4 - Inches(0.4), Inches(0.55),
             h, size=14.5, color=NAVY, bold=True, line_spacing=1.0)
    add_text(s, x + Inches(0.2), y + Inches(0.75), cw4 - Inches(0.4), Inches(1.45),
             b, size=11, color=INK, line_spacing=1.14)

add_text(s, MARGIN, Inches(6.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
         [("Pattern:  ", {"bold": True, "color": NAVY}),
          ("the first production wins are back- and middle-office — high volume, clear policy, "
           "measurable unit cost.", {"color": MUTED})],
         size=12.5)
footer(s)

# =============================================================== SLIDE 12 ==
# Measured impact
s = new_slide()
fill_bg(s, NAVY)
content_header(s, "AI agents in finance", "Measured impact where agents reached production", dark=True)

stats = [
    ("\u221240\u201370%", "false positives in fraud & AML screening", TEAL),
    ("\u221250%+", "SAR filing cycle time in compliance", GOLD),
    ("+30\u201360%", "analyst & investigator productivity", TEAL),
    ("\u221230\u201350%", "manual workload in KYC / AML reviews (McKinsey)", GOLD),
    ("8\u201318 mo", "typical payback period on deployment", TEAL),
    ("100%", "policy adherence in agent decisioning vs. <95% for human four-eyes review (Sardine)", GOLD),
]
cw5 = Inches(4.0)
ch5 = Inches(1.62)
for i, (v, l, c) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = MARGIN + col * (cw5 + Inches(0.14))
    y = Inches(1.8) + row * (ch5 + Inches(0.16))
    stat_card(s, x, y, cw5, ch5, v, l, value_color=c, dark=True, value_size=27,
              label_size=11)

add_rect(s, MARGIN, Inches(5.5), SLIDE_W - 2 * MARGIN, Inches(1.35), NAVY_2,
         rounded=True, radius=0.08)
add_text(s, MARGIN + Inches(0.25), Inches(5.68), Inches(2.6), Inches(0.9),
         "Case in point:\nKlarna", size=16, color=GOLD, bold=True, line_spacing=1.05)
add_text(s, MARGIN + Inches(2.9), Inches(5.66), Inches(9.2), Inches(1.1),
         "Its OpenAI-powered assistant handled 2.3M conversations in month one — two-thirds of all "
         "service chats, the workload of ~700 full-time agents — cutting resolution time from 11 minutes "
         "to under 2 and adding an estimated $40M in annual profit improvement.",
         size=12.5, color=RGBColor(0xC5, 0xD0, 0xE0), line_spacing=1.18)
footer(s, dark=True)

# =============================================================== SLIDE 13 ==
# Real deployments
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "AI agents in finance", "Who's already doing it: named deployments")

deps = [
    ("JPMorgan Chase", "LAW — Legal Agentic Workflows: specialized agents processing custody and fund-services "
     "contracts at 92.9% accuracy across query types."),
    ("BNY", "Autonomous agents working in coding and payment-instruction validation — agents as digital "
     "employees with supervised authority."),
    ("Citi", "Framed the \u201Cdo-it-for-me economy\u201D research agenda; piloting agent assistants "
     "(e.g. voice-enabled market insights) across wealth and markets."),
    ("Mastercard", "Agent Pay: agentic payments rails so consumer and business agents can transact — built with "
     "IBM and Microsoft partnerships."),
    ("PayPal & Visa", "Agent toolkits and \u201CIntelligent Commerce\u201D APIs enabling agents to buy securely "
     "on behalf of customers."),
    ("Klarna", "Consumer-scale proof point: AI assistant absorbing two-thirds of customer service volume "
     "with measured profit impact."),
]
cw6 = Inches(6.0)
ch6 = Inches(1.42)
for i, (h, b) in enumerate(deps):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (cw6 + Inches(0.14))
    y = Inches(1.78) + row * (ch6 + Inches(0.16))
    info_card(s, x, y, cw6, ch6, h, b, heading_size=13.5, body_size=11)

add_text(s, MARGIN, Inches(6.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
         [("Signal:  ", {"bold": True, "color": NAVY}),
          ("the most-regulated institutions are past experimentation — naming systems, "
           "publishing accuracy, and building payment rails.",
           {"color": MUTED})],
         size=12)
footer(s)

# =============================================================== SLIDE 14 ==
# Risk & governance
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "AI agents in finance", "The gating factor: governance, trust and the EU AI Act")

add_rect(s, MARGIN, Inches(1.75), Inches(5.9), Inches(4.75), NAVY, rounded=True, radius=0.05)
add_text(s, MARGIN + Inches(0.3), Inches(2.0), Inches(5.3), Inches(0.35),
         "REGULATORY CLOCK", size=12, color=GOLD, bold=True)
add_text(s, MARGIN + Inches(0.3), Inches(2.38), Inches(5.3), Inches(0.8),
         "EU AI Act high-risk regime — August 2026", size=19, color=WHITE, bold=True,
         line_spacing=1.05)
add_bullets(s, MARGIN + Inches(0.3), Inches(3.25), Inches(5.35), Inches(3.0), [
    ("Credit scoring and fraud detection ", "are explicitly classified high-risk AI systems."),
    ("Institutions must demonstrate ", "explainability, human oversight, audit readiness and documented risk management."),
    ("Penalties reach ", "\u20AC35M or 7% of global annual turnover."),
    ("Only ~15% of CFOs ", "say they are ready to deploy and govern agents at scale today."),
], size=12, color=RGBColor(0xE6, 0xEB, 0xF3), bullet_color=GOLD, space_after=8)

rx = Inches(7.0)
rw = Inches(5.7)
add_text(s, rx, Inches(1.78), rw, Inches(0.35),
         "What regulators and boards now expect", size=15, color=NAVY, bold=True)
add_bullets(s, rx, Inches(2.25), rw, Inches(3.6), [
    ("Human-in-the-loop by design — ", "autonomy graduated by risk tier; high-impact decisions gated on human validation."),
    ("Immutable audit trails — ", "every agent action logged, replayable and attributable; explainable reasoning for any decision that materially affects a customer."),
    ("Agent identity & entitlements — ", "agents as first-class identities with least-privilege access, spend limits and revocable credentials."),
    ("Adversarial resilience — ", "~50% of fraud already involves AI; deepfakes up 2,000%+ in three years. Agents are both the defense and a new attack surface."),
], size=12.5, space_after=9)
add_text(s, rx, Inches(6.15), rw, Inches(0.6),
         [("Read:  ", {"bold": True, "color": NAVY}),
          ("governance is not a tax on agentic AI in finance — it is the moat. Firms that built it first "
           "are the ones scaling now.", {"color": MUTED})], size=12, line_spacing=1.12)
footer(s)

# =============================================================== SLIDE 15 ==
# Outlook & recommendations
s = new_slide()
fill_bg(s, WHITE)
content_header(s, "Outlook", "2026\u20132030: what happens next, and what to do about it")

add_text(s, MARGIN, Inches(1.7), Inches(5.9), Inches(0.35),
         "Where the market goes", size=15, color=NAVY, bold=True)
add_bullets(s, MARGIN, Inches(2.15), Inches(5.9), Inches(4.3), [
    ("Multi-agent orchestration ", "becomes the default architecture: specialized agents chained end-to-end, humans supervising by exception."),
    ("Agentic commerce scales — ", "agents intermediating B2B and consumer transactions on new payment rails from 2027\u201328."),
    ("Outcome-based pricing ", "displaces seats; ROI becomes measurable per completed task, tightening the link between spend and value."),
    ("Consolidation: ", ">40% of today's projects fail, and spend concentrates on platforms with governance built in."),
    ("81% of financial institutions ", "expect meaningful agentic deployment by 2030 — the frontier moves from whether to how fast."),
], size=12.5, space_after=9)

rx = Inches(7.0)
add_text(s, rx, Inches(1.7), Inches(5.7), Inches(0.35),
         "A pragmatic playbook for finance leaders", size=15, color=NAVY, bold=True)
steps = [
    ("1", "Pick governed, high-volume workflows first", "Alert triage, KYC refresh, reconciliations — clear policy, measurable unit cost, contained blast radius."),
    ("2", "Build the audit layer before the agent", "Logging, HITL checkpoints and entitlements up front make every later deployment faster to approve."),
    ("3", "Measure at task level; scale what pays back", "Target 8\u201318-month payback. Kill pilots that can't show task-level ROI — most of the 40% failure rate is avoidable selection error."),
    ("4", "Stay portable via open standards", "MCP, A2A and AGENTS.md reduce switching costs as the vendor landscape consolidates."),
]
y = Inches(2.15)
for num, h, b in steps:
    add_rect(s, rx, y, Inches(0.42), Inches(0.42), TEAL, rounded=True, radius=0.5)
    add_text(s, rx, y + Inches(0.04), Inches(0.42), Inches(0.35), num, size=16,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, rx + Inches(0.6), y - Inches(0.02), Inches(5.1), Inches(0.35),
             h, size=13, color=NAVY, bold=True)
    add_text(s, rx + Inches(0.6), y + Inches(0.32), Inches(5.1), Inches(0.75),
             b, size=10.8, color=MUTED, line_spacing=1.1)
    y += Inches(1.14)
footer(s)

# =============================================================== SLIDE 16 ==
# Sources
s = new_slide()
fill_bg(s, NAVY)
content_header(s, "Appendix", "Sources & notes", dark=True)

sources_left = [
    "Gartner, AI Spending Forecast 1Q26 (May 2026) — agent software & total AI spend",
    "MarketsandMarkets, Agentic AI Market Report 2026\u20132033 (Aug 2026)",
    "Grand View Research / Keyhole Software, Enterprise Agentic AI Market (Q2 2026)",
    "Fortune Business Insights, Agentic AI Market 2026\u20132034",
    "Cambridge Centre for Alternative Finance & industry partners, 2026 Global AI in Financial Services Report (Apr 2026)",
    "McKinsey, State of AI (2025) & AI Trust Maturity Survey (2026)",
]
sources_right = [
    "Deloitte Insights, Agentic AI in Banking (2025)",
    "Citi GPS, Agentic AI: Finance & the \u201CDo It For Me\u201D Economy (Jan 2025)",
    "SymphonyAI, FinCrime Frontier Survey 2025\u201326",
    "CB Insights, AI 100 (2025); PitchBook, Agentic AI Analyst Notes (Q2 2026)",
    "Cisco (Oct 2025) & Axis Intelligence, Agentic AI Statistics 2026",
    "Company disclosures: Salesforce FY26 earnings; JPMorgan, BNY, Mastercard, PayPal, Visa, Klarna announcements",
]
add_bullets(s, MARGIN, Inches(1.85), Inches(5.9), Inches(4.3), sources_left,
            size=12, color=RGBColor(0xC5, 0xD0, 0xE0), bullet_color=GOLD,
            space_after=10)
add_bullets(s, Inches(7.0), Inches(1.85), Inches(5.7), Inches(4.3), sources_right,
            size=12, color=RGBColor(0xC5, 0xD0, 0xE0), bullet_color=GOLD,
            space_after=10)

add_rect(s, MARGIN, Inches(6.2), SLIDE_W - 2 * MARGIN, Pt(1), NAVY_2)
add_text(s, MARGIN, Inches(6.4), SLIDE_W - 2 * MARGIN, Inches(0.5),
         "Note: market-size estimates differ by scope; figures are cited to their source's perimeter. "
         "Vendor metrics are self-reported. All data as published through August 2026.",
         size=10.5, color=RGBColor(0x8C, 0x9B, 0xB5), line_spacing=1.15)
footer(s, dark=True)

# ----------------------------------------------------------------- output --
OUT = "agentic_ai_market_finance.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides")
