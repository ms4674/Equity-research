from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BG = RGBColor(0x1B, 0x1F, 0x2B)
ACCENT_ORANGE = RGBColor(0xFF, 0x6F, 0x00)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x3D)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)
ACCENT_PURPLE = RGBColor(0xA7, 0x55, 0xF7)
ACCENT_TEAL = RGBColor(0x14, 0xB8, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
CARD_BG = RGBColor(0x25, 0x2A, 0x3A)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line_info.get("text", "")
        p.text = text
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.alignment = line_info.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_info.get("space_after", 6))
    return txBox


def add_stat_card(slide, left, top, width, height, stat_value, stat_label, accent_color):
    card = add_card(slide, left, top, width, height)
    add_accent_bar(slide, left, top, width, Pt(4), accent_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.6),
                 stat_value, font_size=28, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), height - Inches(0.9),
                 stat_label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return card


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── SLIDE 1: TITLE ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_accent_bar(slide, Inches(1), Inches(2.8), Inches(1.5), Pt(5), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
             "DATABRICKS DATA + AI SUMMIT", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
             "AI Agents & Token Usage", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(1),
             "Enterprise Adoption, Cost Economics & Production Insights", font_size=22, color=LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "Consolidated from Data + AI Summit 2024, 2025 & 2026 Sessions", font_size=14, color=LIGHT_GRAY)

cards_y = Inches(5.0)
card_w = Inches(2.5)
card_h = Inches(1.6)
stats = [
    ("22,000+", "Summit 2025 Attendees", ACCENT_ORANGE),
    ("700+", "Sessions Presented", ACCENT_BLUE),
    ("20,000+", "Databricks Customers", ACCENT_GREEN),
    ("60%+", "Fortune 500 on Platform", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(stats):
    add_stat_card(slide, Inches(1) + i * (card_w + Inches(0.3)), cards_y, card_w, card_h, val, label, color)


# ── SLIDE 2: AGENDA ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Agenda", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.2), Pt(4), ACCENT_ORANGE)

agenda_items = [
    ("01", "Enterprise AI Agent Landscape", "Market trends, adoption statistics, and growth trajectory", ACCENT_ORANGE),
    ("02", "Key Summit Sessions on AI Agents", "Highlights from 2024, 2025, and 2026 summit presentations", ACCENT_BLUE),
    ("03", "Token Usage & Cost Economics", "Token consumption patterns, pricing, and optimization strategies", ACCENT_GREEN),
    ("04", "Company Case Studies", "Adidas, Walmart, AstraZeneca, Royal Bank of Canada, and more", ACCENT_PURPLE),
    ("05", "State of AI Agents Report", "Multi-agent growth, governance impact, and production readiness", ACCENT_TEAL),
    ("06", "Production Architecture & Best Practices", "Compound AI systems, monitoring, and deployment patterns", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(1.6) + i * Inches(0.9)
    add_accent_bar(slide, Inches(0.8), y, Pt(4), Inches(0.7), color)
    add_text_box(slide, Inches(1.1), y, Inches(0.6), Inches(0.7), num, font_size=24, color=color, bold=True)
    add_text_box(slide, Inches(1.8), y - Inches(0.02), Inches(5), Inches(0.4), title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.35), Inches(8), Inches(0.4), desc, font_size=14, color=LIGHT_GRAY)


# ── SLIDE 3: ENTERPRISE AI AGENT LANDSCAPE ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Enterprise AI Agent Landscape", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2), Pt(4), ACCENT_ORANGE)

stats_row1 = [
    ("$58B → $474B", "Enterprise AI Market\n2025 → 2030 (52% CAGR)", ACCENT_ORANGE),
    ("$644B", "GenAI Spend\nForecast 2025", ACCENT_BLUE),
    ("$1.5T", "Total AI Spend\nForecast 2025", ACCENT_GREEN),
    ("75%", "C-Suite Executives Rank\nAI as Top-3 Priority", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(stats_row1):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(1.6), Inches(2.7), Inches(1.6), val, label, color)

stats_row2 = [
    ("11x", "More AI Models\nDeployed YoY", ACCENT_TEAL),
    ("3x", "More Efficient at\nProduction Deployment", ACCENT_RED),
    ("377%", "Vector Database\nUsage Growth YoY", ACCENT_BLUE),
    ("76%", "Companies Using LLMs\nChoose Open Source", ACCENT_ORANGE),
]
for i, (val, label, color) in enumerate(stats_row2):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(3.6), Inches(2.7), Inches(1.6), val, label, color)

add_multi_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.5), [
    {"text": "Key Insight: Financial Services leads GPU usage growth at 88% over 6 months,", "size": 14, "color": ACCENT_ORANGE, "bold": True},
    {"text": "followed by Healthcare & Life Sciences — highly regulated industries are the surprise early adopters of GenAI.", "size": 14, "color": LIGHT_GRAY},
    {"text": "Source: Databricks State of Data + AI Report 2024 — 10,000+ global customers including 300+ Fortune 500 companies", "size": 11, "color": RGBColor(0x70, 0x78, 0x90)},
])


# ── SLIDE 4: KEY SUMMIT SESSIONS – AI AGENTS (2024) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Key Summit Sessions: AI Agents (2024)", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(4), ACCENT_ORANGE)

sessions_2024 = [
    ("Mosaic AI: Compound AI Systems",
     "Databricks introduced Mosaic AI for building production-quality compound AI systems. "
     "Emphasis shifted from monolithic models to multi-component systems using tuned models, "
     "retrieval, tool use, and reasoning agents.",
     ["New SDK for building, deploying & evaluating AI agents",
      "Foundation Model APIs with pay-per-token pricing",
      "Mosaic AI Tool Catalog for enterprise function registries",
      "Model Serving support for agents and RAG applications"],
     ACCENT_ORANGE),
    ("GenAI Announcements Keynote",
     "16,000+ in-person attendees and 60,000+ virtual participants from 140 countries. "
     "Focus on production-grade compound AI and fine-tuning support.",
     ["Fine-tuning support for foundation models",
      "70% of GenAI companies use vector databases for RAG",
      "Llama 3 captured 39% of open-source usage within 4 weeks of launch",
      "Hugging Face Transformers jumped from #4 to #2 most-adopted product"],
     ACCENT_BLUE),
]

for i, (title, desc, bullets, color) in enumerate(sessions_2024):
    y_start = Inches(1.5) + i * Inches(2.8)
    card = add_card(slide, Inches(0.8), y_start, Inches(11.5), Inches(2.5))
    add_accent_bar(slide, Inches(0.8), y_start, Pt(5), Inches(2.5), color)
    add_text_box(slide, Inches(1.2), y_start + Inches(0.15), Inches(10), Inches(0.5),
                 title, font_size=22, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y_start + Inches(0.6), Inches(5), Inches(1.8),
                 desc, font_size=13, color=LIGHT_GRAY)

    bullet_lines = []
    for b in bullets:
        bullet_lines.append({"text": f"•  {b}", "size": 13, "color": WHITE, "space_after": 4})
    add_multi_text(slide, Inches(6.5), y_start + Inches(0.55), Inches(5.5), Inches(2), bullet_lines)


# ── SLIDE 5: KEY SUMMIT SESSIONS – AI AGENTS (2025) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Key Summit Sessions: AI Agents (2025)", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "June 9–12, 2025 · San Francisco · 22,000+ Attendees · 700+ Sessions", font_size=13, color=LIGHT_GRAY)

sessions_2025 = [
    ("Autonomous AI Agents in AI Infrastructure", "Apurva Kumar — Walmart Global Tech",
     "Architecture & design principles for autonomous agents: containerization, orchestration, robustness, and real-time feedback loops.",
     ACCENT_ORANGE),
    ("Agent Bricks: Building Multi-Agent Systems", "Databricks",
     "Natural language interactions with structured databases and unstructured documents. KBQA evaluation, data governance, and vector search + structured DB retrieval.",
     ACCENT_BLUE),
    ("Multi-Agents in Production", "Tony Holdstock-Brown — Inngest CEO",
     "Techniques for building effective multi-agent networks deterministically while reducing compounding error rates across AI calls.",
     ACCENT_GREEN),
    ("Measure What Matters: Quality Monitoring", "Databricks",
     "Mosaic Agent Monitoring: user feedback integration, monitoring dashboards, quality evaluation using AI judges, and operational metrics tracking.",
     ACCENT_PURPLE),
    ("Building Responsible AI Agents", "Databricks",
     "Fairness, transparency, regulatory compliance. Bias monitoring, explainability (SHAP/LIME), Unity Catalog governance, and LLM security guardrails.",
     ACCENT_TEAL),
    ("Agentic AI in Finance", "Qubika (Sponsored)",
     "AI finance agents using LangChain, RAG, and Databricks unified platform including MLFlow and Mosaic AI integration.",
     ACCENT_RED),
]

card_w = Inches(3.6)
card_h = Inches(2.2)
for i, (title, speaker, desc, color) in enumerate(sessions_2025):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (card_w + Inches(0.2))
    y = Inches(1.8) + row * (card_h + Inches(0.2))
    card = add_card(slide, x, y, card_w, card_h)
    add_accent_bar(slide, x, y, card_w, Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.5),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), card_w - Inches(0.3), Inches(0.3),
                 speaker, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.95), card_w - Inches(0.3), Inches(1.1),
                 desc, font_size=11, color=LIGHT_GRAY)


# ── SLIDE 6: KEY SUMMIT SESSIONS – AI AGENTS (2026) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Key Summit Sessions: AI Agents (2026)", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(4), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "June 15–18, 2026 · San Francisco · 800+ Sessions Expected", font_size=13, color=LIGHT_GRAY)

sessions_2026 = [
    ("Beyond the Trace: adidas' Agent Digital Twin",
     "adidas",
     "Maps agent→tool→prompt→retrieval→model call→post-processing chains. "
     "Rollup of risk, quality, and unit economics across the fleet. "
     "Uses MLflow Tracing + Unity Catalog for audit evidence and DBU/cost signals.",
     ["200+ serving endpoints", "300+ data & AI products", "6,000+ registered models", "600,000+ pipeline runs"],
     ACCENT_ORANGE),
    ("Deploying and Monitoring Agents on Databricks",
     "Databricks",
     "Hands-on course covering batch (AI Functions) and real-time (Model Serving) "
     "agent deployment. MLflow tracing, monitoring dashboards, and observability.",
     ["Batch mode via AI Functions", "Real-time REST endpoints", "Built-in trace collection", "Monitoring dashboards"],
     ACCENT_BLUE),
    ("Building RAG Agents with Agent Bricks",
     "Databricks",
     "Hands-on training using Vector Search, AI Playground, and Agent Bricks "
     "for building production RAG agents.",
     ["Vector Search integration", "AI Playground testing", "Auto-optimized agents", "No-code building"],
     ACCENT_GREEN),
    ("Building Tool-Calling Agents with MCP",
     "Databricks",
     "Creating agents with tool-calling capabilities using the Databricks Agent "
     "Framework and Model Context Protocol (MCP).",
     ["Model Context Protocol", "Tool-calling capabilities", "Agent Framework SDK", "Enterprise integration"],
     ACCENT_PURPLE),
]

for i, (title, company, desc, bullets, color) in enumerate(sessions_2026):
    col = i % 2
    row = i // 2
    cw = Inches(5.7)
    ch = Inches(2.5)
    x = Inches(0.8) + col * (cw + Inches(0.2))
    y = Inches(1.8) + row * (ch + Inches(0.2))
    card = add_card(slide, x, y, cw, ch)
    add_accent_bar(slide, x, y, cw, Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), cw - Inches(0.3), Inches(0.5),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(0.8), Inches(0.3),
                 company, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.85), Inches(3.0), Inches(1.5),
                 desc, font_size=11, color=LIGHT_GRAY)
    bullet_lines = []
    for b in bullets:
        bullet_lines.append({"text": f"•  {b}", "size": 11, "color": WHITE, "space_after": 3})
    add_multi_text(slide, x + Inches(3.3), y + Inches(0.85), Inches(2.2), Inches(1.5), bullet_lines)


# ── SLIDE 7: TOKEN USAGE & COST ECONOMICS ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Token Usage & Cost Economics", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2), Pt(4), ACCENT_GREEN)

# Pricing section
card = add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(5.5), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Databricks Inference Pricing (DBRX)", font_size=18, color=ACCENT_ORANGE, bold=True)
pricing_lines = [
    {"text": "Input Tokens:     $0.0005 / 1K tokens", "size": 15, "color": WHITE, "bold": True, "space_after": 8},
    {"text": "Output Tokens:   $0.0015 / 1K tokens", "size": 15, "color": WHITE, "bold": True, "space_after": 8},
    {"text": "Agent Evaluation:  $0.018 / judge request", "size": 15, "color": WHITE, "bold": True, "space_after": 8},
    {"text": "(5,000 tokens per judge request increment)", "size": 12, "color": LIGHT_GRAY, "space_after": 4},
]
add_multi_text(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(1.8), pricing_lines)

# Production costs
card = add_card(slide, Inches(6.6), Inches(1.5), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(6.6), Inches(1.5), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(1.65), Inches(5), Inches(0.4),
             "Production Agent Costs (Monthly)", font_size=18, color=ACCENT_BLUE, bold=True)
cost_lines = [
    {"text": "Inference (1M conversations):  $3,000 – $7,000", "size": 15, "color": WHITE, "bold": True, "space_after": 8},
    {"text": "Base Compute:                          $5,000 – $15,000", "size": 15, "color": WHITE, "bold": True, "space_after": 8},
    {"text": "+ Storage, networking, and monitoring overhead", "size": 13, "color": LIGHT_GRAY, "space_after": 4},
    {"text": "AI Functions: 3x faster, 4x lower cost vs. competitors", "size": 13, "color": ACCENT_GREEN, "bold": True, "space_after": 4},
]
add_multi_text(slide, Inches(6.9), Inches(2.1), Inches(5), Inches(1.8), cost_lines)

# Token optimization stats
stats_token = [
    ("80%", "New Databases Created\nby AI Agents", ACCENT_ORANGE),
    ("97%", "Test & Dev Environments\nCreated by Agents", ACCENT_BLUE),
    ("98.5%", "Token Efficiency Achieved\n(adidas: 200K → 3K tokens)", ACCENT_GREEN),
    ("91.67%", "Cost Savings Through\nToken Optimization", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(stats_token):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(4.4), Inches(2.7), Inches(1.6), val, label, color)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "Gartner predicts >40% of agentic AI projects will be canceled by end-2027 as costs rise and controls lag.",
             font_size=12, color=ACCENT_RED, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
             "Sources: Databricks Summit Sessions, Getmonetizely Analysis, adidas Summit Presentation",
             font_size=10, color=RGBColor(0x70, 0x78, 0x90))


# ── SLIDE 8: COMPANY CASE STUDIES ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Company Case Studies", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.8), Pt(4), ACCENT_PURPLE)

companies = [
    ("adidas", "Agent Digital Twin for Governance & Cost",
     ["200+ serving endpoints, 300+ data & AI products, 6,000+ registered models",
      "98.5% token efficiency (200K → 3K input tokens)",
      "91.67% cost savings through LLM transition & optimization",
      "Per-hop unit economics and cost leak forensics",
      "MLflow Tracing + Unity Catalog for audit evidence"],
     ACCENT_ORANGE),
    ("Walmart", "Autonomous AI Agents & Self-Service Analytics",
     ["90% reduction in time-to-value with AI/BI Genie",
      "$5.6M annual savings from FTE hours",
      "AI chatbot 'Sparky' driving increased average spend",
      "Plans for conversational AI agents for store managers",
      "Presented architecture for autonomous agents at 2025 Summit"],
     ACCENT_BLUE),
    ("AstraZeneca", "Clinical Trial Document Processing",
     ["400,000+ clinical trial documents reviewed in under 60 minutes",
      "No coding required using Agent Bricks",
      "Surpassed open-source alternatives in accuracy metrics",
      "Production-grade agent with enterprise data integration"],
     ACCENT_GREEN),
    ("Royal Bank of Canada", "AI Agents for Equity Research",
     ["Reduced equity research work from days to minutes",
      "Agents work with proprietary enterprise financial data",
      "Highlighted by CEO Ali Ghodsi as key enterprise success story",
      "Demonstrates value of agents with proprietary data"],
     ACCENT_PURPLE),
    ("North Dakota University", "Legislative Data Extraction",
     ["Saved ~30 days of manual labor",
      "Automated legislative data extraction and processing",
      "Built using Agent Bricks no-code platform"],
     ACCENT_TEAL),
    ("Hawaiian Electric", "AI Agent Accuracy",
     ["Surpassed open-source alternatives in accuracy",
      "Production agent deployment on Databricks platform",
      "Enterprise-grade reliability metrics achieved"],
     ACCENT_RED),
]

card_w = Inches(3.7)
card_h = Inches(2.6)
for i, (company, subtitle, bullets, color) in enumerate(companies):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (card_w + Inches(0.2))
    y = Inches(1.5) + row * (card_h + Inches(0.2))
    card = add_card(slide, x, y, card_w, card_h)
    add_accent_bar(slide, x, y, card_w, Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.12), card_w - Inches(0.3), Inches(0.35),
                 company, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.45), card_w - Inches(0.3), Inches(0.3),
                 subtitle, font_size=11, color=WHITE, bold=True)

    bullet_lines = []
    for b in bullets:
        bullet_lines.append({"text": f"•  {b}", "size": 10, "color": LIGHT_GRAY, "space_after": 2})
    add_multi_text(slide, x + Inches(0.15), y + Inches(0.8), card_w - Inches(0.3), card_h - Inches(1.0), bullet_lines)


# ── SLIDE 9: STATE OF AI AGENTS REPORT ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "State of AI Agents Report (2026)", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2), Pt(4), ACCENT_TEAL)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Based on 20,000+ global customers · Published January 2026", font_size=13, color=LIGHT_GRAY)

# Row 1 stats
report_stats = [
    ("327%", "Multi-Agent Workflow\nGrowth in 4 Months", ACCENT_ORANGE),
    ("80%", "New Databases Built\nby AI Agents", ACCENT_BLUE),
    ("19%", "Organizations with Agents\nDeployed at Scale", ACCENT_GREEN),
    ("37%", "Use Supervisor Agent\nArchitecture", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(report_stats):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(1.8), Inches(2.7), Inches(1.6), val, label, color)

# Governance Impact section
card = add_card(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(3))
add_accent_bar(slide, Inches(0.8), Inches(3.8), Inches(5.5), Pt(4), ACCENT_TEAL)
add_text_box(slide, Inches(1.1), Inches(3.95), Inches(5), Inches(0.4),
             "Governance Impact on Production Success", font_size=18, color=ACCENT_TEAL, bold=True)
gov_lines = [
    {"text": "12x", "size": 36, "color": ACCENT_ORANGE, "bold": True, "space_after": 2},
    {"text": "More AI projects in production for companies using governance tools", "size": 14, "color": WHITE, "space_after": 16},
    {"text": "6x", "size": 36, "color": ACCENT_BLUE, "bold": True, "space_after": 2},
    {"text": "More production deployments for companies using evaluation tools", "size": 14, "color": WHITE, "space_after": 16},
    {"text": "7x", "size": 36, "color": ACCENT_GREEN, "bold": True, "space_after": 2},
    {"text": "Growth in AI Gateway governance product usage in 9 months", "size": 14, "color": WHITE, "space_after": 4},
]
add_multi_text(slide, Inches(1.1), Inches(4.35), Inches(5), Inches(2.3), gov_lines)

# Use cases section
card = add_card(slide, Inches(6.6), Inches(3.8), Inches(5.5), Inches(3))
add_accent_bar(slide, Inches(6.6), Inches(3.8), Inches(5.5), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(6.9), Inches(3.95), Inches(5), Inches(0.4),
             "Top AI Agent Use Cases", font_size=18, color=ACCENT_PURPLE, bold=True)
use_case_lines = [
    {"text": "40%  Customer Experience", "size": 16, "color": ACCENT_ORANGE, "bold": True, "space_after": 6},
    {"text": "       Support, onboarding, and engagement agents", "size": 12, "color": LIGHT_GRAY, "space_after": 12},
    {"text": "Market Intelligence", "size": 16, "color": ACCENT_BLUE, "bold": True, "space_after": 6},
    {"text": "       Competitive analysis, trend monitoring, and research", "size": 12, "color": LIGHT_GRAY, "space_after": 12},
    {"text": "Predictive Maintenance", "size": 16, "color": ACCENT_GREEN, "bold": True, "space_after": 6},
    {"text": "       Equipment monitoring, failure prediction, and scheduling", "size": 12, "color": LIGHT_GRAY, "space_after": 12},
    {"text": "Information Extraction", "size": 16, "color": ACCENT_PURPLE, "bold": True, "space_after": 6},
    {"text": "       Document processing, data extraction, and classification", "size": 12, "color": LIGHT_GRAY, "space_after": 4},
]
add_multi_text(slide, Inches(6.9), Inches(4.35), Inches(5), Inches(2.3), use_case_lines)


# ── SLIDE 10: PRODUCTION ARCHITECTURE ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Production Architecture & Best Practices", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2.5), Pt(4), ACCENT_RED)

# Compound AI Systems
card = add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(5.5), Pt(4), ACCENT_ORANGE)
add_text_box(slide, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Compound AI Systems (Mosaic AI)", font_size=18, color=ACCENT_ORANGE, bold=True)
compound_lines = [
    {"text": "•  Tuned foundation models + fine-tuning", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Retrieval-Augmented Generation (RAG)", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Tool use via Model Context Protocol (MCP)", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Reasoning agents with multi-step planning", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Multi-agent orchestration (supervisor pattern)", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "Production systems increasingly use multiple components", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
    {"text": "rather than monolithic models.", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
]
add_multi_text(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(1.8), compound_lines)

# MLflow 3.0
card = add_card(slide, Inches(6.6), Inches(1.5), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(6.6), Inches(1.5), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(1.65), Inches(5), Inches(0.4),
             "MLflow 3.0 — Redesigned for GenAI", font_size=18, color=ACCENT_BLUE, bold=True)
mlflow_lines = [
    {"text": "•  Agent observability and tracing", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Prompt versioning and management", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Cross-platform monitoring (AWS, GCP, on-prem)", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Works for agents deployed outside Databricks", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Integrated with Unity Catalog for governance", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "Open-source, cloud-agnostic agent monitoring", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
    {"text": "that scales with enterprise needs.", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
]
add_multi_text(slide, Inches(6.9), Inches(2.1), Inches(5), Inches(1.8), mlflow_lines)

# Agent Bricks
card = add_card(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(0.8), Inches(4.3), Inches(5.5), Pt(4), ACCENT_GREEN)
add_text_box(slide, Inches(1.1), Inches(4.45), Inches(5), Inches(0.4),
             "Agent Bricks — No-Code Agent Building", font_size=18, color=ACCENT_GREEN, bold=True)
bricks_lines = [
    {"text": "•  Auto-generates domain-specific agents", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Accepts high-level task descriptions", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Connects to enterprise data automatically", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Auto-builds evaluations for quality assurance", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Optimizes for both quality and cost", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "Use cases: information extraction, knowledge assistance,", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
    {"text": "and multi-agent systems.", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
]
add_multi_text(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(1.8), bricks_lines)

# Monitoring best practices
card = add_card(slide, Inches(6.6), Inches(4.3), Inches(5.5), Inches(2.5))
add_accent_bar(slide, Inches(6.6), Inches(4.3), Inches(5.5), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(6.9), Inches(4.45), Inches(5), Inches(0.4),
             "Monitoring & Observability Best Practices", font_size=18, color=ACCENT_PURPLE, bold=True)
monitor_lines = [
    {"text": "•  Collect end-user feedback for iterative improvement", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  AI judges for continuous quality evaluation", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Dimension-based analysis across time periods", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Token usage tracking via AI Gateway system tables", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "•  Per-hop cost forensics for agent chains", "size": 13, "color": WHITE, "space_after": 4},
    {"text": "system.ai_gateway.usage tables provide built-in", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
    {"text": "dashboards for cost, latency & performance.", "size": 12, "color": LIGHT_GRAY, "space_after": 2},
]
add_multi_text(slide, Inches(6.9), Inches(4.9), Inches(5), Inches(1.8), monitor_lines)


# ── SLIDE 11: CEO VISION & GROWTH TRAJECTORY ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "CEO Vision & Growth Trajectory", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(2), Pt(4), ACCENT_ORANGE)

# Ali Ghodsi quote card
card = add_card(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.8))
add_accent_bar(slide, Inches(0.8), Inches(1.5), Pt(5), Inches(1.8), ACCENT_ORANGE)
quote_lines = [
    {"text": '"It\'s very elusive to get AI that really works and understands that proprietary data', "size": 18, "color": WHITE, "bold": False, "space_after": 2},
    {"text": 'that\'s inside enterprise... this could maybe get all the way to a trillion."', "size": 18, "color": WHITE, "bold": False, "space_after": 8},
    {"text": "— Ali Ghodsi, CEO & Co-Founder, Databricks", "size": 14, "color": ACCENT_ORANGE, "bold": True, "space_after": 4},
]
add_multi_text(slide, Inches(1.3), Inches(1.65), Inches(10.5), Inches(1.5), quote_lines)

# Growth stats
growth_stats = [
    ("$134B", "Latest Valuation\n($7B Equity + Debt Raise)", ACCENT_ORANGE),
    ("20,000+", "Global Customers\nAcross Industries", ACCENT_BLUE),
    ("60%+", "Fortune 500 Companies\non Platform", ACCENT_GREEN),
    ("80%", "Databases Now\nCreated by AI Agents", ACCENT_PURPLE),
]
for i, (val, label, color) in enumerate(growth_stats):
    add_stat_card(slide, Inches(0.8) + i * Inches(3.0), Inches(3.7), Inches(2.7), Inches(1.6), val, label, color)

# Three pillars
pillars = [
    ("1. Enterprise Data + AI Agents",
     "AI agents that work with proprietary enterprise data represent the key differentiator. "
     "General AI knowledge is commoditized, but enterprise-specific AI creates unique value.",
     ACCENT_ORANGE),
    ("2. Explosive Database Growth",
     "80% of databases launched on Databricks are now created by AI agents, not humans. "
     "This drove the creation of Lakebase — a new elastic database for agent-scale workloads.",
     ACCENT_BLUE),
    ("3. Platform Openness",
     "Cloud-agnostic, open-source approach (MLflow, Unity Catalog) attracts enterprises. "
     "Customers value flexibility to run on AWS, GCP, Azure, or on-premise.",
     ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(pillars):
    x = Inches(0.8) + i * Inches(3.95)
    y = Inches(5.6)
    card = add_card(slide, x, y, Inches(3.7), Inches(1.5))
    add_accent_bar(slide, x, y, Inches(3.7), Pt(4), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.12), Inches(3.4), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.45), Inches(3.4), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GRAY)


# ── SLIDE 12: KEY TAKEAWAYS ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Key Takeaways", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.5), Pt(4), ACCENT_ORANGE)

takeaways = [
    ("AI Agents Are Reshaping Enterprise Data",
     "80% of databases on Databricks are now created by AI agents. Multi-agent workflows grew 327% in just 4 months. "
     "The shift from single chatbots to autonomous multi-agent systems is accelerating.",
     ACCENT_ORANGE),
    ("Token Optimization Is Critical for ROI",
     "Companies like adidas achieved 98.5% token efficiency (200K → 3K tokens) and 91.67% cost savings. "
     "Without governance, Gartner predicts >40% of agentic AI projects will be canceled by 2027.",
     ACCENT_BLUE),
    ("Governance Multiplies Production Success",
     "Companies using AI governance tools deploy 12x more projects to production. Evaluation tools yield 6x more deployments. "
     "AI Gateway governance usage grew 7x in 9 months.",
     ACCENT_GREEN),
    ("Enterprise Case Studies Show Massive Value",
     "Walmart saved $5.6M annually; AstraZeneca reviewed 400K documents in 60 minutes; Royal Bank of Canada reduced "
     "research from days to minutes; North Dakota University saved 30 days of manual labor.",
     ACCENT_PURPLE),
    ("The Technology Stack Is Maturing",
     "Mosaic AI for compound systems, MLflow 3.0 for observability, Agent Bricks for no-code building, "
     "and MCP for tool-calling — the enterprise agent platform is becoming production-ready.",
     ACCENT_TEAL),
]

for i, (title, desc, color) in enumerate(takeaways):
    y = Inches(1.5) + i * Inches(1.15)
    card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.0))
    add_accent_bar(slide, Inches(0.8), y, Pt(5), Inches(1.0), color)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.4), Inches(10.8), Inches(0.55),
                 desc, font_size=12, color=LIGHT_GRAY)


# ── SLIDE 13: SOURCES ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_ORANGE)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Sources & References", font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.8), Pt(4), ACCENT_ORANGE)

sources = [
    "Databricks Data + AI Summit 2024 — Mosaic AI Compound Systems Keynote & Sessions",
    "Databricks Data + AI Summit 2025 — Agent Sessions (June 9–12, San Francisco)",
    "Databricks Data + AI Summit 2026 — Agent Sessions (June 15–18, San Francisco)",
    "Databricks State of Data + AI Report 2024 (10,000+ customers, 300+ Fortune 500)",
    "Databricks State of AI Agents Report 2026 (20,000+ global customers)",
    "Databricks Blog: Enterprise AI Agent Trends — Use Cases, Governance & Evaluations",
    "Databricks Blog: Mosaic AI Announcements at Data + AI Summit 2025",
    "adidas: Beyond the Trace — Agent Digital Twin for Governance, Cost, and ROI (Summit 2026)",
    "Walmart: Autonomous AI Agents in AI Infrastructure (Summit 2025)",
    "Walmart: Self-Service Assortment and Space Analytics at Walmart Scale (Summit 2025)",
    "AstraZeneca: Clinical Trial Document Processing via Agent Bricks",
    "Royal Bank of Canada: AI Agents for Equity Research (CEO Keynote Reference)",
    "Ali Ghodsi (Databricks CEO): Fortune Brainstorm AI Interview, December 2025",
    "Getmonetizely: Production-Grade AI Agents on Databricks Mosaic — Cost Analysis",
    "Databricks Docs: AI Gateway Usage Tracking (system.ai_gateway.usage tables)",
]

source_lines = []
for i, s in enumerate(sources):
    source_lines.append({"text": f"{i+1}.  {s}", "size": 12, "color": LIGHT_GRAY, "space_after": 5})

add_multi_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5.5), source_lines)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "All data sourced from publicly available Databricks summit sessions, blog posts, and reports.",
             font_size=11, color=RGBColor(0x70, 0x78, 0x90))

# Save
output_path = "/workspace/databricks_summit_ai_agents_token_usage.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
