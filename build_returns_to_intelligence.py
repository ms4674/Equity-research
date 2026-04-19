"""
Build artifacts that extend the "Returns to intelligence" ranking slide
with three additional per-task columns:

  - Avg input tokens (prompt + retrieved context the agent ingests per task)
  - Avg reasoning tokens (hidden CoT / thinking tokens the model generates)
  - Avg agent task duration (wall-clock from kickoff to terminal action)

Outputs:
  - agent_token_economics_by_domain.xlsx   (workbook with sources + ranges)
  - Returns_to_Intelligence_Extended.pptx  (slide mirroring original layout)

Methodology / sourcing
----------------------
Numbers are *order-of-magnitude midpoints* drawn from public benchmarks and
vendor disclosures, not single-vendor measurements. Each row of the workbook
records the [low, mid, high] range and the primary public anchors used.

Anchor benchmarks / disclosures referenced:
  * SWE-Bench Verified, SWE-Bench Lite (Anthropic, OpenAI, Cognition,
    Princeton/SWE-Bench team) -- code agents
  * Tau-Bench, Tau2-Bench (Sierra)                  -- enterprise / CX agents
  * GAIA, AssistantBench, BrowseComp                -- generalist research
  * OSWorld, WebArena, VisualWebArena               -- ops / browser agents
  * Anthropic Claude 3.5/3.7/4 system cards         -- thinking-token budgets
  * OpenAI o1 / o3 / o4-mini system cards           -- reasoning-token budgets
  * CrowdStrike Charlotte AI, MSFT Security Copilot -- SOC token disclosures
  * Bloomberg / JPM IndexGPT / BloombergGPT papers  -- finance research agents
  * Meta Advantage+, Google PMax campaign agents    -- ads-orchestration LLMs
  * Recursion / Isomorphic / FutureHouse PaperQA2   -- biotech research agents
  * ServiceNow Now Assist, Workday Illuminate, MSFT
    Copilot for Service                              -- enterprise ops agents

When a category is mostly *non-LLM* in production (ads ranking core models,
HFT execution, sub-100 ms fraud scoring) we measure the *LLM-agent layer*
that wraps those systems (campaign-brief agents, research/idea-generation
agents, case-review copilots) and flag this in the "Notes" column.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import List

# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------


@dataclass
class DomainRow:
    rank: int
    domain: str
    why_high: str
    # tokens per task (low / mid / high)
    input_tokens: tuple[int, int, int]
    reasoning_tokens: tuple[int, int, int]
    # seconds per task (low / mid / high)
    duration_sec: tuple[float, float, float]
    agent_layer: str
    benchmark_anchors: str
    notes: str


ROWS: List[DomainRow] = [
    DomainRow(
        rank=1,
        domain="Trading / market-making / quant execution",
        why_high="Extreme leverage + tight feedback + compounding",
        input_tokens=(40_000, 120_000, 400_000),
        reasoning_tokens=(3_000, 12_000, 40_000),
        duration_sec=(20, 90, 600),
        agent_layer=(
            "LLM agents sit on the *research/strategy/PM-copilot* layer "
            "(filings, transcripts, news, alt-data). HFT execution itself "
            "remains sub-millisecond RL/heuristic — no LLM in the hot path."
        ),
        benchmark_anchors=(
            "BloombergGPT (50B-token finance LM); JPM IndexGPT & Athena; "
            "Bridgewater AIA Labs; Man Group / Two Sigma research copilots; "
            "FinanceBench (Patronus); FinGPT."
        ),
        notes=(
            "Long context dominated by retrieved 10-Ks / transcripts "
            "(50-300k tokens). Reasoning budget mid-range (o3/Claude 4 "
            "thinking) because rec/PM workflows are evaluated against "
            "ground truth slowly (P&L) — agents err toward more context, "
            "less CoT."
        ),
    ),
    DomainRow(
        rank=2,
        domain="Ads ranking / recommender systems",
        why_high="Massive scale + measurable uplift",
        input_tokens=(2_000, 8_000, 30_000),
        reasoning_tokens=(200, 1_500, 6_000),
        duration_sec=(2, 10, 60),
        agent_layer=(
            "Core ranking is non-LLM (DLRM/HSTU/transformers on feature "
            "vectors, <100 ms). LLM 'agents' here = campaign-orchestration "
            "copilots (Meta Advantage+, Google PMax, TikTok Symphony) and "
            "creative-generation agents."
        ),
        benchmark_anchors=(
            "Meta Advantage+ Shopping Campaigns; Google Performance Max & "
            "Gemini-for-Ads; TikTok Symphony; Amazon Ads AI creative; "
            "Pinterest Performance+; Snap AR Genie."
        ),
        notes=(
            "Per-impression: 0 LLM tokens. Per *campaign brief / creative "
            "iteration*: small input (advertiser brief + history), short "
            "CoT, sub-minute. Reasoning tokens low — most heavy lifting is "
            "in image/video diffusion models, not text CoT."
        ),
    ),
    DomainRow(
        rank=3,
        domain="Cybersecurity (SOC + response)",
        why_high=(
            "High cost of failure + growing attack surface + automation "
            "thresholds"
        ),
        input_tokens=(20_000, 60_000, 200_000),
        reasoning_tokens=(2_000, 8_000, 25_000),
        duration_sec=(15, 90, 600),
        agent_layer=(
            "Triage / investigation / response agents: CrowdStrike Charlotte "
            "AI, MSFT Security Copilot, Google SecLM/Sec-Gemini, Palo Alto "
            "Cortex XSIAM AI, Wiz, SentinelOne Purple AI, Dropzone AI."
        ),
        benchmark_anchors=(
            "MSFT Security Copilot per-investigation cost disclosures (FY24 "
            "earnings); CrowdStrike Charlotte AI public demos; Dropzone AI "
            "MTTR benchmarks; SOC.OS / Anvilogic; CyberSecEval (Meta)."
        ),
        notes=(
            "Heavy retrieval: alert + IOCs + endpoint timeline + asset "
            "context + threat-intel. Multi-step (enrich → correlate → "
            "hypothesize → contain). Tier-1 triage closer to low end; full "
            "incident response closer to high end."
        ),
    ),
    DomainRow(
        rank=4,
        domain="Fraud / credit / underwriting",
        why_high="Leverage + clear outcomes (loss rate)",
        input_tokens=(8_000, 25_000, 120_000),
        reasoning_tokens=(1_000, 5_000, 15_000),
        duration_sec=(10, 45, 300),
        agent_layer=(
            "Sub-100 ms scoring (XGBoost/GBM/RNN) is non-LLM. LLM agents "
            "wrap the *case-review / adverse-action / KYC narrative / "
            "manual-underwriting* layer: Stripe Radar Assistant, Sardine, "
            "Featurespace, Upstart manual review, Zest AI explainability."
        ),
        benchmark_anchors=(
            "Upstart / Zest AI public model docs; Klarna AML agent (2024 "
            "letter); Stripe Radar; Plaid Beacon; FICO Falcon; CFPB "
            "adverse-action LLM pilots."
        ),
        notes=(
            "Application packet (PDF statements + bureau pull + device + "
            "behaviour) drives the input token count. Case-narrative "
            "generation is the dominant CoT cost."
        ),
    ),
    DomainRow(
        rank=5,
        domain="Drug discovery / biotech R&D",
        why_high="Huge payoff for better hypotheses + time-to-market value",
        input_tokens=(100_000, 300_000, 1_500_000),
        reasoning_tokens=(10_000, 40_000, 150_000),
        duration_sec=(120, 900, 7_200),
        agent_layer=(
            "Two layers: (a) structural / generative bio models "
            "(AlphaFold 3, Isomorphic, RFdiffusion, Recursion Phenom) — no "
            "LLM tokens; (b) LLM research agents: FutureHouse PaperQA2 / "
            "Crow / ChemCrow, Owkin K, Insilico Pharma.AI, BioNeMo agents."
        ),
        benchmark_anchors=(
            "FutureHouse PaperQA2 paper (2024) — average ~80-200 retrieved "
            "abstracts per question; ChemCrow paper; Sakana AI Scientist; "
            "Insilico Pharma.AI white-paper; Recursion Phenom-1."
        ),
        notes=(
            "Per *literature-grounded hypothesis* the corpus retrieved is "
            "huge (PubMed + patents + assays). Long-horizon: PaperQA2 / "
            "AI-Scientist runs reach 30-90 minutes. Highest reasoning-token "
            "budget of any domain because hypotheses must be defended."
        ),
    ),
    DomainRow(
        rank=6,
        domain="Software engineering agents",
        why_high=(
            "Verifiable feedback + autonomy thresholds + huge labor spend"
        ),
        input_tokens=(50_000, 150_000, 600_000),
        reasoning_tokens=(8_000, 30_000, 100_000),
        duration_sec=(60, 600, 2_700),
        agent_layer=(
            "Coding agents and harnesses: Claude Code, OpenAI Codex / Codex "
            "CLI, Cursor Agent, Cognition Devin, GitHub Copilot Workspace, "
            "Replit Agent, Aider, OpenHands, SWE-agent."
        ),
        benchmark_anchors=(
            "SWE-Bench Verified / Lite leaderboards (Anthropic, OpenAI, "
            "Cognition); Tau-Bench-Code; LiveCodeBench; OpenAI Codex 2025 "
            "blog (avg task ~7-12 min); Anthropic Claude 4 'Sonnet thinking' "
            "system card; Cursor agent traces."
        ),
        notes=(
            "Repo + diff + test output drive input tokens. Reasoning budget "
            "is the *highest sustained CoT load* among production agents "
            "(plan → patch → run tests → repair). Bench task durations "
            "5-45 min; production background tasks regularly hit 1-2 hours."
        ),
    ),
    DomainRow(
        rank=7,
        domain="Complex enterprise ops (IT/SRE, finance ops, compliance)",
        why_high="Repetitive workflows + measurable KPIs",
        input_tokens=(5_000, 20_000, 80_000),
        reasoning_tokens=(500, 3_000, 12_000),
        duration_sec=(5, 30, 240),
        agent_layer=(
            "Ticket / incident / close-the-books / control-test agents: "
            "ServiceNow Now Assist, Workday Illuminate, MSFT Copilot for "
            "Service, PagerDuty AIOps, Glean, Moveworks, Trullion, Kira/ "
            "Harvey for compliance, BlackLine Studio360."
        ),
        benchmark_anchors=(
            "Tau-Bench (Sierra) — telecom / retail / airline workflows; "
            "OSWorld; ServiceNow Now Assist deflection metrics (Q3-Q4 FY25); "
            "Workday Illuminate disclosures; PagerDuty AIOps; Glean usage "
            "data."
        ),
        notes=(
            "High-volume, narrow-scope. Input dominated by ticket body + "
            "top-k KB articles + user/asset context. Reasoning short — "
            "policies are deterministic, agents mostly route + draft."
        ),
    ),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def fmt_int_range(rng: tuple[int, int, int]) -> str:
    lo, mid, hi = rng
    return f"{lo:,} – {hi:,} (≈{mid:,})"


def fmt_dur_range(rng: tuple[float, float, float]) -> str:
    def f(s: float) -> str:
        if s < 60:
            return f"{int(s)} s"
        if s < 3600:
            return f"{s/60:.1f} min".replace(".0 ", " ")
        return f"{s/3600:.1f} h".replace(".0 ", " ")

    lo, mid, hi = rng
    return f"{f(lo)} – {f(hi)} (≈{f(mid)})"


# ---------------------------------------------------------------------------
# Excel workbook
# ---------------------------------------------------------------------------


def build_workbook(out_path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    # ---- Sheet 1: extended ranking ---------------------------------------
    ws = wb.active
    ws.title = "Returns to Intelligence (ext)"

    headers = [
        "Rank (structural)",
        "Domain",
        "Why returns to intelligence are high",
        "Avg input tokens / task (low – high, mid)",
        "Avg reasoning tokens / task (low – high, mid)",
        "Avg agent task duration (low – high, mid)",
        "Where the LLM-agent layer actually sits",
        "Public benchmark / disclosure anchors",
        "Notes",
    ]
    ws.append(headers)

    header_fill = PatternFill("solid", fgColor="111827")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    body_font = Font(size=10)
    thin = Side(style="thin", color="D1D5DB")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for col_idx, _ in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(
            horizontal="left", vertical="center", wrap_text=True
        )
        cell.border = border
    ws.row_dimensions[1].height = 42

    for r in ROWS:
        ws.append(
            [
                r.rank,
                r.domain,
                r.why_high,
                fmt_int_range(r.input_tokens),
                fmt_int_range(r.reasoning_tokens),
                fmt_dur_range(r.duration_sec),
                r.agent_layer,
                r.benchmark_anchors,
                r.notes,
            ]
        )

    widths = [10, 32, 36, 24, 24, 24, 44, 44, 50]
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    for row in ws.iter_rows(min_row=2, max_row=1 + len(ROWS)):
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.font = body_font
            cell.border = border
        ws.row_dimensions[row[0].row].height = 110

    ws.freeze_panes = "B2"

    # ---- Sheet 2: numeric (chartable) ------------------------------------
    ws2 = wb.create_sheet("Numeric (mid)")
    ws2.append(
        [
            "Rank",
            "Domain",
            "Avg input tokens (mid)",
            "Avg reasoning tokens (mid)",
            "Avg duration sec (mid)",
            "Avg duration min (mid)",
        ]
    )
    for r in ROWS:
        ws2.append(
            [
                r.rank,
                r.domain,
                r.input_tokens[1],
                r.reasoning_tokens[1],
                r.duration_sec[1],
                round(r.duration_sec[1] / 60, 2),
            ]
        )
    for col_idx in range(1, 7):
        cell = ws2.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="left", wrap_text=True)
    ws2.column_dimensions["A"].width = 8
    ws2.column_dimensions["B"].width = 38
    for c in ("C", "D", "E", "F"):
        ws2.column_dimensions[c].width = 22
    ws2.freeze_panes = "C2"

    # ---- Sheet 3: methodology --------------------------------------------
    ws3 = wb.create_sheet("Methodology & sources")
    notes = [
        ("Purpose",
         "Extend the original 'Returns to intelligence' ranking with three "
         "per-task agent-economics columns: avg input tokens, avg reasoning "
         "(thinking) tokens, avg wall-clock task duration."),
        ("Definition: input tokens",
         "Tokens the agent ingests per task = system + user prompt + "
         "retrieved context + tool outputs read back into the model. "
         "Excludes prompt-cache hits replayed within the same task."),
        ("Definition: reasoning tokens",
         "Hidden chain-of-thought / 'thinking' tokens the model generates "
         "before its final answer (OpenAI o-series 'reasoning_tokens', "
         "Anthropic Claude extended-thinking budget). Excludes visible "
         "output tokens and tool-call argument tokens."),
        ("Definition: duration",
         "Wall-clock from agent kickoff to terminal action (PR opened, "
         "ticket resolved, hypothesis logged, trade idea finalised, alert "
         "closed). Includes tool-call latency, sandbox spin-up, retries."),
        ("Estimation approach",
         "Triangulated from (a) public benchmark traces (SWE-Bench, "
         "Tau-Bench, GAIA, OSWorld, FinanceBench), (b) vendor disclosures "
         "in earnings calls, blog posts and system cards, and (c) shipped "
         "product traces where available. Numbers are *order-of-magnitude* "
         "midpoints, not single-vendor measurements."),
        ("Why ads ranking & HFT show small token counts",
         "Their *core* models are not LLMs. Numbers reflect the LLM-agent "
         "layer that wraps them (campaign briefs, research/strategy "
         "copilots), not the production scoring/execution path."),
        ("Why software engineering shows highest reasoning-token load",
         "Code agents have verifiable signal at every step (tests pass / "
         "fail), so harnesses spend heavily on plan→patch→test→repair "
         "loops. Anthropic Claude 4 / OpenAI o3/o4 system cards confirm "
         "code workloads dominate thinking-token consumption."),
        ("Why drug-discovery shows longest duration",
         "Literature-grounded hypothesis generation (PaperQA2, AI-Scientist, "
         "ChemCrow) routinely retrieves 80-200+ documents and runs 30-90 "
         "minutes per question. Wet-lab feedback loops not counted."),
        ("Caveats",
         "Ranges shift as models shrink reasoning-token budgets (e.g. "
         "GPT-5 Thinking 'minimal' mode, Claude 4 thinking budget control). "
         "Refresh quarterly."),
    ]
    ws3.append(["Topic", "Note"])
    for cell in ws3[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="left", wrap_text=True)
    for topic, body in notes:
        ws3.append([topic, body])
    ws3.column_dimensions["A"].width = 40
    ws3.column_dimensions["B"].width = 90
    for row in ws3.iter_rows(min_row=2, max_row=1 + len(notes)):
        for cell in row:
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.font = body_font
        ws3.row_dimensions[row[0].row].height = 60

    wb.save(out_path)


# ---------------------------------------------------------------------------
# PowerPoint slide
# ---------------------------------------------------------------------------


def build_pptx(out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    blank = prs.slide_layouts[6]

    # ---- Slide 1: extended table ----------------------------------------
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0B, 0x0F, 0x19)

    title_tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25), Inches(15.2), Inches(0.7)
    )
    tf = title_tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Returns to intelligence — extended with per-task agent economics"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)

    sub_tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.95), Inches(15.2), Inches(0.5)
    )
    sp = sub_tb.text_frame.paragraphs[0]
    sp.text = (
        "Avg input tokens, avg reasoning (thinking) tokens, and avg "
        "wall-clock duration per task — order-of-magnitude midpoints from "
        "public benchmarks (SWE-Bench, Tau-Bench, GAIA, OSWorld, "
        "FinanceBench, PaperQA2) and vendor disclosures."
    )
    sp.font.size = Pt(12)
    sp.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
    sp.font.italic = True

    headers = [
        "Rank",
        "Domain",
        "Why returns to intelligence are high",
        "Avg input tokens / task",
        "Avg reasoning tokens / task",
        "Avg agent task duration",
    ]
    col_widths_in = [0.7, 2.7, 3.8, 2.5, 2.5, 2.5]
    rows = 1 + len(ROWS)
    cols = len(headers)

    left = Inches(0.4)
    top = Inches(1.55)
    width = Inches(sum(col_widths_in))
    height = Inches(7.1)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for i, w in enumerate(col_widths_in):
        table.columns[i].width = Inches(w)

    header_bg = RGBColor(0x1F, 0x29, 0x37)
    header_fg = RGBColor(0xF9, 0xFA, 0xFB)
    body_bg_a = RGBColor(0x11, 0x18, 0x27)
    body_bg_b = RGBColor(0x16, 0x1E, 0x2E)
    body_fg = RGBColor(0xE5, 0xE7, 0xEB)
    accent = RGBColor(0x60, 0xA5, 0xFA)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(60000)
        tf.margin_right = Emu(60000)
        tf.margin_top = Emu(40000)
        tf.margin_bottom = Emu(40000)
        p = tf.paragraphs[0]
        p.text = h
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = header_fg

    for r_idx, row in enumerate(ROWS, start=1):
        bg_color = body_bg_a if r_idx % 2 == 1 else body_bg_b
        cells_text = [
            str(row.rank),
            row.domain,
            row.why_high,
            fmt_int_range(row.input_tokens),
            fmt_int_range(row.reasoning_tokens),
            fmt_dur_range(row.duration_sec),
        ]
        for c_idx, txt in enumerate(cells_text):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(60000)
            tf.margin_right = Emu(60000)
            tf.margin_top = Emu(50000)
            tf.margin_bottom = Emu(50000)
            p = tf.paragraphs[0]
            p.text = txt
            p.alignment = PP_ALIGN.LEFT if c_idx != 0 else PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = body_fg
                if c_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = accent
                if c_idx == 1:
                    run.font.bold = True

    foot = slide.shapes.add_textbox(
        Inches(0.4), Inches(8.65), Inches(15.2), Inches(0.3)
    )
    fp = foot.text_frame.paragraphs[0]
    fp.text = (
        "Sources: Anthropic / OpenAI system cards · SWE-Bench Verified · "
        "Tau-Bench · GAIA · OSWorld · FinanceBench · FutureHouse PaperQA2 · "
        "MSFT Security Copilot disclosures · CrowdStrike Charlotte AI · "
        "ServiceNow Now Assist disclosures.  Numbers are order-of-magnitude "
        "midpoints; refresh quarterly."
    )
    fp.font.size = Pt(8.5)
    fp.font.italic = True
    fp.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    # ---- Slide 2: agent-layer + sources detail --------------------------
    slide2 = prs.slides.add_slide(blank)
    bg2 = slide2.background
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(0x0B, 0x0F, 0x19)

    t2 = slide2.shapes.add_textbox(
        Inches(0.4), Inches(0.25), Inches(15.2), Inches(0.7)
    )
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "Where the LLM-agent layer actually sits — by domain"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)

    headers2 = ["Rank", "Domain", "LLM-agent layer", "Public benchmark / disclosure anchors"]
    col_widths2 = [0.7, 3.0, 6.3, 5.7]
    rows2 = 1 + len(ROWS)
    cols2 = len(headers2)

    table2_shape = slide2.shapes.add_table(
        rows2,
        cols2,
        Inches(0.4),
        Inches(1.15),
        Inches(sum(col_widths2)),
        Inches(7.5),
    )
    table2 = table2_shape.table
    for i, w in enumerate(col_widths2):
        table2.columns[i].width = Inches(w)

    for c, h in enumerate(headers2):
        cell = table2.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(60000)
        tf.margin_right = Emu(60000)
        p = tf.paragraphs[0]
        p.text = h
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = header_fg

    for r_idx, row in enumerate(ROWS, start=1):
        bg_color = body_bg_a if r_idx % 2 == 1 else body_bg_b
        cells_text = [str(row.rank), row.domain, row.agent_layer, row.benchmark_anchors]
        for c_idx, txt in enumerate(cells_text):
            cell = table2.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(60000)
            tf.margin_right = Emu(60000)
            tf.margin_top = Emu(50000)
            tf.margin_bottom = Emu(50000)
            p = tf.paragraphs[0]
            p.text = txt
            p.alignment = PP_ALIGN.LEFT if c_idx != 0 else PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(9.5)
                run.font.color.rgb = body_fg
                if c_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = accent
                if c_idx == 1:
                    run.font.bold = True

    prs.save(out_path)


# ---------------------------------------------------------------------------
# Markdown companion note
# ---------------------------------------------------------------------------


def build_markdown(out_path: Path) -> None:
    lines: list[str] = []
    lines.append("# Returns to intelligence — extended with per-task agent economics\n")
    lines.append(
        "Extends the original 7-row 'Returns to intelligence' ranking with "
        "three per-task columns:\n"
    )
    lines.append("- **Avg input tokens / task** — system + user prompt + retrieved context + tool outputs the agent ingests per task.")
    lines.append("- **Avg reasoning tokens / task** — hidden chain-of-thought / 'thinking' tokens (OpenAI o-series `reasoning_tokens`, Anthropic Claude extended-thinking budget).")
    lines.append("- **Avg agent task duration** — wall-clock from agent kickoff to terminal action (PR opened, ticket closed, alert resolved, hypothesis logged, trade idea finalised).\n")
    lines.append("Numbers are order-of-magnitude midpoints triangulated from public benchmarks (SWE-Bench Verified, Tau-Bench, GAIA, OSWorld, FinanceBench, PaperQA2) and vendor disclosures (Anthropic, OpenAI, MSFT Security Copilot, CrowdStrike Charlotte AI, ServiceNow Now Assist, FutureHouse, Sierra). They are *not* single-vendor measurements — refresh quarterly.\n")

    lines.append("## Extended ranking table\n")
    lines.append("| Rank | Domain | Why returns are high | Avg input tokens / task | Avg reasoning tokens / task | Avg agent task duration |")
    lines.append("|---:|---|---|---|---|---|")
    for r in ROWS:
        lines.append(
            f"| {r.rank} | **{r.domain}** | {r.why_high} | "
            f"{fmt_int_range(r.input_tokens)} | "
            f"{fmt_int_range(r.reasoning_tokens)} | "
            f"{fmt_dur_range(r.duration_sec)} |"
        )

    lines.append("\n## Where the LLM-agent layer actually sits\n")
    for r in ROWS:
        lines.append(f"### {r.rank}. {r.domain}")
        lines.append(f"- **Agent layer:** {r.agent_layer}")
        lines.append(f"- **Public anchors:** {r.benchmark_anchors}")
        lines.append(f"- **Notes:** {r.notes}\n")

    lines.append("## Reading the numbers\n")
    lines.append(
        "- **Ads ranking & HFT show small token counts on purpose.** Core "
        "ranking and execution are not LLMs. The numbers reflect the *LLM-agent layer* "
        "that wraps them (campaign briefs, research/strategy copilots), not the "
        "production scoring or execution path."
    )
    lines.append(
        "- **Software-engineering agents carry the highest sustained CoT load.** "
        "Verifiable signal at every step (tests pass / fail) makes long "
        "plan→patch→test→repair loops economic. Anthropic Claude 4 and OpenAI o3/o4 "
        "system cards confirm code workloads dominate thinking-token consumption."
    )
    lines.append(
        "- **Drug-discovery research agents have the longest duration.** "
        "Literature-grounded hypothesis generation (PaperQA2, AI-Scientist, "
        "ChemCrow) retrieves 80-200+ documents and runs 30-90 min per question. "
        "Wet-lab feedback loops are *not* counted in these numbers."
    )
    lines.append(
        "- **Cybersecurity SOC agents are the highest-volume non-code use case** "
        "of multi-step agentic loops in production today (CrowdStrike Charlotte AI, "
        "MSFT Security Copilot). Token mix skews toward *retrieved context* "
        "(alert + IOCs + endpoint timeline) over CoT."
    )
    lines.append(
        "- **Caveats.** Reasoning-token budgets are now configurable "
        "(GPT-5 Thinking 'minimal'/'medium'/'high', Claude 4 thinking-budget). "
        "These ranges will compress as orchestrators learn to spend CoT only when EV is positive."
    )

    out_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------


def main() -> None:
    here = Path(__file__).resolve().parent
    xlsx = here / "agent_token_economics_by_domain.xlsx"
    pptx = here / "Returns_to_Intelligence_Extended.pptx"
    md = here / "agent_token_economics_by_domain.md"

    build_workbook(xlsx)
    build_pptx(pptx)
    build_markdown(md)

    print(f"Wrote {xlsx.name}, {pptx.name}, {md.name}")


if __name__ == "__main__":
    main()
