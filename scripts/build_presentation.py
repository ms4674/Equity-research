#!/usr/bin/env python3
"""Build a PowerPoint deck of the most useful time series from the workbook.

Time series are extracted from the 'Growth trajectory' / token columns of
agentic_ai_revenue_token_usage_2026.xlsx (see scripts/build_spreadsheet.py).
Source numbers refer to the workbook's Sources sheet.

Regenerate with:  python3 scripts/build_presentation.py
Output:           agentic_ai_time_series_2026.pptx
"""

import os
import tempfile
from datetime import date

import matplotlib
matplotlib.use("Agg")
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

AS_OF = "August 23, 2026"

NAVY = "#1F3864"
BLUE = "#2E75B6"
ORANGE = "#ED7D31"
GREEN = "#548235"
RED = "#C00000"
PURPLE = "#7030A0"
GRAY = "#7F7F7F"
LIGHT = "#F2F6FB"

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "font.size": 12,
    "axes.edgecolor": "#BFBFBF",
    "axes.linewidth": 0.8,
    "axes.grid": True,
    "grid.color": "#E0E6EF",
    "grid.linewidth": 0.7,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "figure.facecolor": "white",
    "axes.facecolor": "white",
})

TMP = tempfile.mkdtemp(prefix="agentic_charts_")


def save(fig, name):
    path = os.path.join(TMP, name)
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def label_last(ax, x, y, text, color, dy=0):
    ax.annotate(text, (x[-1], y[-1]), textcoords="offset points",
                xytext=(8, dy), fontsize=11, fontweight="bold", color=color)


# ---------------------------------------------------------------------------
# Chart 1: Claude Code vs Codex tracked ARR
# ---------------------------------------------------------------------------
def chart_coding_arr():
    cc_x = [date(2025, 11, 15), date(2026, 2, 12), date(2026, 3, 31), date(2026, 4, 30),
            date(2026, 6, 15), date(2026, 7, 6), date(2026, 7, 13), date(2026, 7, 20),
            date(2026, 7, 27), date(2026, 8, 3), date(2026, 8, 10)]
    cc_y = [1.0, 2.5, 6.3, 10.3, 14.0, 14.26, 14.38, 14.47, 14.35, 14.50, 15.12]
    cx_x = [date(2026, 1, 5), date(2026, 7, 6), date(2026, 8, 10)]
    cx_y = [0.1, 5.88, 8.83]

    fig, ax = plt.subplots(figsize=(11.4, 5.3))
    ax.plot(cc_x, cc_y, marker="o", lw=2.6, ms=5, color=ORANGE, label="Claude Code (Anthropic)")
    ax.plot(cx_x, cx_y, marker="s", lw=2.6, ms=5, color=BLUE, label="Codex (OpenAI)")
    label_last(ax, cc_x, cc_y, "$15.1B", ORANGE)
    label_last(ax, cx_x, cx_y, "$8.8B", BLUE)
    ax.annotate("$2.5B company-disclosed anchor\n(Series G, Feb 12, 2026)",
                (cc_x[1], cc_y[1]), textcoords="offset points", xytext=(30, -52),
                fontsize=10, color=GRAY,
                arrowprops=dict(arrowstyle="-", color=GRAY, lw=0.8))
    ax.annotate("~$0 at start\nof 2026", (cx_x[0], cx_y[0]), textcoords="offset points",
                xytext=(-8, 10), fontsize=10, color=GRAY, ha="right")
    ax.set_ylabel("Tracked ARR ($B, annualized)")
    ax.set_ylim(0, 17)
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b '%y"))
    ax.legend(loc="upper left", frameon=False, fontsize=12)
    return save(fig, "coding_arr.png")


# ---------------------------------------------------------------------------
# Chart 2: Cursor ARR ladder (log scale)
# ---------------------------------------------------------------------------
def chart_cursor():
    x = [date(2023, 10, 15), date(2025, 1, 15), date(2025, 6, 5), date(2025, 11, 13),
         date(2026, 2, 15), date(2026, 5, 15)]
    y = [0.001, 0.1, 0.5, 1.0, 2.0, 4.0]
    labels = ["$1M", "$100M", "$500M", "$1B", "$2B", "$4B"]

    fig, ax = plt.subplots(figsize=(11.4, 5.3))
    ax.plot(x, y, marker="o", lw=2.6, ms=6, color=NAVY)
    ax.set_yscale("log")
    for xi, yi, lb in zip(x, y, labels):
        ax.annotate(lb, (xi, yi), textcoords="offset points", xytext=(0, 10),
                    fontsize=11, fontweight="bold", color=NAVY, ha="center")
    ax.axvline(date(2026, 6, 16), color=RED, lw=1.4, ls="--")
    ax.annotate("Jun 16, 2026: SpaceX agrees to acquire\nAnysphere for $60B all-stock (15x ARR)",
                (date(2026, 6, 16), 0.004), fontsize=10.5, color=RED, ha="right",
                xytext=(-8, 0), textcoords="offset points")
    ax.set_ylabel("ARR ($B, log scale)")
    ax.set_ylim(0.0007, 9)
    ax.set_yticks([0.001, 0.01, 0.1, 1.0, 4.0])
    ax.set_yticklabels(["$1M", "$10M", "$100M", "$1B", "$4B"])
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b '%y"))
    return save(fig, "cursor.png")


# ---------------------------------------------------------------------------
# Chart 3: Model provider run-rates
# ---------------------------------------------------------------------------
def chart_providers():
    oa_x = [date(2025, 12, 31), date(2026, 2, 15), date(2026, 7, 31)]
    oa_y = [20, 25, 40]
    an_x = [date(2026, 2, 12), date(2026, 7, 31)]
    an_y = [14, 65]

    fig, ax = plt.subplots(figsize=(11.4, 5.3))
    ax.plot(oa_x, oa_y, marker="o", lw=2.6, ms=6, color=BLUE, label="OpenAI")
    ax.plot(an_x, an_y, marker="s", lw=2.6, ms=6, color=ORANGE, label="Anthropic")
    for xi, yi in zip(oa_x, oa_y):
        ax.annotate(f"${yi}B", (xi, yi), textcoords="offset points", xytext=(0, -18),
                    fontsize=11, fontweight="bold", color=BLUE, ha="center")
    for xi, yi in zip(an_x, an_y):
        ax.annotate(f"${yi}B", (xi, yi), textcoords="offset points", xytext=(0, 10),
                    fontsize=11, fontweight="bold", color=ORANGE, ha="center")
    ax.annotate("Q2 2026: Anthropic quarterly revenue ($11.5B)\novertakes OpenAI for the first time;\nClaude Code contributed ~$8B of the quarter",
                (date(2026, 3, 20), 4), fontsize=10.5, color=GRAY, va="bottom")
    ax.set_ylabel("Annualized revenue run-rate ($B)")
    ax.set_ylim(0, 74)
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b '%y"))
    ax.legend(loc="upper left", frameon=False, fontsize=12)
    return save(fig, "providers.png")


# ---------------------------------------------------------------------------
# Chart 4: Vertical & knowledge agents ARR ($M)
# ---------------------------------------------------------------------------
def chart_verticals():
    series = [
        ("Harvey (legal)", GREEN, "o", (8, 4),
         [date(2024, 12, 31), date(2026, 1, 15), date(2026, 6, 15)], [50, 190, 300]),
        ("OpenEvidence (healthcare)", RED, "^", (8, -16),
         [date(2025, 7, 1), date(2025, 12, 14), date(2026, 7, 21)], [50, 150, 300]),
        ("Glean (enterprise knowledge)", PURPLE, "D", (-14, 12),
         [date(2025, 2, 15), date(2026, 5, 28)], [100, 300]),
        ("Sierra (customer service)", BLUE, "s", (8, -3),
         [date(2024, 12, 31), date(2025, 11, 15), date(2025, 12, 31), date(2026, 2, 5),
          date(2026, 5, 15)], [26, 100, 130, 150, 200]),
    ]
    fig, ax = plt.subplots(figsize=(11.4, 5.3))
    for name, color, mk, off, x, y in series:
        ax.plot(x, y, marker=mk, lw=2.4, ms=5.5, color=color, label=name)
        ax.annotate(f"${y[-1]}M", (x[-1], y[-1]), textcoords="offset points",
                    xytext=off, fontsize=10.5, fontweight="bold", color=color,
                    ha="right" if off[0] < 0 else "left")
    ax.set_ylabel("ARR ($M, annualized)")
    ax.set_ylim(0, 340)
    ax.set_xlim(date(2024, 11, 1), date(2026, 11, 15))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b '%y"))
    ax.legend(loc="upper left", frameon=False, fontsize=11.5)
    return save(fig, "verticals.png")


# ---------------------------------------------------------------------------
# Chart 5: Agentforce ARR + tokens (two panels)
# ---------------------------------------------------------------------------
def chart_agentforce():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.4, 5.0))
    q = ["Q4 FY25\n(Jan '25)*", "Q1 FY26\n(Apr '25)*", "Q4 FY26\n(Jan '26)", "Q1 FY27\n(Apr '26)"]
    arr = [297, 393, 800, 1200]
    colors = ["#9DC3E6", "#9DC3E6", BLUE, NAVY]
    bars = ax1.bar(q, arr, color=colors, width=0.62)
    for b, v in zip(bars, arr):
        ax1.annotate(f"${v:,}M", (b.get_x() + b.get_width() / 2, v), xytext=(0, 5),
                     textcoords="offset points", ha="center", fontsize=11, fontweight="bold",
                     color=NAVY)
    ax1.set_title("Agentforce ARR ($M)  |  +205% YoY", fontsize=13, fontweight="bold", color=NAVY)
    ax1.set_ylim(0, 1400)
    ax1.tick_params(axis="x", labelsize=10)

    q2 = ["Q4 FY26\n(Jan '26)*", "Q1 FY27\n(Apr '26)"]
    tok = [11.3, 28.6]
    bars2 = ax2.bar(q2, tok, color=["#F4B183", ORANGE], width=0.5)
    for b, v in zip(bars2, tok):
        ax2.annotate(f"{v}T", (b.get_x() + b.get_width() / 2, v), xytext=(0, 5),
                     textcoords="offset points", ha="center", fontsize=11, fontweight="bold",
                     color="#833C00")
    ax2.set_title("Tokens processed per quarter (trillions)  |  +152% QoQ",
                  fontsize=13, fontweight="bold", color=NAVY)
    ax2.set_ylim(0, 34)
    ax2.tick_params(axis="x", labelsize=10)
    fig.text(0.01, -0.04, "* Derived from disclosed growth rates (+169%/+205% YoY ARR; +152% QoQ tokens). "
             "3.8B agentic work units delivered in Q1 FY27 (+111%).", fontsize=9.5, color=GRAY)
    fig.tight_layout(w_pad=3)
    return save(fig, "agentforce.png")


# ---------------------------------------------------------------------------
# Chart 6: Google monthly tokens (log)
# ---------------------------------------------------------------------------
def chart_tokens():
    fig, ax = plt.subplots(figsize=(11.4, 5.3))
    x = ["May 2024", "May 2025", "May 2026"]
    y = [9.7, 480, 3200]
    bars = ax.bar(x, y, color=["#9DC3E6", BLUE, NAVY], width=0.5)
    ax.set_yscale("log")
    for b, v, lb in zip(bars, y, ["9.7T", "480T  (~49x YoY)", "3.2 quadrillion  (~7x YoY)"]):
        ax.annotate(lb, (b.get_x() + b.get_width() / 2, v), xytext=(0, 6),
                    textcoords="offset points", ha="center", fontsize=12, fontweight="bold",
                    color=NAVY)
    ax.set_ylabel("Tokens per month (trillions, log scale)")
    ax.set_ylim(3, 9000)
    ax.set_yticks([10, 100, 1000])
    ax.set_yticklabels(["10T", "100T", "1,000T"])
    ax.annotate("Reference points (monthly volume):\n"
                "  Gemini models: ~1,000T (Q2 '26, +120% in 6 months)\n"
                "  OpenAI API: ~260T (Oct '25, excl. ChatGPT)\n"
                "  OpenRouter: ~85T   |   Salesforce Agentforce: ~9.5T",
                xy=(0.02, 0.96), xycoords="axes fraction", va="top",
                fontsize=11, color=GRAY,
                bbox=dict(boxstyle="round,pad=0.5", fc=LIGHT, ec="#BFBFBF"))
    return save(fig, "tokens.png")


# ---------------------------------------------------------------------------
# Chart 7: Microsoft adoption (two panels)
# ---------------------------------------------------------------------------
def chart_microsoft():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.4, 5.0))
    q = ["Q2 FY26\n(Dec '25)", "Q3 FY26\n(Mar '26)", "Q4 FY26\n(Jun '26)"]
    seats = [15, 20, 30]
    bars = ax1.bar(q, seats, color=["#9DC3E6", BLUE, NAVY], width=0.55)
    for b, v in zip(bars, seats):
        ax1.annotate(f"{v}M+" if v == 30 else f"{v}M", (b.get_x() + b.get_width() / 2, v),
                     xytext=(0, 5), textcoords="offset points", ha="center",
                     fontsize=12, fontweight="bold", color=NAVY)
    ax1.set_title("M365 Copilot paid seats (millions)\n2x in six months", fontsize=13,
                  fontweight="bold", color=NAVY)
    ax1.set_ylim(0, 36)
    ax1.tick_params(axis="x", labelsize=10)

    labels = ["Paid subs\nFY24", "Paid subs\nJan '26", "Users\nJul '25", "Users\nJul '26"]
    vals = [1.8, 4.7, 20, 50]
    colors = ["#F4B183", ORANGE, "#C5E0B4", GREEN]
    bars2 = ax2.bar(labels, vals, color=colors, width=0.6)
    for b, v in zip(bars2, vals):
        ax2.annotate(f"{v}M", (b.get_x() + b.get_width() / 2, v), xytext=(0, 5),
                     textcoords="offset points", ha="center", fontsize=12,
                     fontweight="bold", color=NAVY)
    ax2.set_title("GitHub Copilot (millions)\npaid +75% YoY; revenue +60% QoQ in Q4",
                  fontsize=13, fontweight="bold", color=NAVY)
    ax2.set_ylim(0, 58)
    ax2.tick_params(axis="x", labelsize=10)
    fig.tight_layout(w_pad=3)
    return save(fig, "microsoft.png")


# ---------------------------------------------------------------------------
# Chart 8: GitHub commits by agent (two panels)
# ---------------------------------------------------------------------------
def chart_github():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.4, 5.0),
                                   gridspec_kw={"width_ratios": [1.35, 1]})
    agents = ["Codex", "OpenHands", "CodeRabbit", "Devin", "Aider", "Jules",
              "Replit", "Claude Code"]
    commits = [731, 25676, 34940, 98493, 196132, 215804, 314779, 886122]
    colors = [BLUE, GRAY, GREEN, GRAY, GRAY, GRAY, ORANGE, NAVY]
    bars = ax1.barh(agents, commits, color=colors, height=0.62)
    ax1.set_xscale("log")
    ax1.set_xlim(300, 4e6)
    fmt = lambda v: f"{v/1e3:.0f}K" if v >= 10000 else f"{v:,}"
    for b, v in zip(bars, commits):
        ax1.annotate(fmt(v), (v, b.get_y() + b.get_height() / 2), xytext=(5, 0),
                     textcoords="offset points", va="center", fontsize=10.5,
                     fontweight="bold", color=NAVY)
    ax1.set_title("Attributed commits, Dec '25 - Apr '26 (log scale)\nClaude Code = 50% of all AI-attributed commits",
                  fontsize=12.5, fontweight="bold", color=NAVY)
    ax1.annotate("Codex is PR-native: 814,522 PRs\nbut near-zero commit traces",
                 xy=(0.42, 0.06), xycoords="axes fraction", fontsize=10, color=BLUE)

    owners = ["Devin +\nWindsurf", "GitHub\nCopilot", "Codex", "Claude\nCode"]
    units = [0.256, 3.8, 5.1, 12.5]
    colors2 = [GRAY, "#9DC3E6", BLUE, NAVY]
    bars2 = ax2.bar(owners, units, color=colors2, width=0.6)
    for b, v in zip(bars2, units):
        ax2.annotate(f"{v}M", (b.get_x() + b.get_width() / 2, v), xytext=(0, 5),
                     textcoords="offset points", ha="center", fontsize=11.5,
                     fontweight="bold", color=NAVY)
    ax2.set_title("All-time attributed units, commits + PRs\n(millions, Aug '26)",
                  fontsize=12.5, fontweight="bold", color=NAVY)
    ax2.set_ylim(0, 14.5)
    ax2.tick_params(axis="x", labelsize=10)
    fig.text(0.01, -0.04, "Attribution undercounts silent agents: Cursor and IDE Copilot do not sign commits; "
             "Codex's PR marker stopped appearing Nov '25. Copilot SWE agent: 1.13M commits cumulative through Oct '25 (different window).",
             fontsize=9.5, color=GRAY)
    fig.tight_layout(w_pad=3)
    return save(fig, "github.png")


# ---------------------------------------------------------------------------
# Chart 9: Reasoning tokens by agent workload (two panels)
# ---------------------------------------------------------------------------
def chart_reasoning():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.4, 5.0))

    cats = ["Q1 2025", "Late 2025"]
    share = [2, 52]
    bars = ax1.bar(cats, share, color=["#9DC3E6", NAVY], width=0.45)
    for b, v, lb in zip(bars, share, ["~0-2%", ">50%"]):
        ax1.annotate(lb, (b.get_x() + b.get_width() / 2, v), xytext=(0, 5),
                     textcoords="offset points", ha="center", fontsize=13,
                     fontweight="bold", color=NAVY)
    ax1.set_title("Share of OpenRouter tokens served by\nreasoning-optimized models",
                  fontsize=12.5, fontweight="bold", color=NAVY)
    ax1.set_ylim(0, 62)
    ax1.set_ylabel("% of routed tokens")
    ax1.annotate("Completion tokens per request ~3x,\nmostly reasoning tokens;\ntop reasoning models are code-oriented\n(Grok Code Fast 1, Gemini 2.5 Pro/Flash)",
                 xy=(0.04, 0.95), xycoords="axes fraction", va="top", fontsize=10,
                 color=GRAY)

    labels = ["Agentic SWE task\n(token type)", "Initial coding\n(stage)", "Code review\n(stage)",
              "Documentation\n(stage)"]
    reasoning = [21.6, 35, 23, 6]
    other = [100 - r for r in reasoning]
    x = range(len(labels))
    ax2.bar(x, reasoning, color=ORANGE, width=0.58, label="Reasoning tokens")
    ax2.bar(x, other, bottom=reasoning, color="#D9D9D9", width=0.58, label="Input + output tokens")
    for i, r in enumerate(reasoning):
        ax2.annotate(f"{r}%", (i, r), xytext=(0, 3), textcoords="offset points",
                     ha="center", fontsize=11.5, fontweight="bold", color="#833C00")
    ax2.set_xticks(list(x))
    ax2.set_xticklabels(labels, fontsize=9.5)
    ax2.set_title("Reasoning share of agentic coding tokens\n(ChatDev + GPT-5, 30 real tasks)",
                  fontsize=12.5, fontweight="bold", color=NAVY)
    ax2.set_ylim(0, 112)
    ax2.legend(loc="upper right", frameon=False, fontsize=10)
    fig.text(0.01, -0.04, "Full task mix: input 53.9% / output 24.4% / reasoning 21.6%; 17K-40K reasoning tokens per task. "
             "Code review alone consumes 59.4% of all task tokens; reads/navigation are 76% of single-agent loop tokens.",
             fontsize=9.5, color=GRAY)
    fig.tight_layout(w_pad=3)
    return save(fig, "reasoning.png")


# ---------------------------------------------------------------------------
# PPTX assembly
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

NAVY_RGB = RGBColor(0x1F, 0x38, 0x64)
GRAY_RGB = RGBColor(0x7F, 0x7F, 0x7F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_bar(slide, color=NAVY_RGB, top=Inches(0), height=Inches(0.16)):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, prs.slide_width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size, bold=False,
             color=NAVY_RGB, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb


def chart_slide(title, subtitle, img, source):
    slide = prs.slides.add_slide(BLANK)
    add_bar(slide)
    add_text(slide, Inches(0.55), Inches(0.32), Inches(12.3), Inches(0.6), title, 26, bold=True)
    add_text(slide, Inches(0.55), Inches(0.92), Inches(12.3), Inches(0.4), subtitle, 13,
             color=GRAY_RGB)
    from PIL import Image
    with Image.open(img) as im:
        w, h = im.size
    disp_w = Inches(12.1)
    disp_h = disp_w * h / w
    max_h = Inches(5.35)
    if disp_h > max_h:
        disp_h = max_h
        disp_w = disp_h * w / h
    slide.shapes.add_picture(img, (prs.slide_width - disp_w) / 2, Inches(1.42),
                             width=disp_w, height=disp_h)
    add_text(slide, Inches(0.55), Inches(7.05), Inches(12.3), Inches(0.35),
             f"Source: {source}  |  # refs -> Sources sheet, agentic_ai_revenue_token_usage_2026.xlsx",
             10, color=GRAY_RGB)
    return slide


# --- Slide 1: title ---
slide = prs.slides.add_slide(BLANK)
from pptx.enum.shapes import MSO_SHAPE
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY_RGB
bg.line.fill.background()
add_text(slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4),
         "Agentic AI: Revenue & Token Usage", 44, bold=True, color=WHITE)
add_text(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.7),
         "The most useful time series across sectors", 24, color=RGBColor(0xBD, 0xD7, 0xEE))
add_text(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
         f"Compiled {AS_OF}  |  Data: agentic_ai_revenue_token_usage_2026.xlsx "
         "(Equity-research repo)", 14, color=RGBColor(0xBD, 0xD7, 0xEE))

# --- Slide 2: key takeaways ---
slide = prs.slides.add_slide(BLANK)
add_bar(slide)
add_text(slide, Inches(0.55), Inches(0.32), Inches(12.3), Inches(0.6),
         "Key takeaways from the time series", 26, bold=True)
takeaways = [
    ("Coding agents are the sector", "~$30.4B of ~$35.0B tracked agentic revenue (86.9%). "
     "Claude Code scaled $1B -> $15.1B tracked ARR in ~9 months; Codex ~$0 -> $8.8B in ~7 months."),
    ("Model providers repriced in months", "Anthropic run-rate $14B (Feb '26) -> $65B (Jul '26); "
     "OpenAI $20B (Dec '25) -> $40B+ (Aug '26). Enterprise/agents are now the majority of both mixes."),
    ("Exits followed the ramps", "SpaceX agreed to buy Cursor's parent Anysphere for $60B (Jun '26) "
     "at ~15x its $4B ARR, 31 months after Cursor's first $1M."),
    ("Verticals compound from smaller bases", "Harvey, OpenEvidence, Glean and Sierra each reached "
     "$200-300M ARR, growing 2-6x in 12-18 months on outcome- and usage-based pricing."),
    ("Enterprise platforms hit the S-curve", "Agentforce ARR +205% YoY to $1.2B; tokens +152% QoQ "
     "to 28.6T/quarter. M365 Copilot seats doubled to 30M+ in six months."),
    ("Token demand is the leading indicator", "Google: 9.7T -> 3.2 quadrillion tokens/month in two "
     "years (~330x). Agentic requests burn ~15x the tokens of chat; coding agents ~1,000x."),
]
top = 1.15
for head, body in takeaways:
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.9), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = head + " — "
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = NAVY_RGB
    r2 = p.add_run(); r2.text = body
    r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    top += 0.98

# --- Chart slides ---
chart_slide(
    "Coding agents: Claude Code vs Codex tracked ARR",
    "Third-party tracked estimates (TickerTrends), anchored to Anthropic's $2.5B disclosure; weekly points Jul-Aug 2026",
    chart_coding_arr(), "TickerTrends (Aug 10, 2026); Anthropic Series G; Sacra  [#3, 4, 5]")
chart_slide(
    "Cursor: $1M to $4B ARR in 31 months, then a $60B exit",
    "Fastest-measured B2B software ramp; ~65% of revenue now enterprise contracts (~$2.6B)",
    chart_cursor(), "Anysphere disclosures; Bloomberg; Sacra; SpaceX deal announcement  [#1, 2, 40]")
chart_slide(
    "Model providers: agentic demand repriced both leaders",
    "Annualized revenue run-rate; Anthropic's Q2 2026 quarterly revenue overtook OpenAI's for the first time",
    chart_providers(), "Bloomberg (Aug 13, 2026); CNBC; Forkast/Bloomberg internal docs; Sacra  [#6, 29, 30, 31]")
chart_slide(
    "Vertical & knowledge agents: the $200-300M ARR club",
    "Legal, healthcare, enterprise knowledge and customer service all crossed from pilots to production",
    chart_verticals(), "AgentMarketCap; ARR Club; TechCrunch; Sacra  [#9-14, 35, 39]")
chart_slide(
    "Salesforce Agentforce: revenue and tokens compounding together",
    "First enterprise platform agent to $1B+ ARR; token volume growing faster than revenue",
    chart_agentforce(), "Salesforce Q1 FY2027 results via CX Today / CXM  [#17, 18]")
chart_slide(
    "Microsoft: agent adoption at platform scale",
    "M365 Copilot seats and GitHub Copilot users/subscribers; Agent 365 registered ~40M agents in its first two months",
    chart_microsoft(), "Microsoft FY26 Q2/Q3/Q4 earnings; Axis Intelligence  [#7, 36, 37]")
chart_slide(
    "GitHub commits: which agents actually ship code",
    "Multi-method census of 180M repositories (WoC V2604) plus all-time attribution tracking; channels differ by agent",
    chart_github(), "arXiv 2606.24429 census; Amplifying.ai; JetBrains survey  [#41, 42, 43]")
chart_slide(
    "Reasoning tokens: now the default mode for agentic work",
    "Reasoning-optimized models serve the majority of routed tokens; ~22% of agentic coding task tokens are reasoning",
    chart_reasoning(), "OpenRouter State of AI; arXiv 2601.14470 Tokenomics study  [#26, 44, 45]")
chart_slide(
    "Token throughput: the clearest demand signal",
    "Google platform-wide monthly tokens, ~330x in two years; agentic workloads are the driver (do not sum reference points)",
    chart_tokens(), "Google I/O 2026; Q2 2026 earnings; OpenAI DevDay; OpenRouter; Salesforce  [#17, 21, 22, 23, 26]")

# --- Final slide: notes & caveats ---
slide = prs.slides.add_slide(BLANK)
add_bar(slide)
add_text(slide, Inches(0.55), Inches(0.32), Inches(12.3), Inches(0.6),
         "Notes, methodology & caveats", 26, bold=True)
notes = [
    "All series are assembled from the Growth trajectory and token columns of agentic_ai_revenue_token_usage_2026.xlsx; "
    "source numbers [#] map to its Sources sheet (40 references).",
    "Confidence varies by series: company-disclosed (Agentforce, M365 Copilot seats, Google tokens, Cursor to $2B) vs "
    "third-party tracked (Claude Code, Codex - TickerTrends) vs analyst estimates (Sacra: Sierra, Abridge, Cursor $4B).",
    "Annualized run-rates extrapolate the latest period x12 and can move sharply; several curves will moderate as they lap larger bases.",
    "Agentforce bars marked * are derived from disclosed growth rates, not directly reported figures.",
    "Double-counting warning: provider totals (OpenAI, Anthropic) include their agent products (Codex, Claude Code) and the API "
    "revenue paid to them by application-layer companies (Cursor, Replit, Lovable).",
    "Token metrics are heterogeneous (platform-wide vs API-only vs router-visible) and are shown for scale, not for summation.",
    "Regenerate this deck: python3 scripts/build_presentation.py",
]
top = 1.2
for n in notes:
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.9), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "-  " + n
    p.font.size = Pt(13.5)
    p.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    top += 0.82

prs.save("agentic_ai_time_series_2026.pptx")
print(f"Wrote agentic_ai_time_series_2026.pptx ({len(prs.slides.__iter__.__self__._sldIdLst)} slides); charts in {TMP}")
