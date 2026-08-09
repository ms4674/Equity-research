#!/usr/bin/env python3
"""Build the sellside research deck:

    "Powering the AI Buildout — Data Center Lead Times, Bottlenecks
     and Our Non-Consensus Views" (August 2026).

Generates Datacenter_Sellside_Deck_Aug2026.pptx in this directory.

Usage:
    pip install -r requirements.txt
    python build_deck.py
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- palette ---
NAVY = RGBColor(0x0B, 0x1D, 0x3A)      # primary brand
NAVY2 = RGBColor(0x1A, 0x33, 0x57)     # panel navy
GOLD = RGBColor(0xC2, 0x90, 0x3A)      # accent
INK = RGBColor(0x22, 0x2A, 0x35)       # body text
MUTED = RGBColor(0x66, 0x70, 0x7D)     # secondary text
LIGHT = RGBColor(0xF1, 0xF4, 0xF8)     # panel background
MID = RGBColor(0xD8, 0xDE, 0xE6)       # table banding
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xA8, 0x39, 0x2E)       # risk / non-consensus tag
GREEN = RGBColor(0x1E, 0x6F, 0x4A)     # positive
GRAY = RGBColor(0x9A, 0xA4, 0xB0)      # chart secondary series

FONT = "Calibri"
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

FOOTER_LEFT = "Data Center Infrastructure  |  Sector Research  |  August 2026"
FOOTER_NOTE = "Illustrative sellside research template. Not investment advice."

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

_page = [0]


# ---------------------------------------------------------------- helpers ---
def new_slide():
    _page[0] += 1
    return prs.slides.add_slide(BLANK), _page[0]


def _set_text(tf, runs_by_para, default_size=11, default_color=INK,
              align=PP_ALIGN.LEFT, space_after=4, line_spacing=1.0):
    """runs_by_para: list of paragraphs; each is a list of
    (text, bold, size, color) tuples (size/color may be None)."""
    tf.word_wrap = True
    for i, runs in enumerate(runs_by_para):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for text, bold, size, color in runs:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.bold = bold
            r.font.size = Pt(size if size else default_size)
            r.font.color.rgb = color if color else default_color


def add_box(slide, x, y, w, h, fill=None, line=None, round_=False, radius=0.06):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if round_:
        sp.adjustments[0] = radius
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.text_frame.word_wrap = True
    return sp


def add_text(slide, x, y, w, h, paras, size=11, color=INK, align=PP_ALIGN.LEFT,
             space_after=4, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.vertical_anchor = anchor
    _set_text(tb.text_frame, paras, default_size=size, default_color=color,
              align=align, space_after=space_after, line_spacing=line_spacing)
    return tb


def bullets(slide, x, y, w, h, items, size=11, space_after=6, line_spacing=1.05):
    """items: list of (lead, rest) — lead rendered bold navy, rest regular."""
    paras = []
    for lead, rest in items:
        runs = []
        if lead:
            runs.append((f"{lead}  ", True, size, NAVY))
        runs.append((rest, False, size, INK))
        paras.append(runs)
    return add_text(slide, x, y, w, h, paras, size=size,
                    space_after=space_after, line_spacing=line_spacing)


def header(slide, page, kicker, title, sub=None):
    add_box(slide, 0, 0, SLIDE_W, Inches(0.06), fill=GOLD)
    add_text(slide, Inches(0.55), Inches(0.28), Inches(11.0), Inches(0.3),
             [[(kicker.upper(), True, 10.5, GOLD)]], space_after=0)
    add_text(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.62),
             [[(title, True, 23, NAVY)]], space_after=0)
    if sub:
        add_text(slide, Inches(0.55), Inches(1.13), Inches(12.2), Inches(0.32),
                 [[(sub, False, 12, MUTED)]], space_after=0)
    # footer
    add_box(slide, 0, Inches(7.14), SLIDE_W, Pt(0.8), fill=MID)
    add_text(slide, Inches(0.55), Inches(7.18), Inches(8.6), Inches(0.28),
             [[(f"{FOOTER_LEFT}   —   {FOOTER_NOTE}", False, 8, MUTED)]],
             space_after=0)
    add_text(slide, Inches(12.35), Inches(7.18), Inches(0.5), Inches(0.28),
             [[(str(page), True, 9, MUTED)]], align=PP_ALIGN.RIGHT, space_after=0)


def source(slide, text, y=Inches(6.82)):
    add_text(slide, Inches(0.55), y, Inches(12.2), Inches(0.3),
             [[("Source: ", True, 8, MUTED), (text, False, 8, MUTED)]],
             space_after=0)


def styled_table(slide, x, y, w, rows_data, col_widths, header_fill=NAVY,
                 font_size=9.5, header_size=9.5, row_h=0.34, header_h=0.36,
                 align_map=None, highlight_rows=None):
    """rows_data[0] is the header row. col_widths in inches (sums to w)."""
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    total_h = Inches(header_h + row_h * (n_rows - 1))
    gfx = slide.shapes.add_table(n_rows, n_cols, x, y, w, total_h)
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    for c, cw in enumerate(col_widths):
        tbl.columns[c].width = Inches(cw)
    tbl.rows[0].height = Inches(header_h)
    for r in range(1, n_rows):
        tbl.rows[r].height = Inches(row_h)
    highlight_rows = highlight_rows or {}
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
            elif r in highlight_rows:
                cell.fill.fore_color.rgb = highlight_rows[r]
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if align_map and c in align_map:
                p.alignment = align_map[c]
            run = p.add_run()
            if isinstance(val, tuple):
                text, bold, color = val
            else:
                text, bold, color = val, False, None
            run.text = text
            run.font.name = FONT
            run.font.size = Pt(header_size if r == 0 else font_size)
            run.font.bold = True if r == 0 else bold
            run.font.color.rgb = (WHITE if r == 0
                                  else (color if color else INK))
    return tbl


def style_chart(chart, size=9):
    chart.font.name = FONT
    chart.font.size = Pt(size)
    chart.font.color.rgb = INK


# ------------------------------------------------------------ slide 1: title
slide, page = new_slide()
add_box(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_box(slide, Inches(0.9), Inches(1.55), Inches(1.7), Pt(3), fill=GOLD)
add_text(slide, Inches(0.9), Inches(1.05), Inches(10), Inches(0.4),
         [[("SECTOR RESEARCH  |  TECHNOLOGY INFRASTRUCTURE", True, 12, GOLD)]])
add_text(slide, Inches(0.9), Inches(1.95), Inches(11.6), Inches(1.9),
         [[("Powering the AI Buildout", True, 44, WHITE)]], space_after=0)
add_text(slide, Inches(0.9), Inches(3.05), Inches(11.4), Inches(1.0),
         [[("Data center lead times, bottlenecks, and where we differ from "
            "consensus", False, 22, RGBColor(0xC9, 0xD4, 0xE3))]])
add_text(slide, Inches(0.9), Inches(4.55), Inches(11), Inches(1.3), [
    [("EQUITY RESEARCH  —  DATA CENTER INFRASTRUCTURE", True, 12, WHITE)],
    [("August 2026  |  Initiation framework deck", False, 12,
      RGBColor(0xC9, 0xD4, 0xE3))],
], space_after=4)
add_text(slide, Inches(0.9), Inches(6.65), Inches(11.6), Inches(0.6),
         [[("Illustrative research deck. Figures compiled from public sources "
            "as of August 2026 and may be approximate. Not investment advice.",
            False, 9, RGBColor(0x8A, 0x99, 0xAE))]])

# -------------------------------------------------- slide 2: executive summary
slide, page = new_slide()
header(slide, page, "Executive summary", "Five key calls: demand is not the "
       "question — delivery is", "The constraint has moved downstream from "
       "chips to power, equipment, memory and labor")
calls = [
    ("1. Capex is committed; the binding constraints sit downstream.",
     "Big-5 hyperscaler 2026 capex guidance totals ~$775–800bn (+~65% y/y). The "
     "gating factors are now energized megawatts, electrical equipment, HBM/"
     "packaging and skilled labor — not GPU compute die."),
    ("2. Lead times are the sector's best leading indicator.",
     "Track them like commodity inventories: large power transformers ~128 wks "
     "(specialty 3–5 yrs), GSUs 160+ wks, grid interconnection median ~5 yrs, "
     "CoWoS booked through mid-2027, 2026 HBM4 output 100% sold out."),
    ("3. Non-consensus: 2027 growth is set by megawatts energized, not GPUs "
     "shipped.", "Consensus still models accelerator supply as the gate. "
     "Speed-to-power — behind-the-meter generation, flexible/curtailable "
     "interconnection, queue position — is the scarce capability."),
    ("4. Non-consensus: the depreciation wall is the dated, mechanical risk.",
     "Trailing-4Q hyperscaler capex of ~$434bn versus only ~$149bn of "
     "recognized D&A. Useful-life disclosures (10-Ks, Jan–Feb) are the leading "
     "indicator for a 2027–28 capex moderation."),
    ("5. Non-consensus: HBM pricing power peaks before HBM capacity does; "
     "circularity means the cycle can crack without AI failing.",
     "Spec downgrades (Rubin Ultra stack reductions) are early demand "
     "destruction; vendor-financing loops transmit any single buyer's pullback "
     "across the ecosystem."),
]
y = 1.62
for lead, rest in calls:
    add_box(slide, Inches(0.55), Inches(y), Inches(12.23), Inches(0.94),
            fill=LIGHT, round_=True, radius=0.10)
    add_box(slide, Inches(0.55), Inches(y), Inches(0.07), Inches(0.94),
            fill=GOLD)
    add_text(slide, Inches(0.78), Inches(y + 0.07), Inches(11.85),
             Inches(0.84),
             [[(lead + "  ", True, 11.5, NAVY), (rest, False, 10.5, INK)]],
             space_after=0, line_spacing=1.0)
    y += 1.04
add_text(slide, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.3),
         [[("Positioning: own the toll collectors (grid-equipment oligopoly, "
            "memory near term, power-secured operators); avoid those bearing "
            "cost inflation without pricing power.", True, 10, RED)]])

# ------------------------------------------------ slide 3: capex backdrop
slide, page = new_slide()
header(slide, page, "Market backdrop", "2026: the ~$800bn year — guidance "
       "raised twice, still supply-constrained",
       "Every major buyer raised 2026 capex guidance intra-year; all cite "
       "capacity shortages, not demand, as the limiter")

chart_data = CategoryChartData()
chart_data.categories = ["Amazon", "Alphabet", "Microsoft", "Meta", "Oracle"]
chart_data.add_series("2025A", (131.8, 91.4, 100.6, 69.7, 21.2))
chart_data.add_series("2026E (guide midpoint)", (220, 200, 175, 137.5, 55.7))
gframe = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.65),
    Inches(6.15), Inches(4.55), chart_data)
chart = gframe.chart
style_chart(chart)
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.plots[0].gap_width = 90
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = GRAY
chart.series[1].format.fill.solid()
chart.series[1].format.fill.fore_color.rgb = NAVY
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(8.5)
plot.data_labels.number_format = "#,##0"
plot.data_labels.number_format_is_linked = False
chart.value_axis.has_major_gridlines = True
chart.value_axis.tick_labels.font.size = Pt(8.5)
chart.category_axis.tick_labels.font.size = Pt(9.5)
add_text(slide, Inches(0.55), Inches(6.25), Inches(6.2), Inches(0.3),
         [[("Capex, $bn. 2026E = latest company guidance (Jul 2026 calls).",
            False, 8.5, MUTED)]])

rows = [
    ["Buyer", "2025A", "2026 guide", "What changed in July 2026"],
    ["Amazon", "$132bn", "~$220bn", "Raised from ~$200bn; memory costs cited; "
     "short capacity into 2027"],
    ["Alphabet", "$91bn", "$195–205bn", "Raised from $180–190bn; \"supply-"
     "constrained environment\""],
    ["Microsoft", "$101bn", "~$175bn", "Lease reclass + useful life 15→25 yrs "
     "(DC/office); intent unchanged"],
    ["Meta", "$70bn", "$130–145bn", "Floor raised; explicitly cites memory "
     "price inflation"],
    ["Oracle", "$21bn", "~$56bn", "Capex/revenue ~83% — most levered to the "
     "build"],
    [("Big-5", True, None), ("~$429bn", True, None),
     ("~$775–800bn", True, None), ("+~65% y/y; ~75% AI-specific (~$545bn)",
                                   True, None)],
]
styled_table(slide, Inches(7.0), Inches(1.75), Inches(5.8), rows,
             [0.95, 0.78, 1.05, 3.02], font_size=8.8, header_size=9.5,
             row_h=0.62, header_h=0.34,
             highlight_rows={6: RGBColor(0xEA, 0xE2, 0xD0)})
add_text(slide, Inches(7.0), Inches(6.05), Inches(5.8), Inches(0.75),
         [[("Read-through:  ", True, 10, RED),
           ("guidance revisions are now driven by input costs (memory) and "
            "delivery schedules (TSMC, SK Hynix, power), not by demand "
            "planning. The print depends on suppliers, not CFOs.",
            False, 10, INK)]])
source(slide, "Company earnings releases and calls (Jul 2026); Axis "
       "Intelligence AI Capex Tracker; AL Capital Advisory compilation.")

# ------------------------------------------------ slide 4: where money goes
slide, page = new_slide()
header(slide, page, "Market backdrop", "Where the money goes: silicon is "
       "half the bill; power sets the date",
       "Estimated 2026 hyperscaler capex mix — a 1GW AI campus runs $5–10bn "
       "and draws a mid-sized city's load")

mix = CategoryChartData()
mix.categories = ["GPUs & AI accelerators", "Server & rack infrastructure",
                  "Data center shell", "Power & cooling",
                  "Networking & fiber", "Land, leases, other"]
mix.add_series("Share of 2026E capex (%)", (44, 18, 13, 11, 7, 7))
gframe = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(1.7),
    Inches(6.3), Inches(4.4), mix)
chart = gframe.chart
style_chart(chart)
chart.has_legend = False
chart.plots[0].gap_width = 70
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = NAVY
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(9)
plot.data_labels.number_format = '0"%"'
plot.data_labels.number_format_is_linked = False
chart.value_axis.has_major_gridlines = False
chart.value_axis.tick_labels.font.size = Pt(8.5)
chart.category_axis.tick_labels.font.size = Pt(9.5)
add_text(slide, Inches(0.55), Inches(6.2), Inches(6.3), Inches(0.3),
         [[("Approximate mix, Big-4; ~$310bn of accelerators at midpoint.",
            False, 8.5, MUTED)]])

bullets(slide, Inches(7.1), Inches(1.8), Inches(5.7), Inches(4.9), [
    ("Silicon is the largest check…", "GPUs/accelerators ~44% of spend (~$310bn"
     "), plus ~18% for the server and rack infrastructure around them."),
    ("…but the smallest schedule risk.", "Compute die supply has largely "
     "normalized. The 11% spent on power & cooling and the grid connection "
     "behind it set the energization date — and therefore the revenue date."),
    ("Campus economics have step-changed.", "A single large AI campus is now a "
     "$5–10bn, 500MW–1GW+ commitment. Microsoft added ~1GW in a single quarter "
     "and plans to roughly double total capacity in two years."),
    ("Electrical content per MW is rising.", "US data-center electrical-"
     "equipment spend is forecast to grow from ~$20bn (2025) to ~$65bn (2030) "
     "— data centers could be ~40% of the power-equipment market by 2030, vs "
     "<2% in 2020."),
], size=10.5, space_after=10)
source(slide, "Analysis Atlas capex mix estimates; Wood Mackenzie; company "
       "disclosures; Value Add VC.")

# ------------------------------------------------ slide 5: lead-time monitor
slide, page = new_slide()
header(slide, page, "Lead-time monitor", "The lead-time board: what you order "
       "today arrives in 2027–2030",
       "Our tracking table for the physical inputs to a gigawatt — updated "
       "monthly; direction matters more than level")
rows = [
    ["Input", "Pre-boom norm", "Current (mid-2026)", "Trend",
     "What to know"],
    ["Grid interconnection (firm)", "~22 mo (2008 median)",
     "Median 61 mo request-to-COD; 4–7+ yrs in NoVA / Phoenix / Dallas",
     ("Worsening", True, RED),
     "~2,300 GW sits in US queues — more than installed US capacity"],
    ["Large power transformer", "24–30 mo",
     "~128 wks average; 3–5 yrs for large/specialty units",
     ("Worsening", True, RED),
     "Oligopoly (Hitachi, Siemens, GEV, ABB) at 48–60 mo; new lines help "
     "2027–28"],
    ["GSU transformer", "~50–80 wks", "144–160+ wks (Q1-26)",
     ("Worsening", True, RED),
     "Demand +274% since 2019 (Wood Mackenzie); winding labor is the floor"],
    ["HV circuit breaker", "77 wks (2023)", "~125 wks",
     ("Worsening", True, RED), "Prices +4–10% expected over the next year"],
    ["MV switchgear", "A few months", "AIS 80–110 wks; GIS 90–130 wks",
     ("Tight", True, RED), "Effectively sold out into 2028"],
    ["Heavy-frame gas turbine", "~2–3 yrs (plant)",
     "Slots booked to end of decade; new CCGT ~5 yrs (vs 3.5 in 2023)",
     ("Worsening", True, RED),
     "GE Vernova backlog 100 GW, guiding 110+ GW; CCGT costs +49%"],
    ["HV cable (subsea/UG)", "~1–2 yrs", "Producers fully booked through 2029",
     ("Tight", True, RED), "Prysmian / Nexans / NKT at capacity"],
    ["Data-center GPU (reserved)", "8–16 wks", "36–52 wks",
     ("Stable-tight", True, GOLD),
     "H100/H200 contract pricing +~40% Oct-25→Mar-26 — driven by memory"],
    ["CoWoS packaging slot", "n/a", "52–78 wks; allocated through mid-2027",
     ("Tight", True, GOLD),
     "Capacity doubling to 120–130k wpm by Q4-26; NVIDIA holds ~60%"],
    ["HBM4 memory", "n/a", "2026 output 100% sold out; fill rates 60–70%",
     ("Worsening", True, RED),
     "Deficit seen through 2027; NVIDIA weighing stack downgrades"],
    ["Skilled labor", "n/a", "~439k US construction-worker shortfall",
     ("Structural", True, RED),
     "Coil winders, electricians, HV engineers; retirements outpace entry"],
]
styled_table(slide, Inches(0.55), Inches(1.6), Inches(12.23), rows,
             [1.72, 1.35, 3.30, 0.94, 4.92], font_size=8.3, header_size=9.5,
             row_h=0.435, header_h=0.32)
source(slide, "Wood Mackenzie; POWER Magazine; LBNL Queued Up (2026); "
       "BloombergNEF; TrendForce; company disclosures. Figures approximate.",
       y=Inches(6.88))

# ------------------------------------------------ slide 6: bottleneck stack
slide, page = new_slide()
header(slide, page, "Bottlenecks", "The constraint stack: rank-ordered, with "
       "time-to-relief",
       "\"The binding constraint is a transformer, not a chip\" — each layer "
       "must clear before the next one binds")
stack = [
    ("#1", "Energized power (interconnection + generation)",
     "Binding", RED,
     "Firm grid connection is the long pole: median ~5 yrs, hubs 4–7+. "
     "Buildings finish in 2–3 yrs and then wait. Relief: 2029+ (queue reform, "
     "behind-the-meter bridges)."),
    ("#2", "Electrical equipment (transformers, breakers, switchgear, "
     "turbines)", "Binding", RED,
     "Multi-year backlogs across the complex; announced factory expansions "
     "become material only in 2027–28 and are labor-limited."),
    ("#3", "HBM memory + CoWoS advanced packaging",
     "Binding thru 2027", RED,
     "Memory — not GPU die — is the silicon constraint. 2026 HBM4 sold out; "
     "CoWoS allocated through mid-2027 even as capacity doubles."),
    ("#4", "Skilled labor (winders, electricians, HV engineers)",
     "Chronic", GOLD,
     "~439k construction shortfall; ~1/3 of critical-infrastructure "
     "technical staff at or near retirement. No cyclical relief mechanism."),
    ("#5", "NOT the constraint: shells, steel, GPU die, land, fiber",
     "Clear", GREEN,
     "Shell construction runs 12–18 months and is rarely the long pole; "
     "compute-die supply has normalized. Don't underwrite scarcity here."),
]
y = 1.62
width = 12.23
for rank, name, tag, tagcolor, desc in stack:
    add_box(slide, Inches(0.55), Inches(y), Inches(width), Inches(0.97),
            fill=LIGHT, round_=True, radius=0.09)
    add_box(slide, Inches(0.55), Inches(y), Inches(0.62), Inches(0.97),
            fill=NAVY, round_=True, radius=0.09)
    add_text(slide, Inches(0.55), Inches(y + 0.26), Inches(0.62), Inches(0.45),
             [[(rank, True, 15, WHITE)]], align=PP_ALIGN.CENTER, space_after=0)
    add_text(slide, Inches(1.35), Inches(y + 0.08), Inches(8.6), Inches(0.34),
             [[(name, True, 12, NAVY)]], space_after=0)
    add_text(slide, Inches(1.35), Inches(y + 0.43), Inches(11.2), Inches(0.5),
             [[(desc, False, 9.5, INK)]], space_after=0, line_spacing=0.98)
    tag_box = add_box(slide, Inches(10.55), Inches(y + 0.08), Inches(2.05),
                      Inches(0.3), fill=tagcolor, round_=True, radius=0.5)
    add_text(slide, Inches(10.55), Inches(y + 0.10), Inches(2.05),
             Inches(0.26), [[(tag.upper(), True, 9, WHITE)]],
             align=PP_ALIGN.CENTER, space_after=0)
    width -= 0.55
    y += 1.07
source(slide, "Artifipedia grid-interconnection analysis; LBNL; TrendForce; "
       "Allianz Research; team assessment.", y=Inches(6.98))

# ------------------------------------------------ slide 7: power deep dive
slide, page = new_slide()
header(slide, page, "Deep dive — power", "Speed-to-power: built in year 3, "
       "energized in year 5–7",
       "The interconnection queue — not the crane — sets the schedule, and "
       "every month in the gap is idle capex")

qchart = CategoryChartData()
qchart.categories = ["2008", "2015", "2025"]
qchart.add_series("Median months, interconnection request → commercial "
                  "operation", (22, 36, 61))
gframe = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.7),
    Inches(5.4), Inches(3.6), qchart)
chart = gframe.chart
style_chart(chart)
chart.has_legend = False
chart.plots[0].gap_width = 80
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = NAVY
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(10)
chart.value_axis.has_major_gridlines = False
chart.value_axis.tick_labels.font.size = Pt(8.5)
chart.category_axis.tick_labels.font.size = Pt(10)
add_text(slide, Inches(0.55), Inches(5.35), Inches(5.4), Inches(0.55),
         [[("Median months from interconnection request to commercial "
            "operation, US (LBNL Queued Up 2026). Studies alone: median 45 "
            "months in 2025.", False, 8.5, MUTED)]])

bullets(slide, Inches(6.35), Inches(1.7), Inches(6.45), Inches(3.5), [
    ("The queue is the scarcest asset.", "~2,300 GW of generation and storage "
     "waits in US queues — more than total installed US capacity. Queue "
     "position belongs to the utility, not the operator: capital, GPUs and "
     "permits don't move you forward."),
    ("Demand shock meets aging grid.", "NERC's 10-year forecast added +224 GW "
     "of summer peak demand (+24% vs prior year), driven by data-center load; "
     "NERC has issued a rare Level 3 alert on DC load oscillations."),
    ("Half the 2026 pipeline is at risk.", "Allianz Research estimates ~30% "
     "supply shortage in key power equipment and expects roughly half of the "
     "~12 GW of planned 2026 US data-center capacity to slip or cancel."),
    ("The bridge: behind-the-meter generation.", "On-site turbines/gensets "
     "convert a 5–7 yr utility wait into a 12–18 month equipment schedule; "
     "flexible/curtailable interconnection can unlock ~98–100 GW of new load "
     "at ~0.5% annual curtailment (Duke/Nicholas Institute)."),
], size=10, space_after=8)
add_box(slide, Inches(6.35), Inches(5.55), Inches(6.45), Inches(1.1),
        fill=NAVY, round_=True, radius=0.08)
add_text(slide, Inches(6.6), Inches(5.68), Inches(6.0), Inches(0.9),
         [[("Analyst take:  ", True, 10.5, GOLD),
           ("\"MW energized\" is the KPI that converts capex into revenue. We "
            "value operators on secured, dated power — not on announced "
            "megawatts.", False, 10.5, WHITE)]], line_spacing=1.05)
source(slide, "LBNL Queued Up 2026; NERC; Allianz Research; Duke Nicholas "
       "Institute; aidatacenterguide.com synthesis.")

# ------------------------------------ slide 8: electrical equipment deep dive
slide, page = new_slide()
header(slide, page, "Deep dive — electrical equipment", "Transformers & "
       "turbines: a supercycle meets a labor-limited oligopoly",
       "Lead times have doubled-to-tripled; announced capacity lands 2027–28 "
       "and does not shorten today's quotes")

lt = CategoryChartData()
lt.categories = ["Large power transformer", "GSU transformer",
                 "HV circuit breaker"]
lt.add_series("Pre-boom (~2021–23)", (65, 65, 77))
lt.add_series("Mid-2026", (128, 152, 125))
gframe = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(1.7),
    Inches(6.0), Inches(3.3), lt)
chart = gframe.chart
style_chart(chart)
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.plots[0].gap_width = 80
chart.series[0].format.fill.solid()
chart.series[0].format.fill.fore_color.rgb = GRAY
chart.series[1].format.fill.solid()
chart.series[1].format.fill.fore_color.rgb = NAVY
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(9)
chart.value_axis.has_major_gridlines = False
chart.value_axis.tick_labels.font.size = Pt(8.5)
chart.category_axis.tick_labels.font.size = Pt(9)
add_text(slide, Inches(0.55), Inches(5.05), Inches(6.0), Inches(0.5),
         [[("US average lead times, weeks. Specialty transformer units quoted "
            "at 3–5 years.", False, 8.5, MUTED)]])
add_box(slide, Inches(0.55), Inches(5.6), Inches(6.0), Inches(1.05),
        fill=LIGHT, round_=True, radius=0.08)
add_text(slide, Inches(0.75), Inches(5.7), Inches(5.65), Inches(0.9),
         [[("The real constraint is a missing trade:  ", True, 10, NAVY),
           ("transformer output is gated by a ~15,000-person US workforce "
            "that hand-winds copper coils — capex cannot compress a "
            "144-week GSU quote in 2026.", False, 10, INK)]],
         line_spacing=1.02)

bullets(slide, Inches(6.95), Inches(1.7), Inches(5.85), Inches(4.8), [
    ("Turbines are scarcer still.", "GE Vernova's gas-turbine backlog reached "
     "100 GW in Q1-26 (guiding ≥110 GW by year-end); Siemens Energy sold 194 "
     "turbines in 2025 vs 100 in 2024, and slots at major makers are booked "
     "into the next decade."),
    ("Plant economics inflating.", "New CCGT lead time is ~5 years (from 3.5 "
     "in 2023) with costs +49% (BloombergNEF). Gas-fired capacity needed now "
     "is effectively rationed by the turbine queue."),
    ("Capacity response is real but late.", "Siemens committed $1bn+ to US "
     "grid/turbine manufacturing; new transformer plants in Virginia and "
     "North Carolina come online 2027–28; GOES steel ramps mid-2028."),
    ("The TAM re-rating.", "US data-center electrical-equipment spend: ~$20bn "
     "(2025) → ~$65bn (2030E). Transformer prices seen +4–10% over the next "
     "year — pricing power sits with the oligopoly."),
], size=10, space_after=9)
source(slide, "GE Vernova and Siemens Energy disclosures; BloombergNEF; Wood "
       "Mackenzie; POWER Magazine; Manufacturing Mag.")

# ------------------------------------ slide 9: silicon & memory deep dive
slide, page = new_slide()
header(slide, page, "Deep dive — silicon & memory", "The silicon gate is "
       "memory and packaging — not the GPU compute die",
       "HBM — not the compute die — is the binding silicon constraint through "
       "2027; packaging slots are the second gate")
rows = [
    ["Layer", "Status (mid-2026)", "Implication"],
    ["HBM4 supply",
     "100% of 2026 output sold out on LTAs; customer fill rates 60–70%; "
     "deficit seen through 2027",
     "Pricing power with SK Hynix (~60% share), Samsung, Micron; HBM opex "
     "now moving hyperscaler capex guides"],
    ["HBM economics",
     "HBM3e contract prices +~20% y/y; HBM4 BOM ~2x HBM3e",
     "Bernstein: HBM prices would need to ~triple to match conventional DRAM "
     "margin per wafer — pricing is cyclical, not structural"],
    ["CoWoS packaging",
     "Doubling to 120–130k wpm by Q4-26, yet fully allocated through "
     "mid-2027; lead times 52–78 wks",
     "NVIDIA holds ~60% of allocation — a moat made of booking windows"],
    ["GPU pricing",
     "H100/H200 contract pricing +~40% Oct-25→Mar-26; reserved lead times "
     "36–52 wks",
     "Driven by memory pass-through, not die scarcity; spot rentals falling "
     "($8→$2/hr H100) while reserved rises"],
    ["Spec responses",
     "NVIDIA evaluating Rubin Ultra HBM downgrades (12-hi→8-hi); 12→8 stacks "
     "= ~50% more accelerators per DRAM input",
     "Output preserved at lower performance/card — early, quantifiable "
     "demand destruction for memory content"],
]
styled_table(slide, Inches(0.55), Inches(1.62), Inches(12.23), rows,
             [1.55, 5.34, 5.34], font_size=9.0, header_size=10,
             row_h=0.84, header_h=0.32)
add_box(slide, Inches(0.55), Inches(6.15), Inches(12.23), Inches(0.62),
        fill=NAVY, round_=True, radius=0.12)
add_text(slide, Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.45),
         [[("Analyst take:  ", True, 10.5, GOLD),
           ("the \"chip shortage\" of 2023–24 was solved; the 2026 shortage is "
            "memory and packaging. Model accelerator units off HBM stacks and "
            "CoWoS wafers, not off GPU die starts.", False, 10.5, WHITE)]],
         space_after=0)
source(slide, "TrendForce; Micron / SK Hynix / Samsung disclosures; Silicon "
       "Analysts; Stratpace Advisory; Value Add VC.", y=Inches(6.9))

# ------------------------------------------- non-consensus helper
def non_consensus_slide(num, title, sub, consensus, ourview, falsify,
                        trade=None):
    slide, page = new_slide()
    header(slide, page, f"Non-consensus view #{num}", title, sub)
    col_w = Inches(5.99)
    # consensus panel
    add_box(slide, Inches(0.55), Inches(1.62), col_w, Inches(0.42),
            fill=MID)
    add_text(slide, Inches(0.75), Inches(1.68), Inches(5.6), Inches(0.3),
             [[("WHAT CONSENSUS BELIEVES", True, 11, INK)]], space_after=0)
    add_box(slide, Inches(0.55), Inches(2.04), col_w, Inches(3.55),
            fill=LIGHT)
    bullets(slide, Inches(0.78), Inches(2.2), Inches(5.55), Inches(3.3),
            consensus, size=10, space_after=8)
    # our view panel
    add_box(slide, Inches(6.79), Inches(1.62), col_w, Inches(0.42),
            fill=NAVY)
    add_text(slide, Inches(6.99), Inches(1.68), Inches(5.6), Inches(0.3),
             [[("OUR VIEW", True, 11, WHITE),
               ("   — variant perception", False, 10, GOLD)]], space_after=0)
    add_box(slide, Inches(6.79), Inches(2.04), col_w, Inches(3.55),
            fill=NAVY2)
    paras = []
    for lead, rest in ourview:
        runs = []
        if lead:
            runs.append((f"{lead}  ", True, 10, GOLD))
        runs.append((rest, False, 10, WHITE))
        paras.append(runs)
    add_text(slide, Inches(7.02), Inches(2.2), Inches(5.55), Inches(3.3),
             paras, space_after=8, line_spacing=1.05)
    # falsification strip
    add_box(slide, Inches(0.55), Inches(5.78), Inches(12.23), Inches(0.55),
            fill=WHITE, line=RED, round_=True, radius=0.12)
    add_text(slide, Inches(0.8), Inches(5.87), Inches(11.8), Inches(0.4),
             [[("What would prove us wrong:  ", True, 9.5, RED),
               (falsify, False, 9.5, INK)]], space_after=0)
    if trade:
        add_text(slide, Inches(0.55), Inches(6.45), Inches(12.23),
                 Inches(0.45),
                 [[("How to express it:  ", True, 10, NAVY),
                   (trade, False, 10, INK)]], space_after=0)
    return slide


# --------------------------------------- slide 10: NC #1 depreciation wall
non_consensus_slide(
    1,
    "The depreciation wall is the dated, mechanical risk to 2027–28 orders",
    "Consensus debates \"bubble vs no bubble\"; we track a public-filings-"
    "visible margin event",
    [
        ("", "Five-to-six-year server schedules are justified by the \"value "
         "cascade\": GPUs move from frontier training to inference and stay "
         "productive for years."),
        ("", "Depreciation is an accounting abstraction; what matters is "
         "demand, which remains supply-constrained."),
        ("", "Useful-life extensions (Microsoft DC/office 15→25 yrs) show "
         "assets last longer, supporting current EPS."),
    ],
    [
        ("The gap is enormous and dated.", "Trailing-4Q Big-4 capex ~$434bn "
         "vs ~$149bn/yr recognized D&A. As recognized D&A compounds 30–40% "
         "annually into 2027–29, margins compress mechanically."),
        ("The cascade is being repriced in real time.", "H100 rentals fell "
         "$8→$2/GPU-hr; Rubin targets ~10x better inference token cost than "
         "Blackwell. Six-year schedules assume revenue rates the market has "
         "already broken."),
        ("Managements are pre-funding it.", "Meta $30bn, Alphabet ~$56bn, "
         "Amazon $40bn of debt plus SPVs; Amazon has already reversed a "
         "useful-life extension. CFOs facing margin questions moderate capex "
         "(Amazon 2022–23 is the template)."),
    ],
    "Take-or-pay renewals in 2026–27 re-price at or above original rates; "
    "older-generation fleets sustain high utilization at stable pricing; "
    "further credible useful-life extensions for accelerators.",
    "Treat useful-life disclosures (10-K estimate-change paragraphs, Jan–Feb) "
    "as the leading indicator for the supplier chain; fade 2028 consensus "
    "capex growth for GPU-levered names, not 2026–27.")

# --------------------------------------- slide 11: NC #2 power sets growth
non_consensus_slide(
    2,
    "Power — not silicon — sets the 2027+ growth rate",
    "Consensus still models accelerator allocations as the gate on compute "
    "revenue; the gate has moved",
    [
        ("", "GPU supply (TSMC wafers, HBM, CoWoS) determines how much AI "
         "capacity comes online each quarter."),
        ("", "Power is a cost line and an ESG topic; utilities will catch up "
         "as they always have."),
        ("", "Announced megawatts and capex guides are good proxies for "
         "future compute revenue."),
    ],
    [
        ("Chips wait on megawatts.", "Buildings finish in 2–3 yrs; firm power "
         "arrives in 5–7. Up to half of planned 2026 US capacity (~12 GW) may "
         "slip or cancel — with GPUs already allocated."),
        ("Speed-to-power is the scarce capability.", "Behind-the-meter "
         "generation (12–18 mo), flexible/curtailable interconnection, and "
         "queue position differentiate operators far more than chip access "
         "in 2027."),
        ("The KPI to model is MW energized, on dates.", "Press-release "
         "megawatts are worthless; energization schedules convert capex to "
         "revenue. This favors the equipment oligopoly, grid E&C, and "
         "power-secured operators over late-queue developers."),
    ],
    "FERC fast-tracking plus flexible-load tariffs materially compress "
    "time-to-power in 2027; turbine/transformer capacity additions land "
    "earlier than the 2027–28 guidance.",
    "Overweight grid equipment and electrical E&C (backlog visibility to "
    "2030); within operators, pay the scarcity premium for secured, dated "
    "power; underweight late-queue speculative developers.")

# --------------------------------------- slide 12: NC #3 memory peak pricing
non_consensus_slide(
    3,
    "HBM pricing power peaks before HBM capacity does",
    "Consensus extrapolates the memory supercycle through 2027; we see the "
    "top forming in the order book",
    [
        ("", "HBM is sold out through 2027; pricing power is structural. "
         "Supplier operating profits re-rate (KRW 85–89tn forecasts, +286–"
         "539% y/y)."),
        ("", "HBM4's ~2x BOM step-up entrenches the premium; AI demand is "
         "price-inelastic."),
        ("", "SK Hynix's ~60% share and qualification lead make allocation "
         "leverage durable."),
    ],
    [
        ("Spec downgrades are demand destruction.", "NVIDIA is evaluating "
         "Rubin Ultra with 8-high instead of 12-high stacks; 12→8 stacks "
         "yields ~50% more accelerators per DRAM input. Customers are "
         "engineering memory OUT at the margin."),
        ("The margin math caps the upside.", "HBM prices would need to "
         "~triple to match conventional DRAM margin per wafer capacity "
         "(Bernstein). As capacity scales, mix shifts back toward commodity "
         "DRAM economics."),
        ("Watch the 2027 crossover.", "All three vendors are HBM4-qualified; "
         "CoWoS doubles by Q4-26. When fill rates normalize above ~85%, "
         "contract-price momentum turns first — before volume does."),
    ],
    "HBM4E 16-high qualifies on schedule with strong yields and NVIDIA "
    "reverts to full-spec configurations; 2027 fill rates stay below ~70% "
    "despite capacity adds.",
    "Stay long memory through the 2026 sold-out window, but roll exposure "
    "toward suppliers with commodity-DRAM torque and packaging/equipment "
    "names as the HBM premium mean-reverts into 2027–28.")

# --------------------------------------- slide 13: NC #4 circularity
non_consensus_slide(
    4,
    "Circular financing means the cycle can crack without AI 'failing'",
    "The right lesson from 1999–2001: usage kept doubling while the balance "
    "sheets built on it collapsed",
    [
        ("", "Demand is real and contracted 18–24 months forward, so the "
         "financing structure is a footnote."),
        ("", "Hyperscalers fund capex from operating cash flow; this is "
         "nothing like dot-com vendor financing."),
        ("", "A correction requires an AI 'disappointment' — falling model "
         "quality or user churn."),
    ],
    [
        ("The loop is documented.", "Chipmakers and hyperscalers take equity "
         "in labs that spend the proceeds on their compute (BIS systemic-risk "
         "work; Bloomberg's Microsoft–OpenAI–NVIDIA deal graph). Revenue "
         "circulates inside the ecosystem before it arrives from end "
         "customers."),
        ("The marginal buyer is unprofitable.", "OpenAI reportedly burns "
         "~$60bn/yr on compute against ~$13bn revenue; neoclouds carry "
         "junk-rated debt against 2–3 yr assets. A 20–30% order cut by one "
         "large buyer cascades through chips, power and construction."),
        ("A correction needs only deceleration.", "Spending must merely stop "
         "rising for write-downs, credit tightening and capex cuts to "
         "feed each other. Telecom traffic doubled straight through 2001 — "
         "infrastructure equities fell anyway."),
    ],
    "Frontier-lab revenue compounds fast enough to close the burn gap; "
    "financing shifts from vendor equity and SPVs to end-customer cash "
    "contracts at stable prices.",
    "This is a risk-management view, not a short call on 2026: demand "
    "conservative multiples for GPU-rental and single-customer exposure; "
    "prefer diversified toll collectors with non-AI demand (grid, utilities).")

# ------------------------------------------------ slide 14: positioning
slide, page = new_slide()
header(slide, page, "Positioning", "Sector scorecard: own the toll "
       "collectors, be selective on the toll payers",
       "Framework stances — pricing power against the bottlenecks determines "
       "who keeps the economics")
rows = [
    ["Subsector", "Stance", "Why", "Key risk"],
    ["Power equipment OEMs (transformers, switchgear, turbines)",
     ("Overweight", True, GREEN),
     "Backlogs to 2030; prices +4–10%; oligopoly with labor moat",
     "Double-ordering unwinds into 2028 capacity adds"],
    ["Electrical E&C / grid contractors", ("Overweight", True, GREEN),
     "Labor-constrained trades = durable pricing power",
     "Wage inflation outpacing contract escalators"],
    ["Utilities / IPPs with DC load growth",
     ("Overweight (selective)", True, GREEN),
     "Rate-base growth; speed-to-power deals command premiums",
     "Regulatory pushback on ratepayer cost allocation"],
    ["Memory (HBM complex)", ("Overweight near term", True, GOLD),
     "2026–27 sold out; ASPs rising; fill rates 60–70%",
     "Our NC#3: spec downgrades + 2027 capacity crossover"],
    ["GPU / accelerator complex", ("Neutral-positive", True, GOLD),
     "Allocation economics intact; CoWoS moat via booking windows",
     "Memory cost pass-through; customer depreciation wall (NC#1)"],
    ["Neoclouds / GPU rental", ("Underweight", True, RED),
     "Commodity service; spot rates $8→$2/hr; debt vs 2–3 yr assets",
     "Take-or-pay renewals 2026–27 mark the model to market"],
    ["Data center REITs / colos", ("Selective", True, GOLD),
     "Power-secured portfolios earn a scarcity premium",
     "Late-queue developments strand capex in the 5–7 yr gap"],
]
styled_table(slide, Inches(0.55), Inches(1.62), Inches(12.23), rows,
             [3.28, 1.85, 3.75, 3.35], font_size=8.8, header_size=10,
             row_h=0.62, header_h=0.34)
add_text(slide, Inches(0.55), Inches(6.35), Inches(12.23), Inches(0.5),
         [[("Unifying logic:  ", True, 10, NAVY),
           ("every bottleneck on slide 6 is somebody's income statement. Own "
            "the constraint; rent the beneficiaries of its relief; avoid "
            "balance sheets that assume the constraint never binds.",
            False, 10, INK)]])
source(slide, "Team framework. Stances are illustrative, not "
       "recommendations.", y=Inches(6.9))

# ------------------------------------------------ slide 15: monitoring
slide, page = new_slide()
header(slide, page, "Monitoring", "The dashboard: ten series that will "
       "flag it before the tape does",
       "Each indicator maps to a specific call; cadence and source noted for "
       "reproducibility")
rows = [
    ["Indicator", "Cadence", "Maps to", "What it tells us"],
    ["Useful-life / estimate-change paragraphs in hyperscaler 10-Ks",
     "Annual (Jan–Feb)", "NC#1",
     "First public sign of the depreciation wall; Amazon already reversed "
     "one extension"],
    ["Hyperscaler capex guides + finance-lease / SPV mix", "Quarterly",
     "NC#1, #4", "Commitment vs accounting optics; pre-funding behavior"],
    ["HBM contract negotiations, customer fill rates", "Quarterly", "NC#3",
     "Fill rates normalizing above ~85% = pricing power rolling over"],
    ["CoWoS capacity and allocation (target 120–130k wpm Q4-26)",
     "Monthly (trade press)", "NC#3", "Packaging constraint easing; NVIDIA "
     "share of allocation"],
    ["GEV / Siemens Energy / Hitachi orders, book-to-bill, slot pricing",
     "Quarterly", "Positioning", "Duration of the equipment cycle; backlog "
     "quality vs double-ordering"],
    ["LBNL Queued Up; FERC interconnection reform dockets", "Annual+",
     "NC#2", "Structural relief in time-to-power (median 61 months today)"],
    ["MW energized vs MW announced, by operator", "Ongoing", "NC#2",
     "The real revenue gate; press releases are not power"],
    ["H100/H200 spot rental rates; take-or-pay renewal pricing", "Ongoing",
     "NC#1", "Health of the 'value cascade' that underpins 5–6 yr schedules"],
    ["Transformer / switchgear lead-time surveys (Wood Mackenzie, POWER)",
     "Quarterly", "Lead-time board", "Directional turn = first sign of the "
     "2027–28 capacity landing"],
    ["NERC reliability assessments and alerts", "Semiannual", "NC#2",
     "Siting constraints; grid stress from DC load oscillations"],
]
styled_table(slide, Inches(0.55), Inches(1.62), Inches(12.23), rows,
             [4.30, 1.42, 1.13, 5.38], font_size=8.4, header_size=9.5,
             row_h=0.475, header_h=0.32)
source(slide, "Team compilation. NC# = non-consensus view number (slides "
       "10–13).", y=Inches(6.9))

# ------------------------------------------------ slide 16: risks
slide, page = new_slide()
header(slide, page, "Risks & disclosures", "What breaks the thesis — in "
       "either direction",
       "We hold the constraint views with conviction and the cycle views "
       "with humility")
add_box(slide, Inches(0.55), Inches(1.62), Inches(5.99), Inches(0.42),
        fill=GREEN)
add_text(slide, Inches(0.75), Inches(1.68), Inches(5.6), Inches(0.3),
         [[("UPSIDE RISKS (to our cautious calls)", True, 11, WHITE)]],
         space_after=0)
add_box(slide, Inches(0.55), Inches(2.04), Inches(5.99), Inches(3.7),
        fill=LIGHT)
bullets(slide, Inches(0.78), Inches(2.2), Inches(5.55), Inches(3.45), [
    ("Inference demand steepens.", "Token demand outgrows efficiency gains, "
     "sustaining the value cascade and validating long useful lives."),
    ("Power reform lands early.", "FERC fast-tracking plus flexible-load "
     "tariffs compress time-to-power; the 2029 relief arrives in 2027."),
    ("Supply-side positive surprise.", "Samsung HBM4 yield ramp and CoWoS "
     "overshoot ease silicon constraints; equipment capacity beats "
     "schedules."),
    ("Monetization broadens.", "Enterprise AI revenue diversifies the buyer "
     "base beyond frontier labs, defusing circularity risk."),
], size=10, space_after=9)
add_box(slide, Inches(6.79), Inches(1.62), Inches(5.99), Inches(0.42),
        fill=RED)
add_text(slide, Inches(6.99), Inches(1.68), Inches(5.6), Inches(0.3),
         [[("DOWNSIDE RISKS (to our constructive calls)", True, 11, WHITE)]],
         space_after=0)
add_box(slide, Inches(6.79), Inches(2.04), Inches(5.99), Inches(3.7),
        fill=LIGHT)
bullets(slide, Inches(7.02), Inches(2.2), Inches(5.55), Inches(3.45), [
    ("One buyer blinks.", "A 20–30% order cut by a single hyperscaler "
     "cascades through equipment backlogs — how much is double-ordered?"),
    ("Memory inflation bites compute demand.", "HBM/DRAM cost pass-through "
     "compresses GPU orders and hyperscaler margins simultaneously."),
    ("Politics of power.", "Ratepayer backlash, local moratoria and "
     "grid-reliability alerts slow siting faster than reform speeds it."),
    ("Backlog quality.", "Equipment 'orders' include reservations and slots; "
     "cancellation terms are untested in a downturn."),
], size=10, space_after=9)
add_box(slide, Inches(0.55), Inches(5.95), Inches(12.23), Inches(1.0),
        fill=NAVY, round_=True, radius=0.06)
add_text(slide, Inches(0.8), Inches(6.06), Inches(11.75), Inches(0.85),
         [[("Disclosures: ", True, 8.5, GOLD),
           ("This document is an illustrative research template prepared for "
            "demonstration purposes. It is not investment research, an offer, "
            "or investment advice. Figures are compiled from public sources "
            "as of August 2026, may be approximate or estimated, and should "
            "be independently verified. Sources include company filings and "
            "calls, LBNL, NERC, Wood Mackenzie, BloombergNEF, TrendForce, "
            "and trade press. Past performance is not indicative of future "
            "results.", False, 8.5, WHITE)]], line_spacing=1.05)

# -------------------------------------------------------------------- save --
OUT = "Datacenter_Sellside_Deck_Aug2026.pptx"
prs.save(OUT)
print(f"Wrote {OUT} with {_page[0]} slides")
